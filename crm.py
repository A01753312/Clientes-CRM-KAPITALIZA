# App para gestión de clientes / CRM
# - Sidebar con filtros/acciones
# - Pestañas ordenadas: Dashboard | Clientes | Documentos | Importar | Historial
# - Auto-refresh tras subir documentos o importar Excel
# - Se ELIMINA la importación de documentos desde ZIP

# === CONFIGURACIÓN PROFESIONAL DEL CRM ===
import streamlit as st
import base64
from datetime import date, datetime
from pathlib import Path

# Configuración de página (DEBE SER LO PRIMERO)
st.set_page_config(
    page_title="Kapitaliza | CRM",
    page_icon="💼",
    layout="wide",
    initial_sidebar_state="expanded",
    menu_items={
        'Get Help': 'https://kapitaliza.com/soporte',
        'Report a bug': 'https://kapitaliza.com/contacto',
        'About': '**Kapitaliza CRM** - Plataforma integral de gestión de clientes y relaciones comerciales'
    }
)

# === AUTENTICACIÓN CON GOOGLE DRIVE ===
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload
import json

if "drive_creds" not in st.session_state:
    st.session_state.drive_creds = None

CLIENT_ID = st.secrets["GOOGLE_CLIENT_ID"]
CLIENT_SECRET = st.secrets["GOOGLE_CLIENT_SECRET"]
REDIRECT_URI = st.secrets["REDIRECT_URI"]

# Scopes actualizados según la documentación oficial de Google
SCOPES = [
    "https://www.googleapis.com/auth/drive.file",
    "https://www.googleapis.com/auth/drive.metadata.readonly"
]

# Mostrar botón de conexión si aún no se ha autenticado
if not st.session_state.drive_creds:
    # Construir URL de autorización manualmente con scopes correctos
    scope_string = "%20".join([scope.replace("https://www.googleapis.com/auth/", "") for scope in SCOPES])
    auth_url = (
        "https://accounts.google.com/o/oauth2/v2/auth"
        f"?response_type=code&client_id={CLIENT_ID}"
        f"&redirect_uri={REDIRECT_URI}"
        f"&scope=https://www.googleapis.com/auth/drive.file%20https://www.googleapis.com/auth/drive.metadata.readonly"
        f"&access_type=offline&prompt=consent"
    )
    st.sidebar.markdown("---")
    st.sidebar.markdown("### 📂 Conexión a Google Drive")
    st.sidebar.markdown(f"[🔐 Conectar con Google Drive]({auth_url})")

else:
    st.sidebar.markdown("---")
    st.sidebar.markdown("### 📂 Google Drive")
    st.sidebar.success("✅ Conectado a Google Drive")
    
    # Botón para desconectar Google Drive
    if st.sidebar.button("🔌 Desconectar Drive", help="Cerrar sesión de Google Drive"):
        st.session_state.drive_creds = None
        if "processed_auth_code" in st.session_state:
            del st.session_state.processed_auth_code
        st.query_params.clear()
        st.sidebar.success("Google Drive desconectado")
        st.rerun()

# Procesar el parámetro de autorización devuelto por Google
query_params = st.query_params
if "code" in query_params and not st.session_state.drive_creds:
    code = query_params["code"]
    
    # Solo procesar si no hemos procesado este código antes
    if "processed_auth_code" not in st.session_state or st.session_state.processed_auth_code != code:
        try:
            # Configuración del cliente OAuth2 actualizada
            client_config = {
                "web": {
                    "client_id": CLIENT_ID,
                    "client_secret": CLIENT_SECRET,
                    "redirect_uris": [REDIRECT_URI],
                    "auth_uri": "https://accounts.google.com/o/oauth2/auth",
                    "token_uri": "https://oauth2.googleapis.com/token",
                    "auth_provider_x509_cert_url": "https://www.googleapis.com/oauth2/v1/certs"
                }
            }
            
            flow = Flow.from_client_config(
                client_config,
                scopes=SCOPES
            )
            flow.redirect_uri = REDIRECT_URI
            flow.fetch_token(code=code)
            
            # Guardar credenciales y marcar el código como procesado
            st.session_state.drive_creds = flow.credentials
            st.session_state.processed_auth_code = code
            
            # Limpiar el código de la URL para evitar reprocessing
            st.query_params.clear()
            # Usar función simple sin dependencias
            if f"shown_drive_auth" not in st.session_state:
                st.success("✅ Autenticación exitosa con Google Drive")
                st.session_state[f"shown_drive_auth"] = True
            st.rerun()
            
        except Exception as e:
            # Quitar mensajes molestos antes del login
            # st.error(f"❌ Error en la autenticación: {str(e)}")
            # Limpiar el código problemático
            st.query_params.clear()
            st.session_state.processed_auth_code = None
            # st.sidebar.error("Error de autenticación. Intenta conectar nuevamente.")

# Si hay un error en la URL, limpiar silenciosamente
if "error" in query_params:
    error = query_params["error"]
    # Quitar mensaje molesto: st.sidebar.error(f"❌ Error de autorización: {error}")
    st.query_params.clear()

# CSS personalizado para look profesional con tema claro Kapitaliza
st.markdown("""
<style>
    /* Ocultar el menú de Streamlit y el footer */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    
    /* Tema claro corporativo Kapitaliza */
    .stApp {
        background-color: #FAFBFC;
        color: #2D3748;
        font-family: 'Inter', 'Roboto', sans-serif;
    }
    
    /* Estilo del header */
    .stApp header {
        background-color: #FFFFFF;
        border-bottom: 2px solid #FFD41D;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
    }
    
    /* Mejorar tabs con tema claro */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
        background-color: #FFFFFF;
        padding: 12px;
        border-radius: 8px;
        border: 1px solid #E2E8F0;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
    }
    
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        padding: 0 24px;
        background-color: #F7FAFC;
        border-radius: 8px;
        border: 2px solid #E2E8F0;
        font-weight: 600;
        color: #4A5568;
        transition: all 0.3s ease;
    }
    
    .stTabs [aria-selected="true"] {
        background-color: #FFD41D;
        color: #2D3748;
        border-color: #FFD41D;
        box-shadow: 0 2px 8px rgba(255, 212, 29, 0.3);
        font-weight: 700;
    }
    
    .stTabs [data-baseweb="tab"]:hover {
        background-color: #EDF2F7;
        color: #2D3748;
        border-color: #FFD41D;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
    }
    
    /* Mejorar métricas con tema claro */
    [data-testid="stMetricValue"] {
        font-size: 32px;
        font-weight: 700;
        color: #2D3748;
        font-family: 'Poppins', sans-serif;
    }
    
    [data-testid="stMetricLabel"] {
        color: #4A5568;
        font-weight: 600;
        font-size: 14px;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }
    
    [data-testid="stMetricDelta"] {
        color: #718096;
        font-weight: 500;
    }
    
    /* Cards para métricas corporativas con tema claro */
    div[data-testid="metric-container"] {
        background-color: #FFFFFF;
        padding: 20px;
        border-radius: 12px;
        border: 2px solid #E2E8F0;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.08);
        transition: transform 0.3s ease;
    }
    
    div[data-testid="metric-container"]:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 25px rgba(0, 0, 0, 0.12);
        border-color: #FFD41D;
    }
    
    /* Dataframes corporativos con tema claro */
    .dataframe {
        border: 2px solid #E2E8F0;
        border-radius: 10px;
        background-color: #FFFFFF;
        color: #2D3748;
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.08);
    }
    
    .dataframe th {
        background-color: #F7FAFC;
        color: #2D3748;
        font-weight: 600;
        border-bottom: 2px solid #E2E8F0;
    }
    
    .dataframe td {
        background-color: #FFFFFF;
        color: #4A5568;
        border-bottom: 1px solid #F1F5F9;
    }
    
    .dataframe tr:hover {
        background-color: #F7FAFC;
    }
    
    /* Botones corporativos Kapitaliza con tema claro */
    .stButton > button {
        border-radius: 8px;
        font-weight: 600;
        font-size: 14px;
        transition: all 0.3s ease;
        background-color: #FFD41D;
        color: #2D3748;
        border: none;
        padding: 12px 24px;
        text-transform: uppercase;
        letter-spacing: 0.5px;
        box-shadow: 0 2px 8px rgba(255, 212, 29, 0.3);
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(255, 212, 29, 0.5);
        background-color: #E6B800;
        color: #1A202C;
    }
    
    /* Sidebar corporativo con tema claro */
    [data-testid="stSidebar"] {
        background-color: #FFFFFF;
        border-right: 2px solid #FFD41D;
        box-shadow: 2px 0 4px rgba(0,0,0,0.05);
    }
    
    [data-testid="stSidebar"] .css-1d391kg {
        color: #2D3748;
    }
    
    /* Expanders corporativos */
    .streamlit-expanderHeader {
        background-color: #F7FAFC;
        border-radius: 8px;
        font-weight: 600;
        color: #2D3748;
        border: 1px solid #E2E8F0;
    }
    
    .streamlit-expanderContent {
        background-color: #FFFFFF;
        border: 1px solid #E2E8F0;
        border-radius: 0 0 8px 8px;
    }
    
    /* Inputs corporativos con tema claro */
    .stTextInput > div > div > input,
    .stSelectbox > div > div > select,
    .stMultiSelect > div > div {
        border-radius: 8px;
        border: 2px solid #E2E8F0;
        background-color: #FFFFFF;
        color: #2D3748;
        font-weight: 500;
        transition: border-color 0.3s ease;
    }
    
    .stTextInput > div > div > input:focus,
    .stSelectbox > div > div > select:focus {
        border-color: #FFD41D;
        box-shadow: 0 0 0 2px rgba(255, 212, 29, 0.2);
        background-color: #FFFBF0;
    }
    
    /* Spinner corporativo */
    .stSpinner > div {
        border-top-color: #FFD41D !important;
    }
    
    /* Títulos y encabezados con tema claro */
    h1, h2, h3 {
        color: #2D3748;
        font-family: 'Poppins', sans-serif;
        font-weight: 700;
    }
    
    h1 {
        font-size: 2.5rem;
        color: #1A202C;
    }
    
    h2 {
        font-size: 2rem;
        color: #2D3748;
    }
    
    h3 {
        font-size: 1.5rem;
        color: #4A5568;
    }
    
    /* Texto secundario con tema claro */
    .stMarkdown p, .stText {
        color: #4A5568;
        line-height: 1.6;
    }
    
    /* Info boxes con tema claro mejorado */
    .stAlert[data-baseweb="notification"] {
        background-color: #FFFFFF;
        border: 2px solid #E2E8F0;
        border-radius: 8px;
        color: #2D3748;
        box-shadow: 0 2px 8px rgba(0,0,0,0.08);
    }
    
    /* Success (Verde suave) */
    .stSuccess {
        background-color: #F0FFF4;
        border: 2px solid #4CAF50;
        color: #22543D;
        border-radius: 8px;
        box-shadow: 0 2px 8px rgba(76, 175, 80, 0.1);
    }
    
    /* Error (Rojo profesional) */
    .stError {
        background-color: #FFF5F5;
        border: 2px solid #D9534F;
        color: #742A2A;
        border-radius: 8px;
        box-shadow: 0 2px 8px rgba(217, 83, 79, 0.1);
    }
    
    /* Warning */
    .stWarning {
        background-color: #FFFBF0;
        border: 2px solid #FFD41D;
        color: #744210;
        border-radius: 8px;
        box-shadow: 0 2px 8px rgba(255, 212, 29, 0.1);
    }
    
    /* Info - Estilo especial para header */
    .stInfo {
        background-color: #F7FAFC;
        border: 2px solid #FFD41D;
        color: #2D3748;
        border-radius: 8px;
        box-shadow: 0 2px 8px rgba(255, 212, 29, 0.1);
        text-align: center;
        font-weight: 500;
    }
    
    /* Selectbox dropdown corporativo */
    .stSelectbox > div > div > div {
        background-color: #FFFFFF;
        color: #2D3748;
        border: 1px solid #E2E8F0;
        border-radius: 8px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.08);
    }
    
    /* Multiselect corporativo */
    .stMultiSelect > div > div > div {
        background-color: #FFFFFF;
        color: #2D3748;
        border: 1px solid #E2E8F0;
        border-radius: 8px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.08);
    }
    
    /* Radio buttons con tema claro */
    .stRadio > div > label {
        color: #2D3748;
        font-weight: 500;
    }
    
    /* Checkbox con tema claro */
    .stCheckbox > label {
        color: #2D3748;
        font-weight: 500;
    }
    
    /* Toast messages */
    .stToast {
        background-color: #FFFFFF;
        color: #2D3748;
        border: 1px solid #FFD41D;
        border-radius: 8px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.1);
    }
    
    /* Progress bar */
    .stProgress > div > div {
        background-color: #FFD41D;
    }
    
    /* Caption text con tema claro */
    .css-1v0mbdj {
        color: #718096;
        font-size: 12px;
    }
    
    /* Mejorar texto general para tema claro */
    .stMarkdown, .stText, p, span {
        color: #4A5568 !important;
    }
    
    /* Títulos de secciones */
    .css-10trblm {
        color: #2D3748;
    }
    
    /* Separadores más suaves */
    hr {
        border-color: #E2E8F0;
        opacity: 0.6;
    }
    
    /* Mejoras adicionales para tema claro */
    .stContainer {
        background-color: #FAFBFC;
    }
    
    /* Mejora de contraste en elementos seleccionados */
    .stSelectbox > div > div > div:hover {
        background-color: #F7FAFC;
    }
    
    /* Estilo para números y datos importantes */
    .highlight-number {
        color: #E6B800;
        font-weight: 700;
    }
    
    /* Botones pequeños y compactos */
    .small-button button, .small-refresh-button button {
        height: 2rem !important;
        padding: 0.25rem 0.75rem !important;
        font-size: 0.8rem !important;
        font-weight: 500 !important;
        min-height: auto !important;
        border-radius: 6px !important;
    }
    
    /* Contenedores de botones pequeños */
    .small-button, .small-refresh-button {
        display: flex;
        justify-content: center;
        align-items: center;
    }
</style>
""", unsafe_allow_html=True)

import io
import re
import json
import zipfile  # (aún se usa para exportar ZIPs, no para importar)  # CHANGED: seguimos exportando
import os
import hashlib
import secrets
import difflib
import unicodedata
import time

import pandas as pd
import gspread
from gspread_dataframe import get_as_dataframe, set_with_dataframe
from google.oauth2.service_account import Credentials
import shutil
import altair as alt
from google.auth.transport.requests import Request

# Debug info removed by user request (sidebar debug block intentionally deleted)

# === FUNCIONES PROFESIONALES DEL CRM ===

# Función para mostrar loading personalizado
def show_loading(message="Procesando..."):
    """Muestra un loading spinner más profesional"""
    return st.spinner(f"⚙️ {message}")

# Función para notificaciones más profesionales
def show_notification(message: str, type: str = "success"):
    """
    Muestra notificaciones profesionales
    type: 'success', 'info', 'warning', 'error'
    """
    icons = {
        "success": "✅",
        "info": "ℹ️",
        "warning": "⚠️",
        "error": "❌"
    }
    
    icon = icons.get(type, "ℹ️")
    
    # Usar toast para notificaciones no intrusivas
    try:
        st.toast(f"{icon} {message}", icon=icon)
    except:
        # Fallback si toast no está disponible
        if type == "success":
            st.success(f"{icon} {message}")
        elif type == "error":
            st.error(f"{icon} {message}")
        elif type == "warning":
            st.warning(f"{icon} {message}")
        else:
            st.info(f"{icon} {message}")

# Funciones de feedback mejorado
def show_success(message, duration=3):
    """Muestra mensaje de éxito usando componentes nativos de Streamlit"""
    st.toast(f"✅ {message}", icon="✅")
    st.success(f"✅ {message}")

def show_error(message):
    """Muestra mensaje de error usando componentes nativos de Streamlit"""
    st.toast(f"❌ {message}", icon="❌")
    st.error(f"❌ {message}")

def show_warning(message):
    """Muestra advertencia usando componentes nativos de Streamlit"""
    st.toast(f"⚠️ {message}", icon="⚠️")
    st.warning(f"⚠️ {message}")

def show_info(message):
    """Muestra información usando componentes nativos de Streamlit"""
    st.toast(f"ℹ️ {message}", icon="ℹ️")
    st.info(f"ℹ️ {message}")

# Función para mostrar mensajes solo una vez por sesión
def show_once(message_key: str, message_func, *args, **kwargs):
    """
    Muestra un mensaje solo una vez por sesión.
    
    Args:
        message_key: Clave única para identificar el mensaje
        message_func: Función de Streamlit para mostrar el mensaje (st.success, st.info, etc.)
        *args, **kwargs: Argumentos para la función de mensaje
    """
    if f"shown_{message_key}" not in st.session_state:
        message_func(*args, **kwargs)
        st.session_state[f"shown_{message_key}"] = True

def show_once_success(message_key: str, message: str):
    """Muestra mensaje de éxito solo una vez por sesión"""
    show_once(message_key, st.success, f"✅ {message}")

def show_once_info(message_key: str, message: str):
    """Muestra mensaje de información solo una vez por sesión"""
    show_once(message_key, st.info, f"ℹ️ {message}")

def show_once_warning(message_key: str, message: str):
    """Muestra mensaje de advertencia solo una vez por sesión"""
    show_once(message_key, st.warning, f"⚠️ {message}")

# Cards profesionales para KPIs
def render_kpi_card(label: str, value: any, delta: str = None, icon: str = "📊", color: str = "#0066cc"):
    """Renderiza una tarjeta KPI profesional"""
    delta_html = ""
    if delta:
        delta_color = "#28a745" if "+" in str(delta) or "↑" in str(delta) else "#dc3545"
        delta_html = f'<p style="color: {delta_color}; margin: 5px 0 0 0; font-size: 14px; font-weight: 500;">{delta}</p>'
    
    card_html = f"""
    <div style="
        background: white;
        padding: 20px;
        border-radius: 12px;
        border-left: 4px solid {color};
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        margin-bottom: 10px;
        height: 140px;
        display: flex;
        flex-direction: column;
        justify-content: space-between;
    ">
        <div style="display: flex; justify-content: space-between; align-items: flex-start;">
            <div style="flex: 1;">
                <p style="color: #6c757d; margin: 0; font-size: 14px; font-weight: 500; text-transform: uppercase; letter-spacing: 0.5px;">
                    {label}
                </p>
                <h2 style="color: #212529; margin: 10px 0 0 0; font-size: 32px; font-weight: 700; line-height: 1;">
                    {value}
                </h2>
            </div>
            <div style="font-size: 32px; opacity: 0.4; margin-left: 15px;">
                {icon}
            </div>
        </div>
        {delta_html}
    </div>
    """
    st.markdown(card_html, unsafe_allow_html=True)

def calcular_analisis_financiero(df: pd.DataFrame) -> dict:
    """Calcula métricas financieras del portfolio de clientes"""
    import re
    
    def limpiar_monto(monto_str):
        """Convierte string de monto a float, manejando diferentes formatos"""
        if pd.isna(monto_str) or str(monto_str).strip() == "":
            return 0.0
        
        # Convertir a string y limpiar
        monto_clean = str(monto_str).strip()
        
        # Remover símbolos de moneda y espacios
        monto_clean = re.sub(r'[,$\s]', '', monto_clean)
        
        # Intentar convertir a float
        try:
            return float(monto_clean)
        except (ValueError, TypeError):
            return 0.0
    
    # Limpiar y convertir montos
    df_temp = df.copy()
    df_temp['monto_propuesta_num'] = df_temp['monto_propuesta'].apply(limpiar_monto)
    df_temp['monto_final_num'] = df_temp['monto_final'].apply(limpiar_monto)
    
    # Calcular métricas
    total_propuesto = df_temp['monto_propuesta_num'].sum()
    total_dispersado = df_temp[df_temp['estatus'] == 'DISPERSADO']['monto_final_num'].sum()
    
    # Promedios
    promedio_propuesto = df_temp['monto_propuesta_num'].mean() if len(df_temp) > 0 else 0
    dispersados_df = df_temp[df_temp['estatus'] == 'DISPERSADO']
    promedio_dispersado = dispersados_df['monto_final_num'].mean() if len(dispersados_df) > 0 else 0
    
    # Efectividad de conversión
    tasa_conversion_financiera = (total_dispersado / total_propuesto * 100) if total_propuesto > 0 else 0
    
    # Análisis por estatus con montos
    montos_por_estatus = df_temp.groupby('estatus').agg({
        'monto_propuesta_num': ['sum', 'mean', 'count'],
        'monto_final_num': ['sum', 'mean']
    }).round(2)
    
    return {
        'total_propuesto': total_propuesto,
        'total_dispersado': total_dispersado,
        'promedio_propuesto': promedio_propuesto,
        'promedio_dispersado': promedio_dispersado,
        'tasa_conversion_financiera': tasa_conversion_financiera,
        'montos_por_estatus': montos_por_estatus,
        'clientes_con_monto': len(df_temp[df_temp['monto_propuesta_num'] > 0]),
        'dispersados_con_monto': len(df_temp[(df_temp['estatus'] == 'DISPERSADO') & (df_temp['monto_final_num'] > 0)])
    }

def formatear_monto(monto: float) -> str:
    """Formatea un monto para mostrar en pesos mexicanos"""
    if monto == 0:
        return "$0"
    elif monto >= 1_000_000:
        return f"${monto/1_000_000:.1f}M"
    elif monto >= 1_000:
        return f"${monto/1_000:.0f}K"
    else:
        return f"${monto:,.0f}"


def generar_presentacion_dashboard(df_cli: pd.DataFrame) -> bytes:
    """Genera una presentación PowerPoint completa del dashboard con gráficas"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from io import BytesIO
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg')  # Backend sin GUI
    
    # Crear presentación — intentar cargar plantilla .pptx si existe
    from pathlib import Path as _Path

    template_path = None
    # Rutas recomendadas donde podrías haber subido la plantilla
    candidates = [
        Path("assets/presentation_template.pptx"),
        Path("data/presentation_template.pptx")
    ]
    for c in candidates:
        if c.exists():
            template_path = c
            break

    # Si no hay plantilla en rutas conocidas, buscar el primer .pptx en el repo
    if template_path is None:
        repo_root = Path(__file__).parent
        for p in repo_root.rglob("*.pptx"):
            # evitar archivos temporales o la propia salida si existiera
            if ".git" in str(p) or "site-packages" in str(p):
                continue
            template_path = p
            break

    try:
        if template_path is not None:
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
    except Exception:
        # Fallback: crear presentación vacía
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # Calcular todas las métricas necesarias
    total_clientes = len(df_cli)
    estatus_counts = df_cli["estatus"].fillna("").value_counts()
    
    dispersados = estatus_counts.get("DISPERSADO", 0)
    rechazados = sum([
        count for estatus, count in estatus_counts.items() 
        if estatus and (estatus.startswith("RECH") or estatus.startswith("REC"))
    ])
    en_proceso = total_clientes - dispersados - rechazados
    
    tasa_exito = (dispersados / total_clientes * 100) if total_clientes > 0 else 0
    tasa_proceso = (en_proceso / total_clientes * 100) if total_clientes > 0 else 0
    tasa_rechazo = (rechazados / total_clientes * 100) if total_clientes > 0 else 0
    
    # Análisis financiero
    analisis_financiero = calcular_analisis_financiero(df_cli)
    total_presupuesto = analisis_financiero['total_propuesto']
    
    # === SLIDE 1: PORTADA ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco
    
    # Fondo de color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 0)  # Amarillo
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Dashboard CRM Kapitaliza"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 0, 0)  # Negro
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo con fecha
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    from datetime import datetime
    subtitle_frame.text = f"Reporte Ejecutivo - {datetime.now().strftime('%d/%m/%Y')}"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = RGBColor(0, 0, 0)  # Negro
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # === SLIDE 2: KPIs PRINCIPALES ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "📊 KPIs Principales"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(33, 37, 41)
    
    # KPIs en cuadros
    kpis = [
        ("Total de Clientes", total_clientes, "👥", ""),
        ("Dispersados (Éxito)", dispersados, "✅", f"{tasa_exito:.1f}%"),
        ("En Proceso", en_proceso, "⏳", f"{tasa_proceso:.1f}%"),
        ("Rechazados", rechazados, "❌", f"{tasa_rechazo:.1f}%")
    ]
    
    x_start = 0.5
    y_pos = 1.5
    width = 2.2
    height = 1.8
    gap = 0.15
    
    for i, (label, value, icon, delta) in enumerate(kpis):
        x_pos = x_start + i * (width + gap)
        
        # Caja con borde
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(x_pos), Inches(y_pos), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
        shape.line.color.rgb = RGBColor(225, 229, 233)
        
        # Etiqueta
        label_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.1), Inches(y_pos + 0.2), Inches(width - 0.2), Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.text = f"{icon} {label}"
        label_p = label_frame.paragraphs[0]
        label_p.font.size = Pt(11)
        label_p.font.color.rgb = RGBColor(108, 117, 125)
        label_p.alignment = PP_ALIGN.CENTER
        
        # Valor
        value_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.1), Inches(y_pos + 0.7), Inches(width - 0.2), Inches(0.6)
        )
        value_frame = value_box.text_frame
        value_frame.text = str(value)
        value_p = value_frame.paragraphs[0]
        value_p.font.size = Pt(36)
        value_p.font.bold = True
        value_p.font.color.rgb = RGBColor(33, 37, 41)
        value_p.alignment = PP_ALIGN.CENTER
        
        # Delta
        if delta:
            delta_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(y_pos + 1.4), Inches(width - 0.2), Inches(0.3)
            )
            delta_frame = delta_box.text_frame
            delta_frame.text = delta
            delta_p = delta_frame.paragraphs[0]
            delta_p.font.size = Pt(10)
            delta_p.font.color.rgb = RGBColor(108, 117, 125)
            delta_p.alignment = PP_ALIGN.CENTER
    
    # Gráfica de distribución de estatus (pie chart)
    fig, ax = plt.subplots(figsize=(4, 3))
    sizes = [dispersados, en_proceso, rechazados]
    labels = ['Dispersados', 'En Proceso', 'Rechazados']
    colors = ['#28a745', '#ffc107', '#dc3545']
    
    if sum(sizes) > 0:
        ax.pie(sizes, labels=labels, autopct='%1.1f%%', colors=colors, startangle=90)
        ax.axis('equal')
    
    # Guardar gráfica en BytesIO
    img_stream = BytesIO()
    plt.tight_layout()
    plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    plt.close()
    img_stream.seek(0)
    
    # Agregar gráfica al slide
    slide.shapes.add_picture(img_stream, Inches(3), Inches(3.8), width=Inches(4))
    
    # === SLIDE 3: TOP ESTATUS POR MONTO ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "💹 Top Estatus por Monto"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(33, 37, 41)
    
    # Total presupuesto
    presupuesto_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.4))
    presupuesto_frame = presupuesto_box.text_frame
    presupuesto_frame.text = f"Total Presupuesto General: {formatear_monto(total_presupuesto)}"
    presupuesto_p = presupuesto_frame.paragraphs[0]
    presupuesto_p.font.size = Pt(18)
    presupuesto_p.font.bold = True
    presupuesto_p.font.color.rgb = RGBColor(33, 37, 41)
    
    # Top estatus
    if not analisis_financiero['montos_por_estatus'].empty:
        estatus_con_monto = analisis_financiero['montos_por_estatus'][
            analisis_financiero['montos_por_estatus'][('monto_propuesta_num', 'sum')] > 0
        ]
        top_estatus = estatus_con_monto.sort_values(
            ('monto_propuesta_num', 'sum'), ascending=False
        ).head(5)
        
        # Crear gráfica de barras
        fig, ax = plt.subplots(figsize=(8, 4))
        estatus_nombres = [str(e)[:30] for e in top_estatus.index]  # Limitar longitud
        montos = [top_estatus.loc[e, ('monto_propuesta_num', 'sum')] for e in top_estatus.index]
        
        bars = ax.barh(estatus_nombres, montos, color='#28a745')
        ax.set_xlabel('Monto ($)', fontsize=11)
        ax.set_title('Top 5 Estatus por Monto', fontsize=13, fontweight='bold')
        
        # Agregar valores en las barras
        for i, (bar, monto) in enumerate(zip(bars, montos)):
            ax.text(bar.get_width(), bar.get_y() + bar.get_height()/2, 
                   f' {formatear_monto(monto)}', 
                   va='center', fontsize=10, fontweight='bold')
        
        plt.tight_layout()
        
        # Guardar gráfica
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        img_stream.seek(0)
        
        # Agregar al slide
        slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.8), width=Inches(8.4))
    
    # === SLIDE 4: ANÁLISIS FINANCIERO ===
    df_temp = df_cli.copy()
    
    def limpiar_monto_simple(monto_str):
        if pd.isna(monto_str) or str(monto_str).strip() == "":
            return 0.0
        try:
            import re
            clean = re.sub(r'[,$\s]', '', str(monto_str))
            return float(clean)
        except:
            return 0.0
    
    df_temp['monto_analisis'] = df_temp.apply(
        lambda row: limpiar_monto_simple(row['monto_final']) if row['estatus'] == 'DISPERSADO' 
        else limpiar_monto_simple(row['monto_propuesta']), axis=1
    )
    df_analisis = df_temp[df_temp['monto_analisis'] > 0].copy()
    
    if not df_analisis.empty:
        # Modelo financiero
        prob_conversion = {
            "DISPERSADO": 1.00, "APROB. CON PROPUESTA": 0.75, "PROPUESTA": 0.75,
            "PEND. ACEPT. CLIENTE": 0.65, "PENDIENTE CLIENTE": 0.65,
            "PEND. DOC. PARA EVALUACION": 0.45, "PENDIENTE DOC": 0.45,
            "EN ONBOARDING": 0.55, "RECH. CLIENTE CANCELA": 0.10,
            "RECH. SOBREENDEUDAMIENTO": 0.05,
        }
        factor_retorno = {
            "DISPERSADO": 1.00, "APROB. CON PROPUESTA": 0.85, "PROPUESTA": 0.85,
            "PEND. ACEPT. CLIENTE": 0.80, "PENDIENTE CLIENTE": 0.80,
            "PEND. DOC. PARA EVALUACION": 0.70, "PENDIENTE DOC": 0.70,
            "EN ONBOARDING": 0.75, "RECH. CLIENTE CANCELA": 0.00,
            "RECH. SOBREENDEUDAMIENTO": 0.00,
        }
        riesgo_pct = {
            "DISPERSADO": 5, "APROB. CON PROPUESTA": 20, "PROPUESTA": 20,
            "PEND. ACEPT. CLIENTE": 30, "PENDIENTE CLIENTE": 30,
            "PEND. DOC. PARA EVALUACION": 45, "PENDIENTE DOC": 45,
            "EN ONBOARDING": 40, "RECH. CLIENTE CANCELA": 90,
            "RECH. SOBREENDEUDAMIENTO": 95,
        }
        
        df_analisis["Probabilidad de Conversión"] = df_analisis["estatus"].map(prob_conversion).fillna(0.5)
        df_analisis["Factor Retorno"] = df_analisis["estatus"].map(factor_retorno).fillna(0.5)
        df_analisis["Riesgo (%)"] = df_analisis["estatus"].map(riesgo_pct).fillna(50)
        
        df_analisis["Monto Esperado"] = df_analisis["monto_analisis"] * df_analisis["Probabilidad de Conversión"]
        df_analisis["Retorno Esperado"] = df_analisis["monto_analisis"] * df_analisis["Factor Retorno"]
        
        total_cartera = df_analisis["monto_analisis"].sum()
        total_monto_esperado = df_analisis["Monto Esperado"].sum()
        total_retorno = df_analisis["Retorno Esperado"].sum()
        prom_riesgo = df_analisis["Riesgo (%)"].mean()
        prom_conversion = df_analisis["Probabilidad de Conversión"].mean() * 100
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "🧠 Diagnóstico Financiero"
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(33, 37, 41)
        
        # Resumen ejecutivo en texto
        resumen_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(2.5))
        resumen_frame = resumen_box.text_frame
        resumen_frame.word_wrap = True
        
        # Agregar párrafos
        p1 = resumen_frame.paragraphs[0]
        p1.text = "📊 Resumen Ejecutivo"
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0, 102, 204)
        
        puntos = [
            f"Cartera total: {formatear_monto(total_cartera)}",
            f"Conversión esperada: {formatear_monto(total_monto_esperado)} ({(total_monto_esperado/total_cartera*100):.1f}% de la cartera)",
            f"Retorno esperado: {formatear_monto(total_retorno)}",
            f"Riesgo promedio: {prom_riesgo:.1f}%",
            f"Conversión media: {prom_conversion:.1f}%"
        ]
        
        for punto in puntos:
            p = resumen_frame.add_paragraph()
            p.text = f"• {punto}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.level = 1
        
        # Gráfica de métricas financieras
        fig, ax = plt.subplots(figsize=(8, 2.5))
        
        categorias = ['Cartera\nTotal', 'Conversión\nEsperada', 'Retorno\nEsperado']
        valores = [total_cartera, total_monto_esperado, total_retorno]
        colores = ['#007bff', '#28a745', '#ffc107']
        
        bars = ax.bar(categorias, valores, color=colores, alpha=0.7, edgecolor='black')
        ax.set_ylabel('Monto ($)', fontsize=11)
        ax.set_title('Análisis Financiero de Cartera', fontsize=13, fontweight='bold')
        
        # Agregar valores en las barras
        for bar, valor in zip(bars, valores):
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{formatear_monto(valor)}',
                   ha='center', va='bottom', fontsize=11, fontweight='bold')
        
        plt.tight_layout()
        
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        img_stream.seek(0)
        
        slide.shapes.add_picture(img_stream, Inches(1), Inches(4.2), width=Inches(8))
    
    # === SLIDE 5: DISTRIBUCIÓN POR ESTATUS ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "📊 Distribución de Clientes por Estatus"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(33, 37, 41)
    
    # Gráfica de barras horizontales con todos los estatus
    fig, ax = plt.subplots(figsize=(8, 5))
    
    # Ordenar por cantidad
    top_10_estatus = estatus_counts.head(10)
    
    estatus_labels = [str(e)[:25] for e in top_10_estatus.index]
    cantidades = top_10_estatus.values
    
    bars = ax.barh(estatus_labels, cantidades, color='#17a2b8')
    ax.set_xlabel('Cantidad de Clientes', fontsize=11)
    ax.set_title('Top 10 Estatus (por cantidad)', fontsize=13, fontweight='bold')
    ax.invert_yaxis()
    
    # Agregar valores
    for bar, cantidad in zip(bars, cantidades):
        ax.text(bar.get_width(), bar.get_y() + bar.get_height()/2,
               f' {int(cantidad)}',
               va='center', fontsize=10, fontweight='bold')
    
    plt.tight_layout()
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    plt.close()
    img_stream.seek(0)
    
    slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.2), width=Inches(8.4))
    
    # === SLIDE 6: DISTRIBUCIÓN POR SUCURSALES ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "🏢 Distribución por Sucursales"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(33, 37, 41)
    
    # Contar clientes por sucursal
    sucursal_counts = df_cli["sucursal"].fillna("Sin sucursal").value_counts()
    
    # Crear gráfica de barras
    fig, ax = plt.subplots(figsize=(8, 5))
    
    sucursal_labels = [str(s)[:30] for s in sucursal_counts.index]
    cantidades = sucursal_counts.values
    
    bars = ax.barh(sucursal_labels, cantidades, color='#6f42c1')
    ax.set_xlabel('Cantidad de Clientes', fontsize=11)
    ax.set_title('Clientes por Sucursal', fontsize=13, fontweight='bold')
    ax.invert_yaxis()
    
    # Agregar valores
    for bar, cantidad in zip(bars, cantidades):
        ax.text(bar.get_width(), bar.get_y() + bar.get_height()/2,
               f' {int(cantidad)}',
               va='center', fontsize=10, fontweight='bold')
    
    plt.tight_layout()
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    plt.close()
    img_stream.seek(0)
    
    slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.2), width=Inches(8.4))
    
    # === SLIDE 7: DISTRIBUCIÓN POR ASESORES ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "👤 Distribución por Asesores"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(33, 37, 41)
    
    # Contar clientes por asesor
    asesor_counts = df_cli["asesor"].fillna("Sin asesor").value_counts()
    
    # Crear gráfica de barras (mostrar top 10 si hay muchos)
    fig, ax = plt.subplots(figsize=(8, 5))
    
    top_asesores = asesor_counts.head(10)
    asesor_labels = [str(a)[:30] for a in top_asesores.index]
    cantidades = top_asesores.values
    
    bars = ax.barh(asesor_labels, cantidades, color='#fd7e14')
    ax.set_xlabel('Cantidad de Clientes', fontsize=11)
    ax.set_title('Top 10 Asesores (por cantidad de clientes)', fontsize=13, fontweight='bold')
    ax.invert_yaxis()
    
    # Agregar valores
    for bar, cantidad in zip(bars, cantidades):
        ax.text(bar.get_width(), bar.get_y() + bar.get_height()/2,
               f' {int(cantidad)}',
               va='center', fontsize=10, fontweight='bold')
    
    plt.tight_layout()
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    plt.close()
    img_stream.seek(0)
    
    slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.2), width=Inches(8.4))
    
    # Guardar presentación en BytesIO
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    
    return pptx_stream.getvalue()

def get_base64_image(image_path):
    """Convierte imagen a base64 para embedding en HTML"""
    import base64
    with open(image_path, "rb") as f:
        return base64.b64encode(f.read()).decode()

# Header profesional simplificado
def render_professional_header():
    """Renderiza un header profesional usando solo componentes nativos de Streamlit"""
    u = current_user()
    user_name = u.get('user') or u.get('email') if u else "Usuario"
    current_time = datetime.now().strftime("%d/%m/%Y %H:%M")
    
    # Container principal para el header
    with st.container():
        # Crear columnas para centrar el contenido
        col1, col2, col3 = st.columns([1, 4, 1])
        
        with col2:
            # Título principal con emoji
            st.markdown("""
            <div style="text-align: center; padding: 20px 0;">
                <h1 style="
                    color: #FFD41D; 
                    font-size: 2.5rem; 
                    font-weight: 700; 
                    margin: 0;
                    text-shadow: 0 2px 4px rgba(0,0,0,0.1);
                ">💼 KAPITALIZA</h1>
            </div>
            """, unsafe_allow_html=True)
            
            # Subtítulo
            st.markdown("""
            <div style="text-align: center; margin-bottom: 20px;">
                <p style="
                    color: #2D3748; 
                    font-size: 1rem; 
                    font-weight: 600; 
                    text-transform: uppercase;
                    letter-spacing: 1px;
                    margin: 0;
                ">Customer Relationship Management</p>
            </div>
            """, unsafe_allow_html=True)
        
        # Separador visual nativo
        st.divider()
        
        # Información de usuario usando componentes nativos
        col_user1, col_user2, col_user3 = st.columns([2, 2, 2])
        
        with col_user2:
            # Mostrar bienvenida solo la primera vez por sesión
            show_once_info("welcome", f"👋 **Bienvenido, {user_name}**  \n📅 {current_time}")

def run_with_progress(func, steps: list, *args, **kwargs):
    """Ejecuta función con barra de progreso"""
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    total_steps = len(steps)
    for i, step in enumerate(steps):
        progress = (i + 1) / total_steps
        progress_bar.progress(progress)
        status_text.text(f"⚙️ {step}...")
        
    result = func(*args, **kwargs)
    
    progress_bar.progress(1.0)
    status_text.text("✅ Completado!")
    time.sleep(0.5)
    progress_bar.empty()
    status_text.empty()
    
    return result

# Paths and data dirs
DATA_DIR = Path("data")
DATA_DIR.mkdir(parents=True, exist_ok=True)
DOCS_DIR = DATA_DIR / "docs"
DOCS_DIR.mkdir(parents=True, exist_ok=True)
CLIENTES_CSV = DATA_DIR / "clientes.csv"
CLIENTES_XLSX = DATA_DIR / "clientes.xlsx"

# === CONFIGURACIÓN GOOGLE SHEETS ===
USE_GSHEETS = True   # pon False si quieres trabajar sólo local
GSHEET_ID      = "10_xueUKm0O1QwOK1YtZI-dFZlNdKVv82M2z29PfM9qk"
GSHEET_TAB     = "clientes"    # tu pestaña principal
GSHEET_HISTTAB = "historial"   # tu pestaña de historial

# Pestañas de catálogos
GSHEET_SUCURSALES_TAB = "sucursales"
GSHEET_ASESORES_TAB = "asesores"  
GSHEET_ESTATUS_TAB = "estatus"
GSHEET_SEGUNDO_ESTATUS_TAB = "segundo_estatus"


# Opcional: pega aquí el contenido JSON del service account si prefieres no usar el archivo
# Si la variable está vacía (""), se seguirá leyendo `service_account.json` desde disco.
SERVICE_ACCOUNT_JSON_STR = ""
# CACHING para gspread: minimizar auth / apertura repetida durante reruns
_GS_CREDS = None
_GS_GC = None
_GS_SH = None
_GS_WS_CACHE: dict = {}

# Variables globales para caché de worksheets con timestamp
_GS_WS_CACHE_TIME = {}

# Variables globales para caché de clientes
_CLIENTES_CACHE = None
_CLIENTES_CACHE_TIME = 0

# Variables globales para caché de historial
_HISTORIAL_CACHE = None
_HISTORIAL_CACHE_TIME = 0

# Variables globales para caché de usuarios
_USUARIOS_CACHE = None
_USUARIOS_CACHE_TIME = 0

# Variables globales para caché de catálogos
_CATALOGS_CACHE = {
    "sucursales": None,
    "estatus": None,
    "segundo_estatus": None,
    "asesores": None
}
_CATALOGS_CACHE_TIME = {
    "sucursales": 0,
    "estatus": 0,
    "segundo_estatus": 0,
    "asesores": 0
}

def _gs_credentials():
    """Carga credenciales desde Streamlit secrets - Versión mejorada para Streamlit Cloud"""
    global _GS_CREDS
    if _GS_CREDS is not None:
        return _GS_CREDS
    try:
        # 1) Streamlit Secrets (PRIMARIO para Streamlit Cloud)
        if hasattr(st, "secrets"):
            secrets_dict = dict(st.secrets)
            
            # Buscar en diferentes estructuras posibles
            sa_info = None
            
            # Opción 1: Secrets directamente en el nivel raíz
            if all(k in secrets_dict for k in ["type", "project_id", "private_key_id", "private_key"]):
                sa_info = secrets_dict
            # Opción 2: Dentro de key 'service_account'
            elif 'service_account' in secrets_dict:
                sa_info = dict(secrets_dict['service_account'])
            
            if sa_info:
                # Asegurar que la private_key tenga los saltos de línea correctos
                if "private_key" in sa_info:
                    sa_info["private_key"] = sa_info["private_key"].replace("\\n", "\n")
                
                scopes = ["https://www.googleapis.com/auth/spreadsheets"]
                _GS_CREDS = Credentials.from_service_account_info(sa_info, scopes=scopes)
                return _GS_CREDS
        
        # 2) Fallback para desarrollo local
        try:
            with open("service_account.json", "r", encoding="utf-8") as f:
                import json
                sa_info = json.load(f)
                if "private_key" in sa_info:
                    sa_info["private_key"] = sa_info["private_key"].replace("\\n", "\n")
                scopes = ["https://www.googleapis.com/auth/spreadsheets"]
                _GS_CREDS = Credentials.from_service_account_info(sa_info, scopes=scopes)
                return _GS_CREDS
        except FileNotFoundError:
            pass
            
        st.error("❌ No se pudieron cargar las credenciales de Google Sheets")
        return None
        
    except Exception as e:
        st.error(f"❌ Error en autenticación Google Sheets: {str(e)}")
        return None

def _gs_open_worksheet(tab_name: str, force_reload: bool = False):
    """Versión con caché temporal para evitar recargas innecesarias"""
    global _GS_GC, _GS_SH, _GS_WS_CACHE, _GS_WS_CACHE_TIME
    
    # Verificar si el caché es reciente (menos de 5 segundos)
    import time
    now = time.time()
    cache_duration = 5  # segundos
    
    if not force_reload and tab_name in _GS_WS_CACHE:
        if tab_name in _GS_WS_CACHE_TIME:
            if (now - _GS_WS_CACHE_TIME[tab_name]) < cache_duration:
                return _GS_WS_CACHE[tab_name]
    
    try:
        creds = _gs_credentials()
        if creds is None:
            return None

        if _GS_GC is None:
            _GS_GC = gspread.authorize(creds)
        
        if _GS_SH is None:
            try:
                _GS_SH = _GS_GC.open_by_key(GSHEET_ID)
            except Exception:
                return None

        try:
            ws = _GS_SH.worksheet(tab_name)
        except gspread.exceptions.WorksheetNotFound:
            ws = _GS_SH.add_worksheet(title=tab_name, rows="5000", cols="50")

        _GS_WS_CACHE[tab_name] = ws
        _GS_WS_CACHE_TIME[tab_name] = now
        return ws
        
    except Exception:
        return None

def limpiar_cache_gsheets():
    """Limpia todos los cachés de Google Sheets para forzar recarga de datos."""
    global _GS_GC, _GS_SH, _GS_WS_CACHE, _CLIENTES_CACHE, _CLIENTES_CACHE_TIME, _GS_WS_CACHE_TIME
    global _HISTORIAL_CACHE, _HISTORIAL_CACHE_TIME, _USUARIOS_CACHE, _USUARIOS_CACHE_TIME
    global _CATALOGS_CACHE, _CATALOGS_CACHE_TIME
    
    _GS_WS_CACHE.clear()
    _GS_WS_CACHE_TIME.clear()
    _GS_GC = None
    _GS_SH = None
    _CLIENTES_CACHE = None
    _CLIENTES_CACHE_TIME = 0
    _HISTORIAL_CACHE = None
    _HISTORIAL_CACHE_TIME = 0
    _USUARIOS_CACHE = None
    _USUARIOS_CACHE_TIME = 0
    
    # Limpiar catálogos
    for key in _CATALOGS_CACHE:
        _CATALOGS_CACHE[key] = None
        _CATALOGS_CACHE_TIME[key] = 0
    
    st.cache_data.clear()
    if 'gs_load_msg_shown' in st.session_state:
        del st.session_state['gs_load_msg_shown']
    if 'gs_first_load' in st.session_state:
        del st.session_state['gs_first_load']

def limpiar_cache_usuarios():
    """Limpia el caché de usuarios para forzar recarga"""
    global _USUARIOS_CACHE, _USUARIOS_CACHE_TIME
    _USUARIOS_CACHE = None
    _USUARIOS_CACHE_TIME = 0

def find_logo_path() -> Path | None:
    # Buscar logo en data/ (logo.png, logo.jpg) o en data/logo subfolder
    candidates = [DATA_DIR / "logo.png", DATA_DIR / "logo.jpg", DATA_DIR / "logo.jpeg"]
    for p in candidates:
        if p.exists():
            return p
    # buscar en carpeta docs o raíz
    altp = DATA_DIR / "logo.png"
    if altp.exists():
        return altp
    return None

SUCURSALES_FILE = DATA_DIR / "sucursales.json"

def load_sucursales() -> list:
    """
    Carga la lista de sucursales desde Google Sheets (prioritario) o JSON local (respaldo).
    """
    # Valores por defecto
    defaults = ["TOXQUI", "COLOKTE", "KAPITALIZA"]
    
    # Intentar cargar desde Google Sheets primero
    if USE_GSHEETS:
        try:
            gsheet_data = load_catalog_from_gsheet(GSHEET_SUCURSALES_TAB, defaults)
            if gsheet_data:
                # Sincronizar con archivo local
                try:
                    SUCURSALES_FILE.write_text(json.dumps(gsheet_data, ensure_ascii=False, indent=2), encoding="utf-8")
                except Exception:
                    pass
                return gsheet_data
        except Exception:
            pass
    
    # Respaldo: cargar desde archivo local
    try:
        if SUCURSALES_FILE.exists():
            data = json.loads(SUCURSALES_FILE.read_text(encoding="utf-8"))
            if isinstance(data, list):
                # Normalizar a strings y limpiar
                return [str(x).strip() for x in data if str(x).strip()]
    except Exception:
        pass
    
    # Si todo falla, usar valores por defecto y crearlos
    try:
        SUCURSALES_FILE.write_text(json.dumps(defaults, ensure_ascii=False, indent=2), encoding="utf-8")
        # También sincronizar con Google Sheets si está disponible
        if USE_GSHEETS:
            sync_catalog_to_gsheet("sucursales", defaults, GSHEET_SUCURSALES_TAB)
    except Exception:
        pass
    return defaults

def save_sucursales(lst: list):
    try:
        clean = [str(x).strip() for x in lst if str(x).strip()]
        SUCURSALES_FILE.write_text(json.dumps(clean, ensure_ascii=False, indent=2), encoding="utf-8")
        # Limpiar cache relacionado
        _cache_data.pop("sucursales", None)
        _cache_timestamp.pop("sucursales", None)
        # Sincronizar con Google Sheets
        if USE_GSHEETS:
            sync_catalog_to_gsheet("sucursales", clean, GSHEET_SUCURSALES_TAB)
    except Exception:
        pass

# Inicializar lista de sucursales desde disco
SUCURSALES = load_sucursales()
ESTATUS_FILE = DATA_DIR / "estatus.json"
SEGUNDO_ESTATUS_FILE = DATA_DIR / "segundo_estatus.json"

def load_estatus() -> list:
    """
    Carga la lista de estatus desde Google Sheets (prioritario) o JSON local (respaldo).
    """
    defaults = ["DISPERSADO","EN ONBOARDING","PENDIENTE CLIENTE","PROPUESTA","PENDIENTE DOC","REC SOBREENDEUDAMIENTO","REC NO CUMPLE POLITICAS","REC EDAD"]
    
    # Intentar cargar desde Google Sheets primero
    if USE_GSHEETS:
        try:
            gsheet_data = load_catalog_from_gsheet(GSHEET_ESTATUS_TAB, defaults)
            if gsheet_data:
                # Sincronizar con archivo local
                try:
                    ESTATUS_FILE.write_text(json.dumps(gsheet_data, ensure_ascii=False, indent=2), encoding="utf-8")
                except Exception:
                    pass
                return gsheet_data
        except Exception:
            pass
    
    # Respaldo: cargar desde archivo local
    try:
        if ESTATUS_FILE.exists():
            data = json.loads(ESTATUS_FILE.read_text(encoding="utf-8"))
            if isinstance(data, list):
                return [str(x).strip() for x in data if str(x).strip()]
    except Exception:
        pass
    
    # Si todo falla, usar valores por defecto
    try:
        ESTATUS_FILE.write_text(json.dumps(defaults, ensure_ascii=False, indent=2), encoding="utf-8")
        # También sincronizar con Google Sheets si está disponible
        if USE_GSHEETS:
            sync_catalog_to_gsheet("estatus", defaults, GSHEET_ESTATUS_TAB)
    except Exception:
        pass
    return defaults

# Cache simple para mejorar rendimiento
_cache_data = {}
_cache_timestamp = {}

def get_cached_data(key: str, loader_func, cache_duration: int = 30):
    """Cache simple con expiración en segundos"""
    import time
    now = time.time()
    
    # Verificar si el cache es válido
    if (key in _cache_data and 
        key in _cache_timestamp and 
        (now - _cache_timestamp[key]) < cache_duration):
        return _cache_data[key]
    
    # Cargar datos frescos
    data = loader_func()
    _cache_data[key] = data
    _cache_timestamp[key] = now
    return data

def clear_cache():
    """Limpiar todo el cache"""
    global _cache_data, _cache_timestamp
    _cache_data.clear()
    _cache_timestamp.clear()

def save_estatus(lst: list):
    try:
        clean = [str(x).strip() for x in lst if str(x).strip()]
        ESTATUS_FILE.write_text(json.dumps(clean, ensure_ascii=False, indent=2), encoding="utf-8")
        # Limpiar cache relacionado
        _cache_data.pop("estatus", None)
        _cache_timestamp.pop("estatus", None)
        # Sincronizar con Google Sheets
        if USE_GSHEETS:
            sync_catalog_to_gsheet("estatus", clean, GSHEET_ESTATUS_TAB)
    except Exception:
        pass

def load_segundo_estatus() -> list:
    """
    Carga la lista de segundo estatus desde Google Sheets (prioritario) o JSON local (respaldo).
    """
    defaults = ["","DISPERSADO","EN ONBOARDING","PEND.ACEPT.CLIENTE","APROB.CON PROPUESTA","PEND.DOC.PARA EVALUACION","RECH.SOBREENDEUDAMIENTO","RECH. TIPO PENSION","RECH.EDAD"]
    
    # Intentar cargar desde Google Sheets primero
    if USE_GSHEETS:
        try:
            gsheet_data = load_catalog_from_gsheet(GSHEET_SEGUNDO_ESTATUS_TAB, defaults)
            if gsheet_data:
                # Sincronizar con archivo local
                try:
                    SEGUNDO_ESTATUS_FILE.write_text(json.dumps(gsheet_data, ensure_ascii=False, indent=2), encoding="utf-8")
                except Exception:
                    pass
                return gsheet_data
        except Exception:
            pass
    
    # Respaldo: cargar desde archivo local
    try:
        if SEGUNDO_ESTATUS_FILE.exists():
            data = json.loads(SEGUNDO_ESTATUS_FILE.read_text(encoding="utf-8"))
            if isinstance(data, list):
                return [str(x).strip() for x in data if str(x).strip() or x == ""]
    except Exception:
        pass
    
    # Si todo falla, usar valores por defecto
    try:
        SEGUNDO_ESTATUS_FILE.write_text(json.dumps(defaults, ensure_ascii=False, indent=2), encoding="utf-8")
        # También sincronizar con Google Sheets si está disponible
        if USE_GSHEETS:
            sync_catalog_to_gsheet("segundo_estatus", defaults, GSHEET_SEGUNDO_ESTATUS_TAB)
    except Exception:
        pass
    return defaults

def save_segundo_estatus(lst: list):
    try:
        clean = [str(x).strip() for x in lst if (str(x).strip() or x == "")]
        SEGUNDO_ESTATUS_FILE.write_text(json.dumps(clean, ensure_ascii=False, indent=2), encoding="utf-8")
        # Limpiar cache relacionado
        _cache_data.pop("segundo_estatus", None)
        _cache_timestamp.pop("segundo_estatus", None)
        # Sincronizar con Google Sheets
        if USE_GSHEETS:
            sync_catalog_to_gsheet("segundo_estatus", clean, GSHEET_SEGUNDO_ESTATUS_TAB)
    except Exception:
        pass

# === FUNCIONES DE SINCRONIZACIÓN CON GOOGLE SHEETS ===

def sync_catalog_to_gsheet(catalog_name: str, catalog_data: list, sheet_tab: str):
    """Sincroniza un catálogo local a Google Sheets"""
    if not USE_GSHEETS:
        return
    try:
        ws = _gs_open_worksheet(sheet_tab)
        if ws is None:
            return
        
        # Crear DataFrame con el catálogo
        df_catalog = pd.DataFrame({"valor": catalog_data})
        
        # Limpiar la hoja y escribir los datos
        ws.clear()
        ws.update("A1", [["valor"]])  # Encabezado
        if not df_catalog.empty:
            rows = df_catalog.values.tolist()
            ws.append_rows(rows, value_input_option="RAW")
            
    except Exception:
        pass  # Fallar silenciosamente

def load_catalog_from_gsheet(sheet_tab: str, default_values: list = None, force_reload: bool = False) -> list:
    """
    Carga un catálogo desde Google Sheets con caché inteligente
    force_reload: True para forzar recarga desde Google Sheets
    """
    global _CATALOGS_CACHE, _CATALOGS_CACHE_TIME
    
    if not USE_GSHEETS:
        return default_values or []
    
    # Determinar qué catálogo es basándose en sheet_tab
    catalog_name = None
    if "sucursal" in sheet_tab.lower():
        catalog_name = "sucursales"
    elif "segundo" in sheet_tab.lower() or "2" in sheet_tab.lower():
        catalog_name = "segundo_estatus"
    elif "estatus" in sheet_tab.lower() or "status" in sheet_tab.lower():
        catalog_name = "estatus"
    elif "asesor" in sheet_tab.lower():
        catalog_name = "asesores"
    
    # Si tenemos el nombre del catálogo, usar caché
    if catalog_name:
        import time
        now = time.time()
        cache_duration = 600  # 10 minutos (los catálogos cambian muy raramente)
        
        # Usar caché si es reciente y no se fuerza recarga
        if not force_reload and _CATALOGS_CACHE.get(catalog_name) is not None:
            if (now - _CATALOGS_CACHE_TIME.get(catalog_name, 0)) < cache_duration:
                return _CATALOGS_CACHE[catalog_name].copy()
    
    try:
        ws = _gs_open_worksheet(sheet_tab, force_reload=force_reload)
        if ws is None:
            # Si hay caché antiguo, usarlo
            if catalog_name and _CATALOGS_CACHE.get(catalog_name) is not None:
                return _CATALOGS_CACHE[catalog_name].copy()
            return default_values or []
        
        # Obtener todos los valores
        data = ws.get_all_records()
        if data:
            values = [str(row.get("valor", "")).strip() for row in data]
            # Filtrar valores vacíos
            values = [v for v in values if v]
            result = values if values else (default_values or [])
        else:
            result = default_values or []
        
        # Actualizar caché si identificamos el catálogo
        if catalog_name:
            import time
            _CATALOGS_CACHE[catalog_name] = result.copy()
            _CATALOGS_CACHE_TIME[catalog_name] = time.time()
        
        return result
            
    except Exception:
        # Si hay error y existe caché, usarlo
        if catalog_name and _CATALOGS_CACHE.get(catalog_name) is not None:
            return _CATALOGS_CACHE[catalog_name].copy()
        pass
    
    return default_values or []

def sync_sucursales_to_gsheet():
    """Sincroniza sucursales a Google Sheets"""
    sync_catalog_to_gsheet("sucursales", SUCURSALES, GSHEET_SUCURSALES_TAB)

def sync_estatus_to_gsheet():
    """Sincroniza estatus a Google Sheets"""
    sync_catalog_to_gsheet("estatus", ESTATUS_OPCIONES, GSHEET_ESTATUS_TAB)

def sync_segundo_estatus_to_gsheet():
    """Sincroniza segundo estatus a Google Sheets"""
    sync_catalog_to_gsheet("segundo_estatus", SEGUNDO_ESTATUS_OPCIONES, GSHEET_SEGUNDO_ESTATUS_TAB)

def sync_asesores_to_gsheet():
    """Sincroniza asesores únicos desde la base de clientes a Google Sheets"""
    try:
        df_cli = cargar_clientes()
        asesores_unicos = sorted(list(set([
            asesor.strip() for asesor in df_cli["asesor"].fillna("").astype(str).tolist() 
            if asesor.strip() and asesor.strip() != "(Sin asesor)"
        ])))
        sync_catalog_to_gsheet("asesores", asesores_unicos, GSHEET_ASESORES_TAB)
    except Exception:
        pass

def load_catalogs_from_gsheet():
    """Carga todos los catálogos desde Google Sheets y actualiza las variables globales"""
    global SUCURSALES, ESTATUS_OPCIONES, SEGUNDO_ESTATUS_OPCIONES
    
    try:
        # Cargar desde Google Sheets con fallback a valores actuales
        new_sucursales = load_catalog_from_gsheet(GSHEET_SUCURSALES_TAB, SUCURSALES)
        new_estatus = load_catalog_from_gsheet(GSHEET_ESTATUS_TAB, ESTATUS_OPCIONES) 
        new_segundo_estatus = load_catalog_from_gsheet(GSHEET_SEGUNDO_ESTATUS_TAB, SEGUNDO_ESTATUS_OPCIONES)
        
        # Actualizar variables globales si hay datos válidos
        if new_sucursales:
            SUCURSALES = new_sucursales
            save_sucursales(SUCURSALES)  # Sincronizar con archivo local
            
        if new_estatus:
            ESTATUS_OPCIONES = new_estatus
            save_estatus(ESTATUS_OPCIONES)  # Sincronizar con archivo local
            
        if new_segundo_estatus:
            SEGUNDO_ESTATUS_OPCIONES = new_segundo_estatus
            save_segundo_estatus(SEGUNDO_ESTATUS_OPCIONES)  # Sincronizar con archivo local
            
    except Exception:
        pass  # Mantener valores actuales si falla

# Inicializar catálogos desde disco (con posible actualización desde Google Sheets)
ESTATUS_OPCIONES = load_estatus()
SEGUNDO_ESTATUS_OPCIONES = load_segundo_estatus()

DOC_CATEGORIAS = {
    "estado_cuenta": ["pdf", "jpg", "jpeg", "png"],
    "buro_credito":  ["pdf", "jpg", "jpeg", "png"],
    "solicitud":     ["pdf", "docx", "jpg", "jpeg", "png"],
    "contrato":      ["pdf", "docx", "jpg", "jpeg", "png"],  # visible si estatus = en dispersión
    "otros":         ["pdf", "docx", "xlsx", "jpg", "jpeg", "png"],
}

# ---------- Helpers ----------
SAFE_NAME_RE = re.compile(r"[^A-Za-z0-9._\\-áéíóúÁÉÍÓÚñÑ ]+")

def sort_df_by_dates(df: pd.DataFrame) -> pd.DataFrame:
    """
    Ordena el DataFrame por las columnas de fecha si existen ('fecha_ingreso', 'fecha_dispersion', 'ts').
    Si ninguna existe, retorna el DataFrame sin cambios.
    Maneja formatos de fecha MM/DD/YYYY y DD/MM/YYYY automáticamente.
    """
    df = df.copy()
    date_cols = [col for col in ["fecha_ingreso", "fecha_dispersion", "ts"] if col in df.columns]
    for col in date_cols:
        try:
            df[col] = parse_dates_flexible(df[col])
        except Exception:
            pass
    if date_cols:
        return df.sort_values(date_cols, ascending=True, na_position="last").reset_index(drop=True)
    return df

def parse_dates_flexible(date_series: pd.Series) -> pd.Series:
    """
    Parsea fechas de manera flexible, manejando formatos MM/DD/YYYY y DD/MM/YYYY.
    Retorna una Serie de datetime o NaT para valores inválidos.
    """
    try:
        # Intentar formato americano (MM/DD/YYYY) primero
        result = pd.to_datetime(date_series, format='%m/%d/%Y', errors='coerce')
        # Si quedan valores NaT, intentar formato europeo (DD/MM/YYYY)
        mask_nat = result.isna()
        if mask_nat.any():
            result.loc[mask_nat] = pd.to_datetime(date_series.loc[mask_nat], format='%d/%m/%Y', errors='coerce')
        # Fallback al parser automático
        mask_nat = result.isna()
        if mask_nat.any():
            result.loc[mask_nat] = pd.to_datetime(date_series.loc[mask_nat], errors='coerce')
        return result
    except Exception:
        # Fallback completo al parser automático
        return pd.to_datetime(date_series, errors='coerce')

def safe_name(s: str) -> str:
    if s is None:
        return ""
    s = str(s).strip()
    s = SAFE_NAME_RE.sub("_", s)
    s = re.sub(r"\s+", " ", s)
    return s[:150]

# NEW: normalización y búsqueda de asesor existente
def _norm_key(s: str) -> str:
    s = (s or "")
    s = str(s).strip()
    s = re.sub(r"\s+", " ", s)
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))
    # usar casefold() en lugar de lower() para una comparación Unicode más robusta
    return s.casefold()

def find_matching_asesor(name: str, df: pd.DataFrame) -> str:
    """
    Si name coincide (normalizado) con algún 'asesor' ya presente en df -> retorna la forma registrada.
    Si no hay coincidencia, retorna name limpio con capitalización de palabras (o '' si vacío).
    """
    name = (name or "").strip()
    if not name:
        return ""
    name_key = _norm_key(name)
    # buscar en el dataframe por la clave normalizada
    for a in df["asesor"].fillna("").unique():
        if not str(a).strip():
            continue
        if _norm_key(a) == name_key:
            return a  # usar la forma ya existente
    # si no existe, devolver una versión "limpia" con Title Case (mínima transformación)
    return " ".join(w.capitalize() for w in name.split())


# ----- Document helpers para manejo de archivos de clientes -----
def carpeta_docs_cliente(cid: str) -> Path:
    """
    Retorna la carpeta donde se almacenan los documentos para el cliente `cid`.
    Crea la carpeta si no existe.
    """
    # Preferir usar el nombre del cliente como carpeta (más legible).
    # Si no hay nombre, o no se encuentra, fallback al id.
    try:
        nombre = get_nombre_by_id(cid) or ""
    except Exception:
        nombre = ""

    name_safe = safe_name(nombre) if nombre else ""
    id_safe = safe_name(str(cid))

    # Si hay un nombre válido, usar/crear esa carpeta. Si existe la carpeta por id,
    # intentar migrar su contenido a la carpeta por nombre para conservar documentos.
    if name_safe:
        name_folder = DOCS_DIR / name_safe
        id_folder = DOCS_DIR / id_safe
        try:
            if id_folder.exists() and id_folder.is_dir() and not name_folder.exists():
                # mover la carpeta entera para preservar archivos previos
                shutil.move(str(id_folder), str(name_folder))
        except Exception:
            # si la migración falla, no bloquear: seguiremos usando/creando name_folder
            pass
        folder = name_folder
    else:
        folder = DOCS_DIR / id_safe

    folder.mkdir(parents=True, exist_ok=True)
    return folder

def canonicalize_from_catalog(
    raw: str,
    catalog: list[str],
    extra_synonyms: dict[str, str] | None = None,
    min_ratio: float = 0.90
) -> str:
    """
    Devuelve el valor 'raw' mapeado al elemento 'canónico' del catálogo más similar:
    - Igualdad exacta tras normalizar (ignora acentos/case/espacios)
    - Sinónimos explícitos (opcional)
    - 'Fuzzy' por similitud (difflib) con umbral min_ratio
    Si no encuentra nada suficientemente parecido → devuelve 'raw' tal cual.
    """
    s = (raw or "").strip()
    if not s:
        return s

    key = _norm_key(s)

    # 1) match exacto normalizado
    for opt in catalog:
        if _norm_key(opt) == key:
            return opt

    # 2) sinónimos opcionales (mapa: "en revision" -> "EN REVISIÓN")
    if extra_synonyms:
        for k, v in extra_synonyms.items():
            if _norm_key(k) == key:
                # devolver el canónico si existe en catálogo; si no, el sinónimo
                for opt in catalog:
                    if _norm_key(opt) == _norm_key(v):
                        return opt
                return v

    # 3) fuzzy: el más parecido por ratio
    best, best_r = None, 0.0
    for opt in catalog:
        r = difflib.SequenceMatcher(None, key, _norm_key(opt)).ratio()
        if r > best_r:
            best_r, best = r, opt

    if best and best_r >= min_ratio:
        return best

    return s


from googleapiclient.http import MediaIoBaseUpload
import io

def crear_carpeta_cliente_drive(cliente_id, cliente_nombre=""):
    """Crea o encuentra la carpeta del cliente en Google Drive."""
    if not st.session_state.drive_creds:
        return None
    
    drive_service = build("drive", "v3", credentials=st.session_state.drive_creds)
    
    # Buscar carpeta principal "CRM Kapitaliza"
    query = "name='CRM Kapitaliza' and mimeType='application/vnd.google-apps.folder' and trashed=false"
    results = drive_service.files().list(q=query, fields="files(id, name)").execute()
    
    if not results.get('files'):
        # Crear carpeta principal
        folder_metadata = {
            'name': 'CRM Kapitaliza',
            'mimeType': 'application/vnd.google-apps.folder'
        }
        main_folder = drive_service.files().create(body=folder_metadata, fields='id').execute()
        main_folder_id = main_folder.get('id')
    else:
        main_folder_id = results['files'][0]['id']
    
    # Buscar carpeta del cliente - usar solo el nombre (sin ID)
    if cliente_nombre:
        # Limpiar el nombre para evitar problemas con caracteres especiales
        import re
        folder_name = re.sub(r'[<>:"/\\|?*]', '_', cliente_nombre.strip())
        folder_name = folder_name[:100]  # Limitar longitud para evitar problemas
    else:
        folder_name = f"Cliente_{cliente_id}"
    
    # Buscar si ya existe una carpeta con este nombre
    query = f"name='{folder_name}' and mimeType='application/vnd.google-apps.folder' and '{main_folder_id}' in parents and trashed=false"
    results = drive_service.files().list(q=query, fields="files(id, name)").execute()
    
    if not results.get('files'):
        # Crear carpeta del cliente
        folder_metadata = {
            'name': folder_name,
            'mimeType': 'application/vnd.google-apps.folder',
            'parents': [main_folder_id]
        }
        client_folder = drive_service.files().create(body=folder_metadata, fields='id').execute()
        return client_folder.get('id')
    else:
        return results['files'][0]['id']

def subir_a_drive(uploaded_file, cliente_id=None, cliente_nombre=""):
    """Sube un archivo a Google Drive en la carpeta del cliente y devuelve el enlace público."""
    if not st.session_state.drive_creds:
        st.warning("⚠ Conecta tu cuenta de Google Drive primero.")
        return None

    drive_service = build("drive", "v3", credentials=st.session_state.drive_creds)
    
    # Si se especifica cliente, subir a su carpeta
    parent_folder_id = None
    if cliente_id:
        parent_folder_id = crear_carpeta_cliente_drive(cliente_id, cliente_nombre)
    
    file_metadata = {"name": uploaded_file.name}
    if parent_folder_id:
        file_metadata["parents"] = [parent_folder_id]
    
    media = MediaIoBaseUpload(io.BytesIO(uploaded_file.getbuffer()), mimetype=uploaded_file.type)
    file = drive_service.files().create(body=file_metadata, media_body=media, fields="id, webViewLink").execute()
    
    # Hacer el archivo público para visualización
    try:
        permission = {
            'type': 'anyone',
            'role': 'reader'
        }
        drive_service.permissions().create(fileId=file.get('id'), body=permission).execute()
    except:
        pass  # Si no se puede hacer público, continuar
    
    return file.get("webViewLink")


def subir_docs(cid: str, files, prefijo: str = "", usar_drive: bool = True) -> list:
    """
    Guarda una lista de archivos subidos por Streamlit en la carpeta del cliente.
    Si usar_drive=True y hay credenciales de Drive, sube a Google Drive.
    Sino, guarda localmente como antes.
    `files` puede ser una lista de UploadedFile o similar; cada objeto debe exponer `.name` y `.read()` / `.getbuffer()`.
    `prefijo` se antepone al nombre del archivo en disco.
    NO escribe en historial; retorna la lista de nombres guardados o enlaces de Drive.
    """
    if not cid:
        return []
    
    # Determinar si usar Google Drive
    use_drive = usar_drive and st.session_state.get('drive_creds') is not None
    
    # Asegurar que `files` sea iterable (Streamlit acepta single file o lista)
    if files is None:
        return []
    files_iter = files if hasattr(files, '__iter__') and not isinstance(files, (bytes, bytearray)) else [files]

    # Obtener nombre del cliente para la carpeta de Drive
    cliente_nombre = ""
    if use_drive:
        try:
            df_clientes = cargar_clientes()
            # La columna se llama 'id' no 'cliente_id'
            cliente_info = df_clientes[df_clientes['id'] == cid]
            if not cliente_info.empty:
                cliente_nombre = cliente_info.iloc[0]['nombre']
                st.info(f"🏷️ Creando carpeta para: {cliente_nombre}")
        except Exception as e:
            st.warning(f"⚠️ No se pudo obtener el nombre del cliente: {str(e)}")
            pass
    
    resultados = []
    
    if use_drive:
        # Subir a Google Drive - mostrar mensaje solo la primera vez
        show_once_info("drive_upload", "Subiendo documentos a Google Drive...")
        exitosos = 0
        errores = 0
        
        for f in files_iter:
            try:
                fname = getattr(f, "name", None) or getattr(f, "filename", None) or "uploaded"
                target_name = safe_name(f"{prefijo}{fname}")
                
                # Crear un objeto temporal con el nombre correcto
                import io
                temp_file = io.BytesIO(f.getbuffer())
                temp_file.name = target_name
                temp_file.type = getattr(f, 'type', 'application/octet-stream')
                temp_file.getbuffer = lambda: f.getbuffer()
                
                drive_link = subir_a_drive(temp_file, cid, cliente_nombre)
                if drive_link:
                    resultados.append(f"[Drive] {target_name}: {drive_link}")
                    exitosos += 1
                else:
                    errores += 1
            except Exception as e:
                errores += 1
                
        # Mostrar resumen al final
        if exitosos > 0:
            st.success(f"✅ {exitosos} documento(s) subido(s) a Google Drive")
        if errores > 0:
            st.error(f"❌ {errores} documento(s) con errores")
    else:
        # Guardar localmente (método original)
        folder = carpeta_docs_cliente(cid)
        
        # Primero: leer todo el contenido en memoria de forma segura
        to_write = []  # list of tuples (target_name, bytes)
        for f in files_iter:
            try:
                fname = getattr(f, "name", None) or getattr(f, "filename", None) or "uploaded"
                target_name = safe_name(f"{prefijo}{fname}")
                data = None
                if hasattr(f, "getbuffer"):
                    try:
                        data = f.getbuffer()
                    except Exception:
                        data = None
                if data is None and hasattr(f, "read"):
                    try:
                        data = f.read()
                    except Exception:
                        data = None
                if data is None:
                    continue
                if isinstance(data, memoryview):
                    data = data.tobytes()
                # ensure bytes
                if isinstance(data, str):
                    data = data.encode("utf-8")
                if not isinstance(data, (bytes, bytearray)):
                    try:
                        data = bytes(data)
                    except Exception:
                        continue
                to_write.append((target_name, data))
            except Exception:
                continue

        # Escribir en paralelo para acelerar (especialmente cuando hay varios archivos)
        saved_files = []
        try:
            import concurrent.futures
            def write_file(item):
                target_name, data = item
                target_path = folder / target_name
                try:
                    with open(target_path, "wb") as out:
                        out.write(data)
                    return target_name
                except Exception:
                    return None

            with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
                results = list(executor.map(write_file, to_write))
                saved_files = [r for r in results if r is not None]
        except ImportError:
            # Fallback secuencial si no hay concurrent.futures
            for target_name, data in to_write:
                target_path = folder / target_name
                try:
                    with open(target_path, "wb") as out:
                        out.write(data)
                    saved_files.append(target_name)
                except Exception:
                    pass
        resultados = saved_files
    
    return resultados


def listar_docs_cliente(cid: str):
    """
    Lista los archivos asociados a un cliente (Path objects), ordenados por nombre.
    Retorna lista vacía si no existe carpeta.
    """
    try:
        nombre = get_nombre_by_id(cid) or ""
    except Exception:
        nombre = ""

    name_safe = safe_name(nombre) if nombre else ""
    id_safe = safe_name(str(cid))

    # Preferir la carpeta por nombre si existe; si no, revisar la carpeta por id.
    if name_safe:
        folder = DOCS_DIR / name_safe
        if folder.exists() and folder.is_dir():
            return sorted([p for p in folder.iterdir() if p.is_file()], key=lambda p: p.name)

    folder_id = DOCS_DIR / id_safe
    if folder_id.exists() and folder_id.is_dir():
        return sorted([p for p in folder_id.iterdir() if p.is_file()], key=lambda p: p.name)

    return []


def nuevo_id_cliente(df: pd.DataFrame) -> str:
    """
    Genera un nuevo ID de cliente único con prefijo 'C' basado en los IDs existentes del DataFrame.
    Si no encuentra IDs del formato C<number>, comienza en C1000.
    """
    base_id = 1000
    try:
        if df is not None and not df.empty and "id" in df.columns:
            nums = []
            for x in df["id"].astype(str):
                if not x:
                    continue
                m = re.match(r"^C(\d+)$", str(x).strip())
                if m:
                    try:
                        nums.append(int(m.group(1)))
                    except Exception:
                        continue
            if nums:
                base_id = max(nums) + 1
            else:
                # fallback: avoid collision con filas existentes
                base_id = base_id + len(df)
    except Exception:
        base_id = base_id
    return f"C{base_id}"

def get_nombre_by_id(cid: str) -> str:
    """Retorna el nombre del cliente por id de forma segura ('' si no existe)."""
    try:
        if cid is None or cid == "":
            return ""
        if 'id' not in df_cli.columns or df_cli.empty:
            return ""
        sel = df_cli.loc[df_cli['id'] == cid, 'nombre']
        if sel is None or sel.empty:
            return ""
        return str(sel.values[0])
    except Exception:
        return ""

def get_field_by_id(cid: str, field: str) -> str:
    """Retorna el valor de `field` para el cliente `cid` de forma segura ('' si no existe)."""
    try:
        if cid is None or cid == "":
            return ""
        if 'id' not in df_cli.columns or df_cli.empty:
            return ""
        if field not in df_cli.columns:
            return ""
        sel = df_cli.loc[df_cli['id'] == cid, field]
        if sel is None or sel.empty:
            return ""
        return str(sel.values[0])
    except Exception:
        return ""

# --- NEW: búsquedas rápidas y cacheadas (preindexado) ---
@st.cache_data(show_spinner=False)
def build_text_index(options: list[str]):
	opts = [str(o) for o in options]
	norms = [_norm_key(o) for o in opts]
	tokens = [set(n.split()) for n in norms]
	inv = {}
	for i, toks in enumerate(tokens):
		for t in toks:
			inv.setdefault(t, set()).add(i)
	initials = ["".join(w[0] for w in n.split() if w) for n in norms]
	buckets = {}
	for i, n in enumerate(norms):
		b = n[:1]
		buckets.setdefault(b, []).append(i)
	return {"opts": opts, "norms": norms, "tokens": tokens, "inv": inv, "initials": initials, "buckets": buckets}

import re as _re

# --- ROBUST SEARCH (reemplaza fast_search) ---
def _parse_query(q: str):
    """
    Soporta:
      - AND por espacios
      - OR por comas (cada parte es un grupo AND)
      - Frases exactas entre "comillas"
      - Exclusiones con -token o !token
      - Prefijos con asterisco: vent*  (== "empieza por vent")
    """
    q = (q or "").strip()
    if not q:
        return []

    parts = [p.strip() for p in q.split(",") if p.strip()]  # OR
    groups = []
    for part in parts:
        phrases = [_norm_key(m) for m in _re.findall(r'"([^"]+)"', part)]
        base = _re.sub(r'"[^"]+"', " ", part)

        req, excl = [], []
        for t in [t for t in _re.split(r"\s+", base) if t]:
            neg = t.startswith("-") or t.startswith("!")
            tt = t[1:] if neg else t
            tt = _norm_key(tt)
            if not tt:
                continue
            (excl if neg else req).append(tt)

        groups.append({"req": req, "phrases": phrases, "exclude": excl})
    return groups

def _score_match(opt_norm: str, opt_tokens: set[str], opt_initials: str, group: dict) -> tuple[bool, float]:
    # exclusiones
    for ex in group["exclude"]:
        ex_base = ex.rstrip("*")
        if any(t.startswith(ex_base) for t in opt_tokens) or ex_base in opt_norm:
            return False, 0.0

    score = 0.0

    # frases exactas (todas)
    for ph in group["phrases"]:
        if ph in opt_norm:
            score += 3.0
        else:
            return False, 0.0

    # requisitos (todas)
    for req in group["req"]:
        is_prefix = req.endswith("*")
        base = req.rstrip("*")
        hit = False

        if base in opt_tokens:                       # token exacto
            score += 2.0; hit = True
        elif any(t.startswith(base) for t in opt_tokens):  # prefijo
            score += 1.6 if is_prefix else 1.4; hit = True
        elif base in opt_norm:                       # substring
            score += 1.2; hit = True
        elif opt_initials.startswith(base):          # iniciales
            score += 1.0; hit = True
        else:
            # fuzzy contra todo el texto normalizado (tolerancia a typos)
            ratio = difflib.SequenceMatcher(None, base, opt_norm).ratio()
            if ratio >= 0.82:
                score += 0.8; hit = True

        if not hit:
            return False, 0.0

    return True, score

def robust_search(q: str, idx: dict, limit: int | None = None) -> list[str]:
    """
    Búsqueda determinista y tolerante:
      - AND (espacios), OR (comas), "frases", -exclusiones, prefijo*
      - Acentos/case ignorados · fuzzy para typos
      - Fallback seguro si no hay matches
    """
    if not q:
        return idx["opts"]

    groups = _parse_query(q)
    if not groups:
        return idx["opts"]

    scored = []
    for i, opt_norm in enumerate(idx["norms"]):
        ok_any = False
        best = 0.0
        for g in groups:
            matched, sc = _score_match(opt_norm, idx["tokens"][i], idx["initials"][i], g)
            if matched:
                ok_any = True
                best = max(best, sc)
        if ok_any:
            best += min(0.5, len(idx["opts"][i]) / 200.0)  # bonus pequeño estable
            scored.append((best, i))

    if not scored:
        # fallback: similitud global contra la query (normalizada)
        q_norm = _norm_key(_re.sub(r'"', "", q))
        pool = idx["norms"]
        close = difflib.get_close_matches(q_norm, pool, n=min(12, len(pool)), cutoff=0.6)
        ids = []
        for val in close:
            try:
                ids.append(pool.index(val))
            except ValueError:
                pass
        out = [idx["opts"][j] for j in ids] or idx["opts"]
        return out[:limit] if limit else out

    scored.sort(key=lambda x: (-x[0], x[1]))
    out = [idx["opts"][i] for _, i in scored]
    return out[:limit] if limit else out

# --- /ROBUST SEARCH ---


def stable_multiselect(
    *,
    title: str,
    idx: dict,
    state_key: str,
    search_key: str,
    help_txt: str,
    all_options: list[str],
    on_all,
    on_none,
    min_len: int = 1,
    display_inline: bool = False,
):
    """
    Multiselect estable con popover/expander:
      - Estado lógico en st.session_state[state_key]
      - Estado visual del widget en st.session_state[w_<state_key>]
      - Botones: Todos / Ninguno / Añadir / Reemplazar / Limpiar / Nueva búsqueda
      - Sincronización widget ← estado en acciones para evitar que el widget pise tu selección.
    """
    # Contenedor
    # Contenedor: por defecto popover/expander, pero si display_inline==True renderizamos directo en sidebar
    if display_inline:
        pop = st.sidebar
    else:
        try:
            pop = st.sidebar.popover(f"{title} · {len(st.session_state.get(state_key, []))}/{len(all_options)}")
        except Exception:
            pop = st.sidebar.expander(f"{title} · {len(st.session_state.get(state_key, []))}/{len(all_options)}", expanded=False)

    # Safety: ensure `pop` is a Streamlit container (DeltaGenerator-like); fallback if not
    try:
        if not hasattr(pop, "fragment_id_queue"):
            pop = st.sidebar.expander(f"{title} · {len(st.session_state.get(state_key, []))}/{len(all_options)}", expanded=False)
    except Exception:
        try:
            pop = st.sidebar.expander(f"{title} · {len(st.session_state.get(state_key, []))}/{len(all_options)}", expanded=False)
        except Exception:
            pop = st.sidebar

    # Estado lógico inicial (una sola vez)
    st.session_state.setdefault(state_key, [o for o in all_options])

    # Widget key separado
    wkey = f"w_{state_key}"
    chk_key = f"{state_key}_all_chk"

    def _set_all():
        st.session_state[state_key] = [o for o in all_options]
        try:
            on_all()
        except Exception:
            pass

    def _set_none():
        st.session_state[state_key] = []
        try:
            on_none()
        except Exception:
            pass

    with pop:
        # Checkbox 'Todas' que determina si está todo seleccionado
        is_now_all = set(st.session_state.get(state_key, [])) == set(all_options)
        if chk_key in st.session_state:
            checked = st.checkbox("Todas", key=chk_key)
        else:
            st.session_state.setdefault(chk_key, is_now_all)
            checked = st.checkbox("Todas", value=is_now_all, key=chk_key)

        # Si el checkbox cambió, actualizar estado lógico
        if checked and not is_now_all:
            _set_all()
        elif (not checked) and is_now_all:
            _set_none()

        # Mostrar multiselect desplegado (sin buscador ni botones)
        opts = [o for o in all_options]
        selected_now = [o for o in st.session_state.get(state_key, []) if o in all_options]
        if wkey in st.session_state:
            sel = st.multiselect("", options=opts, key=wkey, label_visibility="collapsed", help=help_txt)
        else:
            sel = st.multiselect("", options=opts, default=selected_now, key=wkey, label_visibility="collapsed", help=help_txt)

        # Sincronizar widget → estado lógico
        if sel != st.session_state.get(state_key):
            st.session_state[state_key] = [o for o in sel if o in all_options]
            # sincronizar checkbox
            st.session_state[chk_key] = (set(st.session_state[state_key]) == set(all_options))

def _is_dispersion(estatus: str) -> bool:
    e = _norm_key(estatus)
    return e in {_norm_key("DISPERSADO"), _norm_key("EN DISPERSIÓN"), _norm_key("EN DISPERSION")}

# === Helper: multiselección con estética de selectbox ===
# ...existing code...
def selectbox_multi(label: str, options: list[str], state_key: str) -> list[str]:
    opts = [str(o) for o in options]
    opts = list(dict.fromkeys(opts))

    st.session_state.setdefault(state_key, opts.copy())
    selected = [o for o in st.session_state[state_key] if o in opts]
    all_selected = (set(selected) == set(opts))

    # usar popover del sidebar para mejor comportamiento de reruns
    try:
        pop = st.sidebar.popover(label)
    except Exception:
        pop = st.sidebar.expander(label, expanded=False)

    # Safety: if pop is not a Streamlit container (unexpected), fallback to expander or sidebar
    try:
        if not hasattr(pop, "fragment_id_queue"):
            pop = st.sidebar.expander(label, expanded=False)
    except Exception:
        try:
            pop = st.sidebar.expander(label, expanded=False)
        except Exception:
            pop = st.sidebar

    wkey = f"{state_key}_ms"
    chk_key = f"{state_key}_all"

    def _on_checkbox():
        if st.session_state.get(chk_key):
            st.session_state[state_key] = opts.copy()
            st.session_state[wkey] = opts.copy()
        else:
            st.session_state[state_key] = []
            st.session_state[wkey] = []
        st.session_state["_filters_token"] = st.session_state.get("_filters_token", 0) + 1

    def _on_ms_change():
        st.session_state[state_key] = [o for o in st.session_state.get(wkey, []) if o in opts]
        st.session_state[chk_key] = set(st.session_state[state_key]) == set(opts)
        st.session_state["_filters_token"] = st.session_state.get("_filters_token", 0) + 1

    with pop:
        # encabezado resumido
        c1, c2 = st.columns([0.85, 0.15])
        with c1:
            st.caption(label)
            if set(st.session_state.get(state_key, [])) == set(opts):
                st.write("(Todas)")
            elif not st.session_state.get(state_key):
                st.write("— Ninguna —")
            elif len(st.session_state.get(state_key, [])) <= 3:
                st.write(", ".join(st.session_state.get(state_key, [])))
            else:
                st.write(f"{len(st.session_state.get(state_key, []))} seleccionadas")
        with c2:
            st.write("")  # espacio para el botón pequeño

        checked = st.checkbox("Seleccionar todas", value=all_selected, key=chk_key, on_change=_on_checkbox)
        # multiselect con callback para sincronizar y forzar rerun
        if checked:
            st.multiselect("", options=opts, default=opts, disabled=True, label_visibility="collapsed", key=wkey)
        else:
            st.multiselect("", options=opts, default=selected, label_visibility="collapsed", key=wkey, on_change=_on_ms_change)

    return st.session_state[state_key]

# ---------- Sidebar (filtros + acciones) ----------
# Columnas esperadas en el CSV / DataFrame de clientes
COLUMNS = [
    "id","nombre","sucursal","asesor","fecha_ingreso","fecha_dispersion",
    "estatus","monto_propuesta","monto_final","segundo_estatus","observaciones",
    "score","telefono","correo","analista","fuente"
]

def cargar_clientes(force_reload: bool = False) -> pd.DataFrame:
    """
    Lee primero de Google Sheets con caché inteligente
    force_reload: True para forzar recarga desde Google Sheets
    """
    global _CLIENTES_CACHE, _CLIENTES_CACHE_TIME
    
    import time
    now = time.time()
    cache_duration = 3  # segundos
    
    # Usar caché si es reciente y no se fuerza recarga
    if not force_reload and _CLIENTES_CACHE is not None:
        if (now - _CLIENTES_CACHE_TIME) < cache_duration:
            return _CLIENTES_CACHE.copy()
    
    def _ensure_cols(df: pd.DataFrame) -> pd.DataFrame:
        df = df.copy().fillna("")
        for c in COLUMNS:
            if c not in df.columns:
                df[c] = ""
        return df[[c for c in COLUMNS if c in df.columns]]

    # 1) Intentar Google Sheets
    if USE_GSHEETS:
        try:
            ws = _gs_open_worksheet(GSHEET_TAB, force_reload=force_reload)
            if ws is None:
                raise Exception("No connection")
                
            df = get_as_dataframe(ws, evaluate_formulas=True, dtype=str, header=0).dropna(how="all")
            if df is None or df.empty:
                df = pd.DataFrame(columns=COLUMNS)
            else:
                # Mostrar mensaje solo al iniciar sesión
                if 'gs_first_load' not in st.session_state:
                    st.session_state['gs_first_load'] = True
                    show_once_success("gsheets_load", f"Datos cargados desde Google Sheets: {len(df)} registros")
                
            df = df.fillna("")
            for c in COLUMNS:
                if c not in df.columns:
                    df[c] = ""
            
            result = df[COLUMNS].astype(str).fillna("")
            
            # Actualizar caché
            _CLIENTES_CACHE = result.copy()
            _CLIENTES_CACHE_TIME = now
            
            return result
        except Exception as e:
            if 'gs_first_load' not in st.session_state:
                st.warning(f"⚠️ No se pudo cargar desde Google Sheets, usando datos locales")

    # 2) Fallback a archivos locales
    try:
        if CLIENTES_XLSX.exists():
            df = pd.read_excel(CLIENTES_XLSX, dtype=str).fillna("")
            result = _ensure_cols(df)
            _CLIENTES_CACHE = result.copy()
            _CLIENTES_CACHE_TIME = now
            return result
    except Exception:
        pass

    try:
        if CLIENTES_CSV.exists():
            df = pd.read_csv(CLIENTES_CSV, dtype=str).fillna("")
            result = _ensure_cols(df)
            _CLIENTES_CACHE = result.copy()
            _CLIENTES_CACHE_TIME = now
            return result
    except Exception:
        pass

    return pd.DataFrame(columns=COLUMNS)

def guardar_clientes(df: pd.DataFrame):
    """Guarda la base y actualiza caché"""
    global _CLIENTES_CACHE, _CLIENTES_CACHE_TIME
    
    try:
        if df is None:
            return

        for c in COLUMNS:
            if c not in df.columns:
                df[c] = ""
        df_to_save = df[[c for c in COLUMNS if c in df.columns]].copy().fillna("").astype(str)

        # CSV (local)
        df_to_save.to_csv(CLIENTES_CSV, index=False, encoding="utf-8")

        # XLSX (respaldo)
        try:
            engine = None
            try:
                import xlsxwriter
                engine = "xlsxwriter"
            except Exception:
                try:
                    import openpyxl
                    engine = "openpyxl"
                except Exception:
                    engine = None

            if engine:
                with pd.ExcelWriter(CLIENTES_XLSX, engine=engine) as writer:
                    df_to_save.to_excel(writer, index=False, sheet_name="Clientes")
        except Exception:
            pass

        # Actualizar caché inmediatamente
        import time
        _CLIENTES_CACHE = df_to_save.copy()
        _CLIENTES_CACHE_TIME = time.time()

        # Google Sheets (async, sin bloquear)
        if USE_GSHEETS:
            try:
                guardar_clientes_gsheet_append(df_to_save)
            except Exception:
                pass

    except Exception as e:
        try:
            st.error(f"Error guardando clientes: {e}")
        except Exception:
            pass

# --- Helpers para GSheet append/upsert ---
def _ensure_columns(df: pd.DataFrame, cols: list[str]) -> pd.DataFrame:
    df = df.copy().fillna("")
    for c in cols:
        if c not in df.columns:
            df[c] = ""
    return df[cols].astype(str).fillna("")

def _sheet_to_df(ws) -> pd.DataFrame:
    dfsh = get_as_dataframe(ws, evaluate_formulas=True, dtype=str, header=0).dropna(how="all")
    if dfsh is None or dfsh.empty:
        return pd.DataFrame()
    return dfsh.fillna("").astype(str)

def guardar_clientes_gsheet_append(df_nuevo: pd.DataFrame):
    """Versión optimizada: solo actualiza filas modificadas"""
    if df_nuevo is None or df_nuevo.empty:
        return

    try:
        ws = _gs_open_worksheet(GSHEET_TAB)
        if ws is None:
            return
            
        df_nuevo = _ensure_columns(df_nuevo, COLUMNS)

        # Asegurar encabezado
        try:
            header_row = ws.row_values(1)
            header_norm = [str(h).strip() for h in header_row]
            if header_norm[:len(COLUMNS)] != COLUMNS:
                ws.update("A1", [COLUMNS])
        except Exception:
            ws.update("A1", [COLUMNS])

        df_actual = _sheet_to_df(ws)
        
        # Si la hoja está vacía, usar set_with_dataframe (una sola vez)
        if df_actual.empty:
            try:
                set_with_dataframe(ws, df_nuevo, include_index=False, include_column_header=True, resize=True)
            except Exception:
                pass
            return
        
        df_actual = _ensure_columns(df_actual, COLUMNS)

        # Índices por ID para detección eficiente
        idx_actual = {str(r["id"]): i for i, r in df_actual.reset_index(drop=True).iterrows() if str(r["id"]).strip() != ""}
        idx_nuevo  = {str(r["id"]): i for i, r in df_nuevo.reset_index(drop=True).iterrows() if str(r["id"]).strip() != ""}

        # 1) Nuevos: append en lote
        nuevos_ids = [i for i in idx_nuevo.keys() if i not in idx_actual]
        if nuevos_ids:
            rows_to_append = df_nuevo.loc[df_nuevo["id"].astype(str).isin(nuevos_ids), COLUMNS].values.tolist()
            try:
                ws.append_rows(rows_to_append, value_input_option="RAW")
            except Exception:
                pass

        # 2) Actualizados: batch update (más eficiente)
        comunes_ids = [i for i in idx_nuevo.keys() if i in idx_actual]
        if comunes_ids:
            updates = []
            for _id in comunes_ids:
                row_new = df_nuevo.loc[idx_nuevo[_id], COLUMNS]
                row_old = df_actual.loc[idx_actual[_id], COLUMNS]
                if not row_new.equals(row_old):
                    fila = idx_actual[_id] + 2  # +2 por encabezado (1-indexed)
                    # Calcular letra de última columna (ej: P para 16 columnas)
                    import string
                    col_letter = string.ascii_uppercase[len(COLUMNS) - 1] if len(COLUMNS) <= 26 else 'Z'
                    rango = f"A{fila}:{col_letter}{fila}"
                    updates.append({
                        "range": rango,
                        "values": [row_new.tolist()]
                    })

            # Batch update (máximo 100 por lote para evitar límites de API)
            if updates:
                try:
                    for i in range(0, len(updates), 100):
                        batch = updates[i:i+100]
                        ws.batch_update(batch, value_input_option="RAW")
                except Exception:
                    pass
                    
    except Exception:
        pass

# Función cargar_y_corregir_clientes optimizada
# Función cargar_y_corregir_clientes optimizada
def cargar_y_corregir_clientes(force_reload: bool = False) -> pd.DataFrame:
    """Carga los clientes y corrige IDs duplicados/vacíos si es necesario"""
    df_cli = cargar_clientes(force_reload=force_reload)
    
    try:
        df_fixed = _fix_missing_or_duplicate_ids(df_cli)
        try:
            changed = not df_fixed.equals(df_cli)
        except Exception:
            changed = True
        if changed:
            df_cli = df_fixed
            guardar_clientes(df_cli)
        else:
            df_cli = df_fixed
    except Exception:
        pass
    
    return df_cli

# Función para arreglar IDs duplicados/vacíos
def _fix_missing_or_duplicate_ids(df: pd.DataFrame) -> pd.DataFrame:
    """Corrige IDs vacíos o duplicados en el DataFrame"""
    if df is None or df.empty:
        return df
    
    df_fixed = df.copy()
    
    # Generar IDs únicos para registros sin ID o con ID vacío
    for idx, row in df_fixed.iterrows():
        if pd.isna(row.get('id')) or str(row.get('id')).strip() == '':
            # Generar nuevo ID único
            existing_ids = set(df_fixed['id'].dropna().astype(str))
            counter = 1000
            while f"C-{counter}" in existing_ids:
                counter += 1
            df_fixed.at[idx, 'id'] = f"C-{counter}"
    
    return df_fixed

# Funciones de historial
HIST_COLUMNS_DEFAULT = ["fecha","accion","id","nombre","detalle","usuario"]

def append_historial_gsheet(evento: dict):
    """ Agrega un registro al historial (una fila nueva) """
    if not USE_GSHEETS:
        return
        
    try:
        ws = _gs_open_worksheet(GSHEET_HISTTAB)
        if ws is None:
            return
            
        df_actual = _sheet_to_df(ws)
        if df_actual.empty:
            try:
                ws.update(values=[HIST_COLUMNS_DEFAULT], range_name="A1")
            except Exception:
                pass
        fila = [str(evento.get(col, "")) for col in HIST_COLUMNS_DEFAULT]
        try:
            ws.append_rows([fila], value_input_option="RAW")
        except Exception:
            pass
    except Exception:
        pass

# Función append_historial local
    """Corrige IDs vacíos o duplicados en el DataFrame"""
    if df is None or df.empty:
        return df
    df = df.copy()
    if "id" not in df.columns:
        df["id"] = ""

    usados = set()
    def _nuevo_id_local(_df):
        base_id = 1000
        try:
            if not _df.empty and "id" in _df.columns:
                nums = []
                for x in _df["id"].astype(str):
                    if str(x).startswith("C"):
                        try:
                            nums.append(int(str(x).lstrip("C")))
                        except Exception:
                            continue
                if nums:
                    base_id = max(nums) + 1
                else:
                    base_id = base_id + len(_df) + 1
        except Exception:
            base_id = base_id + 1
        return f"C{base_id}"

    for i in df.index:
        cur = str(df.at[i, "id"]).strip()
        if not cur or cur in usados:
            # genera ID nuevo que no choque con usados ni con el df
            nuevo = _nuevo_id_local(df)
            while nuevo in usados or (df["id"] == nuevo).any():
                try:
                    num = int(nuevo[1:]) + 1
                except Exception:
                    num = 1000
                nuevo = f"C{num}"
            df.at[i, "id"] = nuevo
            usados.add(nuevo)
        else:
            usados.add(cur)
    return df

# ---------- Historial y eliminación de clientes ----------
HISTORIAL_CSV = DATA_DIR / "historial.csv"

def cargar_historial(force_reload: bool = False) -> pd.DataFrame:
    """
    Lee el historial desde Google Sheets (prioritario) o CSV local como respaldo.
    Usa caché inteligente para evitar cargas repetitivas.
    Retorna DataFrame con columnas esperadas si no existe.
    """
    global _HISTORIAL_CACHE, _HISTORIAL_CACHE_TIME
    
    # Verificar caché (3-5 segundos de duración)
    import time
    now = time.time()
    cache_duration = 4  # segundos
    
    if not force_reload and _HISTORIAL_CACHE is not None and (now - _HISTORIAL_CACHE_TIME) < cache_duration:
        return _HISTORIAL_CACHE.copy()
    
    # Columnas estándar del historial
    cols = ["id", "nombre", "estatus_old", "estatus_new", "segundo_old", "segundo_new", "observaciones", "action", "actor", "ts"]
    
    # Intentar cargar desde Google Sheets primero
    if USE_GSHEETS:
        try:
            ws = _gs_open_worksheet(GSHEET_HISTTAB)
            if ws:
                # Obtener todos los datos de la hoja
                data = ws.get_all_records()
                if data:
                    dfh = pd.DataFrame(data).fillna("").astype(str)
                    
                    # Mapear columnas de Google Sheets a formato interno
                    # Google Sheets usa: ["fecha","accion","id","nombre","detalle","usuario"]
                    # Formato interno usa: ["id", "nombre", "estatus_old", "estatus_new", "segundo_old", "segundo_new", "observaciones", "action", "actor", "ts"]
                    
                    if len(dfh) > 0 and ('fecha' in dfh.columns or 'accion' in dfh.columns):
                        # Crear DataFrame con formato interno
                        dfh_formatted = pd.DataFrame()
                        dfh_formatted['id'] = dfh.get('id', '').astype(str)
                        dfh_formatted['nombre'] = dfh.get('nombre', '').astype(str)
                        dfh_formatted['estatus_old'] = ''  # No se almacena separadamente en GSheets
                        dfh_formatted['estatus_new'] = ''  # No se almacena separadamente en GSheets  
                        dfh_formatted['segundo_old'] = ''  # No se almacena separadamente en GSheets
                        dfh_formatted['segundo_new'] = ''  # No se almacena separadamente en GSheets
                        dfh_formatted['observaciones'] = dfh.get('detalle', '').astype(str)
                        dfh_formatted['action'] = dfh.get('accion', '').astype(str)
                        dfh_formatted['actor'] = dfh.get('usuario', '').astype(str)
                        dfh_formatted['ts'] = dfh.get('fecha', '').astype(str)
                        
                        # Corregir problemas de encoding común en Google Sheets
                        for col in ['nombre', 'observaciones', 'actor']:
                            if col in dfh_formatted.columns:
                                # Corregir caracteres con tilde
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ã©', 'é', regex=False)
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ã¡', 'á', regex=False)
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ã­', 'í', regex=False)
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ã³', 'ó', regex=False)
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ãº', 'ú', regex=False)
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ã±', 'ñ', regex=False)
                                # Corregir mayúsculas
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ã', 'Á', regex=False)
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ã‰', 'É', regex=False)
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ã"', 'Ó', regex=False)
                                dfh_formatted[col] = dfh_formatted[col].str.replace('Ãš', 'Ú', regex=False)
                                # Ñ mayúscula usando código unicode
                                dfh_formatted[col] = dfh_formatted[col].str.replace('\u00c3\u0091', 'Ñ', regex=False)
                        
                        # Asegurar todas las columnas requeridas
                        for c in cols:
                            if c not in dfh_formatted.columns:
                                dfh_formatted[c] = ""
                        
                        # Ordenar por timestamp de manera descendente (más reciente primero)
                        try:
                            dfh_formatted['_ts_sort'] = pd.to_datetime(dfh_formatted['ts'], errors='coerce')
                            dfh_formatted = dfh_formatted.sort_values('_ts_sort', ascending=False)
                            dfh_formatted = dfh_formatted.drop(columns=['_ts_sort'])
                        except Exception:
                            pass
                        
                        # Actualizar caché
                        result = dfh_formatted[cols].copy()
                        _HISTORIAL_CACHE = result.copy()
                        _HISTORIAL_CACHE_TIME = now
                        return result
        except Exception:
            pass  # Si falla Google Sheets, usar CSV local
    
    # Respaldo: cargar desde CSV local
    try:
        if HISTORIAL_CSV.exists():
            dfh = pd.read_csv(HISTORIAL_CSV, dtype=str).fillna("")
            for c in cols:
                if c not in dfh.columns:
                    dfh[c] = ""
            result = dfh[cols].copy()
            # Actualizar caché también para CSV
            _HISTORIAL_CACHE = result.copy()
            _HISTORIAL_CACHE_TIME = now
            return result
    except Exception:
        pass
    
    # Si todo falla, retornar DataFrame vacío
    result = pd.DataFrame(columns=cols)
    _HISTORIAL_CACHE = result.copy()
    _HISTORIAL_CACHE_TIME = now
    return result

def append_historial_gsheet(evento: dict):
    """
    Versión global para registrar historial en Google Sheets.
    Se crea si no existe encabezado y luego se hace append.
    Silenciosa ante cualquier excepción.
    """
    if not USE_GSHEETS:
        return
    try:
        ws = _gs_open_worksheet(GSHEET_HISTTAB)
        headers = ["fecha","accion","id","nombre","detalle","usuario"]
        try:
            existing_header = ws.row_values(1)
        except Exception:
            existing_header = []
        if not existing_header:
            try:
                ws.update("A1", [headers])
            except Exception:
                pass
        fila = [str(evento.get(col, "")) for col in headers]
        try:
            ws.append_rows([fila], value_input_option="RAW")
        except Exception:
            pass
    except Exception:
        pass

def append_historial(cid: str, nombre: str, estatus_old: str, estatus_new: str, seg_old: str, seg_new: str, observaciones: str = "", action: str = "ESTATUS MODIFICADO", actor: str | None = None):
    """
    Agrega una fila al historial de estatus (archivo CSV).
    action: 'crear'|'modificar'|'eliminar'|'importar' u otro texto libre.
    actor: nombre de usuario que realizó la acción; si no se pasa, se toma el usuario actual.
    """
    try:
        if actor is None:
            cu = current_user() or {}
            actor = cu.get("user") or cu.get("email") or "(sistema)"

        registro = {
            "id": cid,
            "nombre": nombre or "",
            "estatus_old": estatus_old or "",
            "estatus_new": estatus_new or "",
            "segundo_old": seg_old or "",
            "segundo_new": seg_new or "",
            "observaciones": observaciones or "",
            "action": action or "",
            "actor": actor or "",
            "ts": pd.Timestamp.now().isoformat()
        }
        if HISTORIAL_CSV.exists():
            dfh = cargar_historial()
            dfh = pd.concat([dfh, pd.DataFrame([registro])], ignore_index=True)
        else:
            dfh = pd.DataFrame([registro])
        dfh.to_csv(HISTORIAL_CSV, index=False, encoding="utf-8")
        # También intentar escribir en Google Sheets (modo append) si está habilitado
        if USE_GSHEETS:
            try:
                evento = {
                    "fecha": registro.get("ts", ""),
                    "accion": registro.get("action", ""),
                    "id": registro.get("id", ""),
                    "nombre": registro.get("nombre", ""),
                    "detalle": registro.get("observaciones", ""),
                    "usuario": registro.get("actor", "")
                }
                try:
                    append_historial_gsheet(evento)
                except Exception:
                    pass
            except Exception:
                pass
    except Exception:
        # no bloquear la app por errores de historial
        pass

def eliminar_cliente(cid: str, df: pd.DataFrame, borrar_historial: bool = False) -> pd.DataFrame:
    """
    Elimina al cliente del DataFrame `df`, borra su carpeta de documentos y (opcionalmente) las entradas de historial.
    Retorna el DataFrame resultante (y guarda el CSV de clientes).
    """
    try:
        if cid is None or cid == "" or df is None or df.empty or "id" not in df.columns:
            return df
        # Borrar carpeta de documentos del cliente
        try:
            folder = DOCS_DIR / safe_name(str(cid))
            if folder.exists() and folder.is_dir():
                shutil.rmtree(folder)
        except Exception:
            pass

        # Eliminar de df
        df_new = df[df["id"] != cid].reset_index(drop=True)
        guardar_clientes(df_new)

        # --- NUEVO: eliminar también de la hoja de Google Sheets (si está habilitado) ---
        if USE_GSHEETS:
            try:
                ws = _gs_open_worksheet(GSHEET_TAB)
                vals = ws.get_all_values()
                if vals and len(vals) > 0:
                    header = [str(h).strip() for h in vals[0]]
                    # buscar índice de la columna 'id' (case/acentos tolerantemente)
                    id_col = None
                    for idx, h in enumerate(header):
                        if _norm_key(h) == _norm_key("id"):
                            id_col = idx
                            break
                    if id_col is not None:
                        rows_to_delete = []
                        # recorrer filas de datos (vals[1:] corresponde a filas físicas a partir de la 2)
                        for i in range(1, len(vals)):
                            try:
                                cell = vals[i][id_col] if id_col < len(vals[i]) else ""
                            except Exception:
                                cell = ""
                            if str(cell).strip() == str(cid):
                                # i -> índice en vals; la fila en la hoja es i+1 (1-based)
                                rows_to_delete.append(i + 1)
                        # borrar de abajo hacia arriba para no invalidar índices
                        for rownum in sorted(rows_to_delete, reverse=True):
                            try:
                                ws.delete_rows(rownum)
                            except Exception:
                                # fallback no crítico: intentar limpiar la fila en lugar de borrarla
                                try:
                                    # limpiar un rango amplio (A..Z) para evitar errores si delete_rows falla
                                    ws.update(f"A{rownum}:Z{rownum}", [[""] * 26])
                                except Exception:
                                    pass
            except Exception:
                # silenciar cualquier error de GSheets para no romper la app
                pass
        # -------------------------------------------------------

        # Borrar historial asociado si se solicita
        if borrar_historial:
            try:
                if HISTORIAL_CSV.exists():
                    dfh = cargar_historial()
                    dfh = dfh[dfh["id"] != cid].reset_index(drop=True)
                    dfh.to_csv(HISTORIAL_CSV, index=False, encoding="utf-8")
            except Exception:
                pass

        return df_new
    except Exception:
        return df

# --- AUTENTICACIÓN CON ROLES (admin / member) ---
import secrets
import base64

USERS_FILE = DATA_DIR / "users.json"   # { "users":[{"user": "...", "role":"admin|member", "salt":"...", "hash":"..."}] }
# Google Sheets tab name for users
GSHEET_USERSTAB = "users"

PERMISSIONS = {
    "admin":  {"manage_users": True,  "delete_client": True},
    "member": {"manage_users": False, "delete_client": False},
}

def do_rerun():
    """Forzar rerun compatible con varias versiones de Streamlit."""
    try:
        if hasattr(st, "experimental_rerun"):
            st.experimental_rerun()
            return
    except Exception:
        pass
    try:
        from streamlit.runtime.scriptrunner import RerunException  # type: ignore
        raise RerunException("Requested rerun")
    except Exception:
        pass
    try:
        from streamlit.script_runner import RerunException as _RerunOld  # type: ignore
        raise _RerunOld("Requested rerun (old)")
    except Exception:
        pass
    st.session_state["_need_rerun"] = not st.session_state.get("_need_rerun", False)
    try:
        st.stop()
    except Exception:
        return

def _hash_pw_pbkdf2(password: str, salt_hex: str | None = None) -> tuple[str, str]:
    if not salt_hex:
        salt_hex = secrets.token_hex(16)
    salt = bytes.fromhex(salt_hex)
    dk = hashlib.pbkdf2_hmac("sha256", (password or "").encode("utf-8"), salt, 100_000)
    return salt_hex, dk.hex()

def _verify_pw(password: str, salt_hex: str, hash_hex: str) -> bool:
    _, hh = _hash_pw_pbkdf2(password, salt_hex)
    return secrets.compare_digest(hh, (hash_hex or ""))

def cargar_usuarios_gsheet(force_reload: bool = False) -> dict:
    """
    Carga usuarios desde Google Sheets con caché inteligente
    force_reload: True para forzar recarga desde Google Sheets
    """
    global _USUARIOS_CACHE, _USUARIOS_CACHE_TIME
    
    import time
    now = time.time()
    cache_duration = 300  # 5 minutos (los usuarios no cambian frecuentemente)
    
    # Usar caché si es reciente y no se fuerza recarga
    if not force_reload and _USUARIOS_CACHE is not None:
        if (now - _USUARIOS_CACHE_TIME) < cache_duration:
            return _USUARIOS_CACHE.copy()
    
    if not USE_GSHEETS:
        return {"users": []}
    
    try:
        ws = _gs_open_worksheet(GSHEET_USERSTAB, force_reload=force_reload)
        if ws is None:
            return {"users": []}
            
        df = get_as_dataframe(ws, evaluate_formulas=True, dtype=str, header=0).dropna(how="all")
        if df is None or df.empty:
            result = {"users": []}
        else:
            users = []
            for _, row in df.iterrows():
                user_data = {
                    "user": row.get("user", ""),
                    "role": row.get("role", "member"),
                    "salt": row.get("salt", ""),
                    "hash": row.get("hash", "")
                }
                # Solo agregar usuarios válidos
                if user_data["user"] and user_data["salt"] and user_data["hash"]:
                    users.append(user_data)
            result = {"users": users}
        
        # Actualizar caché
        _USUARIOS_CACHE = result.copy()
        _USUARIOS_CACHE_TIME = now
        return result
        
    except Exception as e:
        # Si hay error de API (quota exceeded), usar caché antiguo si existe
        if _USUARIOS_CACHE is not None:
            st.warning(f"⚠️ Usando caché de usuarios (API temporalmente no disponible)")
            return _USUARIOS_CACHE.copy()
        st.error(f"Error cargando usuarios desde Google Sheets: {e}")
        return {"users": []}

def guardar_usuarios_gsheet(users_data: dict):
    """
    Guarda usuarios en Google Sheets
    """
    if not USE_GSHEETS:
        return
        
    try:
        ws = _gs_open_worksheet(GSHEET_USERSTAB)
        if ws is None:
            st.error("No se pudo abrir la pestaña de usuarios en Google Sheets")
            return
            
        # Preparar datos para guardar
        users_list = users_data.get("users", [])
        if not users_list:
            # Limpiar la hoja si no hay usuarios
            try:
                ws.clear()
            except Exception:
                pass
            return
            
        # Crear DataFrame
        df = pd.DataFrame(users_list)
        
        # Asegurar columnas en orden correcto
        column_order = ["user", "role", "salt", "hash"]
        for col in column_order:
            if col not in df.columns:
                df[col] = ""
                
        df = df[column_order]
        
        # Guardar en Google Sheets
        set_with_dataframe(ws, df, include_index=False, include_column_header=True, resize=True)
        
    except Exception as e:
        st.error(f"Error guardando usuarios en Google Sheets: {e}")

def load_users() -> dict:
    """
    Carga usuarios primero desde Google Sheets, luego desde archivo local como fallback
    """
    # 1) Intentar Google Sheets primero
    if USE_GSHEETS:
        gsheet_data = cargar_usuarios_gsheet()
        if gsheet_data and gsheet_data.get("users"):
            return gsheet_data
    
    # 2) Fallback a archivo local
    try:
        if USERS_FILE.exists():
            return json.loads(USERS_FILE.read_text(encoding="utf-8"))
    except Exception:
        pass
    
    return {"users": []}

def save_users(obj: dict):
    """
    Guarda usuarios en Google Sheets Y en archivo local (backup)
    """
    # 1) Guardar en Google Sheets (primario)
    if USE_GSHEETS:
        guardar_usuarios_gsheet(obj)
    
    # 2) Guardar en archivo local (backup)
    try:
        USERS_FILE.write_text(json.dumps(obj, indent=2, ensure_ascii=False), encoding="utf-8")
    except Exception as e:
        st.error(f"Error guardando usuarios localmente: {e}")

def sync_usuarios_local_to_gsheet():
    """
    Sincroniza usuarios desde el archivo local a Google Sheets
    Útil para migración inicial
    """
    try:
        local_data = load_users()
        guardar_usuarios_gsheet(local_data)
    except Exception as e:
        st.error(f"Error sincronizando usuarios: {e}")

def get_user(identifier: str) -> dict | None:
    """Buscar usuario por username o por email (compatibilidad backward).
    identifier: lo que ingresa el usuario (username); busca en 'user' o 'email'.
    """
    ident = (identifier or "").strip().lower()
    data = load_users()
    for u in data.get("users", []):
        if u.get("user", "").lower() == ident or u.get("email", "").lower() == ident:
            return u
    return None

def add_user(username: str, password: str, role: str = "member") -> tuple[bool, str]:
    uname = (username or "").strip()
    if not uname or not password:
        return False, "Usuario y contraseña son obligatorios."
    if role not in ("admin", "member"):
        return False, "Rol inválido."
    data = load_users()
    # comprueba duplicados en 'user' y en el antiguo 'email'
    lower_uname = uname.lower()
    if any((u.get("user","") or u.get("email","" )).lower() == lower_uname for u in data.get("users", [])):
        return False, "Ese usuario ya existe."
    salt_hex, hash_hex = _hash_pw_pbkdf2(password)
    data["users"].append({"user": uname, "role": role, "salt": salt_hex, "hash": hash_hex})
    save_users(data)
    # Limpiar caché para forzar recarga
    limpiar_cache_usuarios()
    return True, "Usuario creado."

def delete_user(username: str) -> tuple[bool, str]:
    name = (username or "").strip().lower()
    if not name:
        return False, "Usuario inválido."
    data = load_users()
    users = data.get("users", [])
    for i, u in enumerate(users):
        if (u.get("user","") or u.get("email","" )).lower() == name:
            users.pop(i)
            data["users"] = users
            save_users(data)
            # Limpiar caché para forzar recarga
            limpiar_cache_usuarios()
            return True, "Usuario eliminado."
    return False, "Usuario no encontrado."

def maybe_migrate_legacy_admin():
    legacy = DATA_DIR / "admin.json"
    if legacy.exists() and not USERS_FILE.exists():
        try:
            obj = json.loads(legacy.read_text(encoding="utf-8"))
            email = (obj.get("email","") or "").strip()
            if email:
                temp_pw = base64.urlsafe_b64encode(secrets.token_bytes(9)).decode("utf-8").rstrip("=")
                ok, _ = add_user(email, temp_pw, role="admin")
                if ok:
                    st.sidebar.info(f"Admin migrado: {email}. Contraseña temporal: {temp_pw}. Inicia y cámbiala.")
                    # Sincronizar con Google Sheets
                    if USE_GSHEETS:
                        try:
                            local_data = load_users()
                            guardar_usuarios_gsheet(local_data)
                        except Exception:
                            pass
                legacy.unlink(missing_ok=True)
        except Exception:
            pass

maybe_migrate_legacy_admin()

# session state for auth
if "auth_user" not in st.session_state:
    st.session_state["auth_user"] = None  # dict: {"email":..., "role":...}

def current_user():
    return st.session_state.get("auth_user")

def is_admin():
    u = current_user()
    return bool(u and u.get("role") == "admin")

def can(action: str) -> bool:
    u = current_user()
    role = (u or {}).get("role", "member")
    return PERMISSIONS.get(role, {}).get(action, False)

# Migración inicial de usuarios locales a Google Sheets si procede
def maybe_migrate_users_to_gsheet():
    """
    Migra usuarios del archivo local a Google Sheets si es la primera vez
    """
    if not USE_GSHEETS:
        return
        
    try:
        # Verificar si ya hay usuarios en Google Sheets
        gsheet_data = cargar_usuarios_gsheet()
        if not gsheet_data.get("users"):
            # No hay usuarios en GS, migrar desde local
            local_data = load_users()
            if local_data.get("users"):
                guardar_usuarios_gsheet(local_data)
                st.sidebar.success("✅ Usuarios migrados a Google Sheets")
    except Exception as e:
        st.sidebar.warning(f"⚠️ Error migrando usuarios: {e}")

# Setup inicial: si no hay usuarios, crear primer admin
maybe_migrate_users_to_gsheet()
users_data = load_users()
# Setup inicial: si no hay usuarios, crear primer admin
users_data = load_users()
if not users_data.get("users"):
    with st.sidebar.expander("Configurar administrador", expanded=True):
        st.warning("No hay usuarios. Crea el primer administrador.")
        _user = st.text_input("Usuario admin", key="setup_user")
        _pw1 = st.text_input("Contraseña", type="password", key="setup_pw1")
        _pw2 = st.text_input("Confirmar", type="password", key="setup_pw2")
        if st.button("Crear administrador"):
            if not _user or not _pw1:
                st.error("Usuario y contraseña obligatorios.")
            elif _pw1 != _pw2:
                st.error("Las contraseñas no coinciden.")
            else:
                ok, msg = add_user(_user, _pw1, role="admin")
                if ok:
                    st.success("Administrador creado. Inicia sesión.")
                    do_rerun()
                else:
                    st.error(msg)

# Login: renderizar en placeholder y eliminarlo al iniciar sesión
if not current_user():
    login_panel = st.sidebar.empty()
    with login_panel.form("login_form", clear_on_submit=True):
        st.markdown("### Iniciar sesión")
        luser = st.text_input("Usuario", key="login_user")
        lpw = st.text_input("Contraseña", type="password", key="login_pw")
        submitted = st.form_submit_button("Entrar")

    if submitted:
        u = get_user(luser)
        if u and _verify_pw(lpw, u.get("salt",""), u.get("hash","")):
            # establecer usuario y limpiar estado sensible
            st.session_state["auth_user"] = {"user": u.get("user") or u.get("email"), "role": u["role"]}
            for _k in ("login_pw", "login_user"):
                st.session_state.pop(_k, None)
            # quitar el formulario al instante; no forzar rerun inmediato (evita pantalla en blanco)
            login_panel.empty()
            st.toast(f"Bienvenido, {st.session_state['auth_user']['user']} ({u['role']}).", icon="✅")
        else:
            st.error("Credenciales inválidas.")

    # Si aún no hay usuario autenticado, detenemos la app aquí.
    if not current_user():
        st.stop()

# Sidebar: info + logout + admin panel
u = current_user()
st.sidebar.markdown(f"**Usuario:** {u.get('user') or u.get('email')} — _{u['role']}_")
if st.sidebar.button("Cerrar sesión"):
    st.session_state["auth_user"] = None
    # Limpiar filtros y campos de login/alta para evitar que queden visibles
    for k in ("f_suc","f_est","f_ases","ases_q","suc_q","est_q",
              "login_user","login_pw",
              "new_user_user","new_user_pw1","new_user_pw2",
              "setup_user","setup_pw1","setup_pw2"):
        st.session_state.pop(k, None)
    do_rerun()

if is_admin():
    # Agregar miembro del equipo: mostrar dentro de un expander para ahorrar espacio
    with st.sidebar.expander("Agregar miembro del equipo", expanded=False):
        with st.form("add_user_form", clear_on_submit=True):
            st.caption("Agregar miembro del equipo")
            nuser = st.text_input("Usuario del miembro", key="new_user_user")
            npw1 = st.text_input("Contraseña", type="password", key="new_user_pw1")
            npw2 = st.text_input("Confirmar contraseña", type="password", key="new_user_pw2")
            nrole = st.selectbox("Rol", ["member", "admin"], index=0, help="Por defecto: member", key="new_user_role")
            submitted = st.form_submit_button("Agregar usuario")

        if submitted:
            # Validaciones y feedback efímero
            if not nuser or not npw1:
                st.toast("Usuario y contraseña obligatorios.", icon="🚫")
            elif npw1 != npw2:
                st.toast("Las contraseñas no coinciden.", icon="🚫")
            else:
                ok, msg = add_user(nuser, npw1, role=nrole)
                st.toast(msg, icon="✅" if ok else "🚫")
                if ok:
                    # Forzar rerun para refrescar el sidebar y limpiar el form
                    do_rerun()

    # Mostrar lista de usuarios opcionalmente (toggle apagado por defecto)
    # -- Sincronización Usuarios (admin) --
    with st.sidebar.expander("🔄 Sincronización Usuarios", expanded=False):
        st.caption("Sincronizar usuarios entre Google Sheets y archivo local")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("📤 A Google Sheets", key="sync_to_gs"):
                local_data = load_users()
                guardar_usuarios_gsheet(local_data)
                st.success("Usuarios sincronizados a Google Sheets")
        
        with col2:
            if st.button("📥 Desde Google Sheets", key="sync_from_gs"):
                gsheet_data = cargar_usuarios_gsheet()
                if gsheet_data.get("users"):
                    USERS_FILE.write_text(
                        json.dumps(gsheet_data, indent=2, ensure_ascii=False), 
                        encoding="utf-8"
                    )
                    st.success("Usuarios sincronizados desde Google Sheets")
                    # Forzar rerun para recargar usuarios
                    do_rerun()

    show_users = st.sidebar.checkbox("Mostrar usuarios registrados", value=False, key="admin_show_users")
    if show_users:
        data = load_users()
        if data.get("users"):
            st.sidebar.caption("Usuarios registrados")
            st.sidebar.caption("Siempre dar doble click para confirmar.")
            # Mostrar tabla con botón eliminar por fila
            for x in data["users"]:
                uname = x.get("user") or x.get("email")
                role = x.get("role")
                col1, col2 = st.sidebar.columns([3,1])
                with col1:
                    st.write(f"{uname} — {role}")
                with col2:
                    # Evitar que el admin se borre a sí mismo
                    cur = current_user() or {}
                    cur_user = (cur.get("user") or cur.get("email"))
                    if uname != cur_user:
                        # flujo de confirmación en dos pasos para evitar borrados accidentales
                        confirm_key = f"confirm_del_{uname}"
                        confirm_input_key = f"confirm_del_input_{uname}"
                        if not st.session_state.get(confirm_key, False):
                            if st.sidebar.button("Eliminar", key=f"del_{uname}"):
                                # activar el modo de confirmación
                                st.session_state[confirm_key] = True
                                do_rerun()
                        else:
                            # Usar columnas del sidebar para mantener todo horizontal y evitar wrapping
                            c1, c2, c3 = st.sidebar.columns([3,1,1])
                            with c1:
                                # solo placeholder para ahorrar espacio visual
                                st.sidebar.text_input("", key=confirm_input_key, placeholder=uname)
                            with c2:
                                if st.sidebar.button("Eliminar", key=f"confirm_eliminar_{uname}"):
                                    st.toast("Escribe el usuario para confirmar.", icon="📝")
                                    typed = st.session_state.get(confirm_input_key, "").strip()
                                    if typed == str(uname):
                                        eliminar, msg = delete_user(uname)
                                        st.toast(msg, icon="✅" if eliminar else "🚫")
                                        # limpiar estado de confirmación
                                        st.session_state.pop(confirm_key, None)
                                        st.session_state.pop(confirm_input_key, None)
                                        if eliminar:
                                            do_rerun()
                                    else:
                                        st.toast("El texto no coincide. Escribe el nombre exacto para confirmar.", icon="🚫")
                            with c3:
                                if st.sidebar.button("Cancelar", key=f"confirm_cancel_{uname}"):
                                    st.session_state.pop(confirm_key, None)
                                    st.session_state.pop(confirm_input_key, None)
                                    do_rerun()
                    else:
                        st.write("")

    # -- Sincronización de Catálogos con Google Sheets --
    st.sidebar.markdown("---")
    with st.sidebar.expander("☁️ Sincronización de Catálogos", expanded=False):
        st.caption("Sincronizar catálogos con Google Sheets")
        
        col_sync1, col_sync2 = st.columns(2)
        
        with col_sync1:
            if st.button("📤 Subir Todo", key="sync_all_to_gsheet", help="Subir todos los catálogos a Google Sheets"):
                with st.spinner("Sincronizando catálogos..."):
                    sync_sucursales_to_gsheet()
                    sync_estatus_to_gsheet()
                    sync_segundo_estatus_to_gsheet()
                    sync_asesores_to_gsheet()
                st.toast("✅ Catálogos sincronizados a Google Sheets", icon="✅")
                
        with col_sync2:
            if st.button("📥 Descargar Todo", key="sync_all_from_gsheet", help="Descargar todos los catálogos desde Google Sheets"):
                with st.spinner("Cargando catálogos..."):
                    load_catalogs_from_gsheet()
                st.toast("✅ Catálogos actualizados desde Google Sheets", icon="✅")
                st.rerun()
        
        # Sincronización individual
        st.markdown("**Sincronización individual:**")
        sync_cols = st.columns(4)
        
        with sync_cols[0]:
            if st.button("🏢", key="sync_suc", help="Sincronizar Sucursales"):
                sync_sucursales_to_gsheet()
                st.toast("✅ Sucursales sincronizadas")
                
        with sync_cols[1]:
            if st.button("👥", key="sync_ases", help="Sincronizar Asesores"):
                sync_asesores_to_gsheet()
                st.toast("✅ Asesores sincronizados")
                
        with sync_cols[2]:
            if st.button("📊", key="sync_est", help="Sincronizar Estatus"):
                sync_estatus_to_gsheet()
                st.toast("✅ Estatus sincronizados")
                
        with sync_cols[3]:
            if st.button("📈", key="sync_seg", help="Sincronizar 2° Estatus"):
                sync_segundo_estatus_to_gsheet()
                st.toast("✅ 2° Estatus sincronizados")

    # -- Gestión unificada (solo admin) -- (OPTIMIZADA)
    with st.sidebar.expander("⚙️ Gestión de Catálogos", expanded=False):
        st.caption("Administrar sucursales, asesores, estatus y segundo estatus")
        
        # Tabs para organizar mejor la gestión
        tab_suc, tab_ases, tab_est, tab_seg = st.tabs(["🏢 Sucursales", "👥 Asesores", "📊 Estatus", "📈 2° Estatus"])
        
        # === TAB SUCURSALES ===
        with tab_suc:
            # Manejar reset de campo si está marcado
            if st.session_state.get("reset_new_suc", False):
                st.session_state["new_suc"] = ""
                st.session_state["reset_new_suc"] = False
            
            # Agregar nueva sucursal
            col1, col2 = st.columns([3, 1])
            with col1:
                nueva_suc = st.text_input("Nueva sucursal:", key="new_suc", placeholder="Ej. CDMX_CENTRO")
            with col2:
                if st.button("➕", key="add_suc", help="Agregar sucursal"):
                    if nueva_suc.strip():
                        if nueva_suc.strip() not in SUCURSALES:
                            SUCURSALES.append(nueva_suc.strip())
                            save_sucursales(SUCURSALES)
                            st.toast(f"✅ Sucursal '{nueva_suc.strip()}' agregada")
                            # Marcar para reset en próximo rerun
                            st.session_state["reset_new_suc"] = True
                            st.rerun()
                        else:
                            st.toast("⚠️ Sucursal ya existe")
                    else:
                        st.toast("⚠️ Nombre vacío")
            
            # Lista de sucursales existentes
            if SUCURSALES:
                st.caption(f"Sucursales ({len(SUCURSALES)}):")
                for i, suc in enumerate(SUCURSALES):
                    col1, col2 = st.columns([4, 1])
                    with col1:
                        st.write(f"• {suc}")
                    with col2:
                        if st.button("🗑️", key=f"del_suc_{i}", help=f"Eliminar {suc}"):
                            # Verificar si está en uso
                            df_check = cargar_clientes()
                            en_uso = False
                            if not df_check.empty and 'sucursal' in df_check.columns:
                                en_uso = (df_check['sucursal'] == suc).any()
                            
                            if en_uso:
                                st.toast(f"⚠️ '{suc}' está en uso por clientes")
                            else:
                                SUCURSALES.remove(suc)
                                save_sucursales(SUCURSALES)
                                st.toast(f"✅ Sucursal '{suc}' eliminada")
                                st.rerun()
        
        # === TAB ASESORES ===
        with tab_ases:
            # Manejar reset de campo si está marcado
            if st.session_state.get("reset_new_asesor", False):
                st.session_state["new_asesor"] = ""
                st.session_state["reset_new_asesor"] = False
            
            # Agregar nuevo asesor
            col1, col2 = st.columns([3, 1])
            with col1:
                nuevo_asesor = st.text_input("Nuevo asesor:", key="new_asesor", placeholder="Ej. Juan Pérez")
            with col2:
                if st.button("➕", key="add_asesor", help="Agregar asesor"):
                    if nuevo_asesor.strip():
                        # Los asesores se manejan dinámicamente, pero podemos mantener una lista maestra
                        st.toast(f"✅ Asesor '{nuevo_asesor.strip()}' listo para usar")
                        st.session_state["reset_new_asesor"] = True
                        # Forzar actualización de filtros para incluir nuevo asesor
                        st.session_state["_filters_token"] = st.session_state.get("_filters_token", 0) + 1
                        st.rerun()
                    else:
                        st.toast("⚠️ Nombre vacío")
            
            # Mostrar asesores existentes (desde la base de datos)
            df_check = cargar_clientes()
            if not df_check.empty and 'asesor' in df_check.columns:
                asesores_existentes = df_check['asesor'].fillna("").unique()
                asesores_existentes = [a for a in asesores_existentes if a.strip()]
                if asesores_existentes:
                    st.caption(f"Asesores en uso ({len(asesores_existentes)}):")
                    for asesor in sorted(asesores_existentes):
                        st.write(f"• {asesor}")
        
        # === TAB ESTATUS ===
        with tab_est: # type: ignore
            # Manejar reset de campo si está marcado
            if st.session_state.get("reset_new_estatus", False):
                st.session_state["new_estatus"] = ""
                st.session_state["reset_new_estatus"] = False
            
            # Agregar nuevo estatus
            col1, col2 = st.columns([3, 1])
            with col1:
                nuevo_estatus = st.text_input("Nuevo estatus:", key="new_estatus", placeholder="Ej. EN_REVISION")
            with col2:
                if st.button("➕", key="add_estatus", help="Agregar estatus"):
                    if nuevo_estatus.strip():
                        if nuevo_estatus.strip() not in ESTATUS_OPCIONES:
                            ESTATUS_OPCIONES.append(nuevo_estatus.strip())
                            save_estatus(ESTATUS_OPCIONES)
                            st.toast(f"✅ Estatus '{nuevo_estatus.strip()}' agregado")
                            st.session_state["reset_new_estatus"] = True
                            st.rerun()
                        else:
                            st.toast("⚠️ Estatus ya existe")
                    else:
                        st.toast("⚠️ Nombre vacío")
            
            # Lista de estatus existentes
            if ESTATUS_OPCIONES:
                st.caption(f"Estatus ({len(ESTATUS_OPCIONES)}):")
                for i, est in enumerate(ESTATUS_OPCIONES):
                    col1, col2 = st.columns([4, 1])
                    with col1:
                        st.write(f"• {est}")
                    with col2:
                        if st.button("🗑️", key=f"del_est_{i}", help=f"Eliminar {est}"):
                            # Verificar si está en uso
                            df_check = cargar_clientes()
                            en_uso = False
                            if not df_check.empty and 'estatus' in df_check.columns:
                                en_uso = (df_check['estatus'] == est).any()
                            
                            if en_uso:
                                st.toast(f"⚠️ '{est}' está en uso por clientes")
                            else:
                                ESTATUS_OPCIONES.remove(est)
                                save_estatus(ESTATUS_OPCIONES)
                                st.toast(f"✅ Estatus '{est}' eliminado")
                                st.rerun()
        
        # === TAB SEGUNDO ESTATUS ===
        with tab_seg:
            # Manejar reset de campo si está marcado
            if st.session_state.get("reset_new_seg_estatus", False):
                st.session_state["new_seg_estatus"] = ""
                st.session_state["reset_new_seg_estatus"] = False
            
            # Agregar nuevo segundo estatus
            col1, col2 = st.columns([3, 1])
            with col1:
                nuevo_seg_estatus = st.text_input("Nuevo 2° estatus:", key="new_seg_estatus", placeholder="Ej. ALTA_PRIORIDAD")
            with col2:
                if st.button("➕", key="add_seg_estatus", help="Agregar segundo estatus"):
                    if nuevo_seg_estatus.strip():
                        if nuevo_seg_estatus.strip() not in SEGUNDO_ESTATUS_OPCIONES:
                            SEGUNDO_ESTATUS_OPCIONES.append(nuevo_seg_estatus.strip())
                            save_segundo_estatus(SEGUNDO_ESTATUS_OPCIONES)
                            st.toast(f"✅ 2° Estatus '{nuevo_seg_estatus.strip()}' agregado")
                            st.session_state["reset_new_seg_estatus"] = True
                            st.rerun()
                        else:
                            st.toast("⚠️ 2° Estatus ya existe")
                    else:
                        st.toast("⚠️ Nombre vacío")
            
            # Lista de segundo estatus existentes
            if SEGUNDO_ESTATUS_OPCIONES:
                st.caption(f"2° Estatus ({len(SEGUNDO_ESTATUS_OPCIONES)}):")
                for i, seg_est in enumerate(SEGUNDO_ESTATUS_OPCIONES):
                    col1, col2 = st.columns([4, 1])
                    with col1:
                        st.write(f"• {seg_est}")
                    with col2:
                        if st.button("🗑️", key=f"del_seg_{i}", help=f"Eliminar {seg_est}"):
                            # Verificar si está en uso
                            df_check = cargar_clientes()
                            en_uso = False
                            if not df_check.empty and 'segundo_estatus' in df_check.columns:
                                en_uso = (df_check['segundo_estatus'] == seg_est).any()
                            
                            if en_uso:
                                st.toast(f"⚠️ '{seg_est}' está en uso por clientes")
                            else:
                                SEGUNDO_ESTATUS_OPCIONES.remove(seg_est)
                                save_segundo_estatus(SEGUNDO_ESTATUS_OPCIONES)
                                st.toast(f"✅ 2° Estatus '{seg_est}' eliminado")
                                st.rerun()

# ---------- Sidebar (filtros + acciones) ----------
st.sidebar.title("👤 CRM")
st.sidebar.caption("Filtros")

# Cargar datos frescos para el sidebar
df_cli = cargar_y_corregir_clientes()

# Opciones base
SUC_LABEL_EMPTY = "(Sin sucursal)"
# CORREGIR: Manejar tanto valores vacíos ("") como None/NaN
sucursal_for_filter = df_cli["sucursal"].fillna("").replace({"": SUC_LABEL_EMPTY})
SUC_ALL  = sorted(set(SUCURSALES + [SUC_LABEL_EMPTY]))

# Recalcular campos derivados para filtros (asegurar que reflejen la versión en disco)
asesor_for_filter = df_cli["asesor"].fillna("").replace({"": "(Sin asesor)"})
# Normalizar variantes de la etiqueta "(Sin asesor)" (mayúsculas/minúsculas/espacios)
def _norm_sin_asesor_label(x: str) -> str:
    try:
        s = (x or "").strip()
    except Exception:
        s = ""
    if s.casefold() == "(sin asesor)":
        return "(Sin asesor)"
    return s

# Aplicar normalización y asegurar unicidad
ASES_ALL = sorted(list(dict.fromkeys([_norm_sin_asesor_label(x) for x in asesor_for_filter.unique().tolist()])))

# IMPORTANTE: Aplicar la misma normalización a asesor_for_filter para que coincida con ASES_ALL
asesor_for_filter_normalized = asesor_for_filter.apply(_norm_sin_asesor_label)

EST_ALL = ESTATUS_OPCIONES.copy()

# --- NEW: Fuentes para filtro (se generan dinámicamente desde la base para que nuevas fuentes aparezcan automáticamente)
fuente_for_filter = df_cli["fuente"].fillna("").replace({"": "(Sin fuente)"})
FUENTE_ALL = sorted(list(dict.fromkeys([ (str(x).strip() if str(x).strip() else "(Sin fuente)") for x in fuente_for_filter.unique().tolist() ])))

# Controles “tipo selectbox” pero multi
f_suc  = selectbox_multi("Sucursales", SUC_ALL,  "f_suc")
f_ases = selectbox_multi("Asesores",   ASES_ALL, "f_ases")
f_est  = selectbox_multi("Estatus",    EST_ALL,  "f_est")
# NEW: añadir filtro de Fuente en el sidebar
f_fuente = selectbox_multi("Fuente", FUENTE_ALL, "f_fuente")

# Validar consistencia de datos automáticamente (método seguro)
if f_ases:
    # Verificar que los asesores seleccionados existen en los datos actuales
    valid_ases = [a for a in f_ases if a in ASES_ALL]
    if len(valid_ases) != len(f_ases):
        # Solo usar los asesores válidos para el filtro, sin modificar session_state aquí
        f_ases = valid_ases
        # Mostrar advertencia opcional (removida para evitar spam)
        # st.sidebar.warning("Algunos asesores seleccionados ya no existen y fueron removidos automáticamente.")

def _reset_filters():
    try:
        st.session_state["f_suc"] = SUC_ALL.copy()
        st.session_state["f_ases"] = ASES_ALL.copy()
        st.session_state["f_est"] = EST_ALL.copy()
        st.session_state["f_suc_all"] = True
        st.session_state["f_ases_all"] = True
        st.session_state["f_est_all"] = True
        st.session_state["f_suc_ms"] = SUC_ALL.copy()
        st.session_state["f_ases_ms"] = ASES_ALL.copy()
        st.session_state["f_est_ms"] = EST_ALL.copy()
        # NEW: reset para fuente
        st.session_state.setdefault("f_fuente", FUENTE_ALL.copy())
        st.session_state["f_fuente_all"] = True
        st.session_state["f_fuente_ms"] = FUENTE_ALL.copy()
        # token opcional para forzar lógica dependiente si la usas
        st.session_state["_filters_token"] = st.session_state.get("_filters_token", 0) + 1
    except Exception:
        pass

def _force_refresh():
    """Fuerza actualización de caché y filtros para mostrar nuevos datos"""
    global _CLIENTES_CACHE, _CLIENTES_CACHE_TIME, _HISTORIAL_CACHE, _HISTORIAL_CACHE_TIME
    try:
        # Limpiar caché de clientes y historial
        _CLIENTES_CACHE = None
        _CLIENTES_CACHE_TIME = 0
        _HISTORIAL_CACHE = None
        _HISTORIAL_CACHE_TIME = 0
        # Reset filtros también
        _reset_filters()
        # Marcar que se necesita actualizar (sin llamar st.rerun() en callback)
        st.session_state["_force_refresh_requested"] = True
    except Exception:
        pass

st.sidebar.button("🔁", key="btn_reset_filters", on_click=_reset_filters)

# Botón actualizar mejorado con feedback visual
col_refresh1, col_refresh2 = st.sidebar.columns([3, 1])
with col_refresh1:
    # Calcular tiempo desde última actualización
    import time
    cache_age = int(time.time() - _CLIENTES_CACHE_TIME) if _CLIENTES_CACHE_TIME > 0 else 0
    if cache_age < 60:
        st.sidebar.caption(f"Última actualización: hace {cache_age}s")
    else:
        st.sidebar.caption(f"Última actualización: hace {cache_age//60}m")
with col_refresh2:
    if st.sidebar.button("🔄", key="btn_force_refresh", help="Recargar todo desde Google Sheets"):
        with st.spinner("Recargando datos..."):
            _force_refresh()
        st.toast("✅ Datos actualizados", icon="✅")
        st.rerun()

# Verificar si se solicitó actualización desde callback
if st.session_state.get("_force_refresh_requested", False):
    st.session_state["_force_refresh_requested"] = False
    st.rerun()

# Aplicar filtros: si no hay selección o está vacía, NO filtrar (mostrar todo)
try:
    # Debug info comentado para optimización
    # st.sidebar.write(f"**Debug Info:**")
    # st.sidebar.write(f"- Total clientes en base: {len(df_cli)}")
    # st.sidebar.write(f"- Asesores disponibles: {len(ASES_ALL)}")
    # st.sidebar.write(f"- Asesores filtrados: {len(f_ases) if f_ases else 0}")
    # if f_ases and len(f_ases) < len(ASES_ALL):
    #     st.sidebar.write(f"- Asesores seleccionados: {f_ases[:3]}{'...' if len(f_ases) > 3 else ''}")
    
    # Sucursal: si está vacío o tiene todos, no filtrar
    if not f_suc or len(f_suc) == 0 or set(f_suc) == set(SUC_ALL):
        suc_mask = pd.Series(True, index=df_cli.index)
    else:
        # Regenerar sucursal_for_filter con datos actuales para evitar desincronización  
        current_sucursal_for_filter = df_cli["sucursal"].fillna("").replace({"": SUC_LABEL_EMPTY})
        suc_mask = current_sucursal_for_filter.isin(f_suc)

    # Asesor: si está vacío o tiene todos, no filtrar
    if not f_ases or len(f_ases) == 0 or set(f_ases) == set(ASES_ALL):
        asesor_mask = pd.Series(True, index=df_cli.index)
    else:
        # Regenerar asesor_for_filter_normalized con datos actuales para evitar desincronización
        current_asesor_for_filter = df_cli["asesor"].fillna("").replace({"": "(Sin asesor)"})
        current_asesor_normalized = current_asesor_for_filter.apply(_norm_sin_asesor_label)
        asesor_mask = current_asesor_normalized.isin(f_ases)

    # Estatus: si está vacío o tiene todos, no filtrar
    if not f_est or len(f_est) == 0 or set(f_est) == set(EST_ALL):
        est_mask = pd.Series(True, index=df_cli.index)
    else:
        est_mask = df_cli["estatus"].isin(f_est)

    # Fuente: si está vacío o tiene todos, no filtrar
    try:
        if not f_fuente or len(f_fuente) == 0 or set(f_fuente) == set(FUENTE_ALL):
            fuente_mask = pd.Series(True, index=df_cli.index)
        else:
            # Regenerar fuente_for_filter con datos actuales para evitar desincronización
            current_fuente_for_filter = df_cli["fuente"].fillna("").apply(lambda x: str(x).strip() if str(x).strip() else "(Sin fuente)")
            fuente_mask = current_fuente_for_filter.isin(f_fuente)
    except Exception:
        fuente_mask = pd.Series(True, index=df_cli.index)

except Exception as e:
    # Fallback seguro: no filtrar si algo falla
    st.sidebar.error(f"Error en filtros: {e}")
    suc_mask = pd.Series(True, index=df_cli.index)
    asesor_mask = pd.Series(True, index=df_cli.index)
    est_mask = pd.Series(True, index=df_cli.index)
    fuente_mask = pd.Series(True, index=df_cli.index)

df_ver = df_cli[suc_mask & est_mask & asesor_mask & fuente_mask].copy()

# Resumen
st.sidebar.markdown("---")
st.sidebar.subheader("📊 Resumen filtrado")
st.sidebar.metric("Clientes visibles", len(df_ver))
st.sidebar.metric("Total en base", len(df_cli))

# Añadir botón para descargar Excel del resumen filtrado (df_ver)
try:
    bio = io.BytesIO()
    engine = None
    try:
        import xlsxwriter  # type: ignore
        engine = "xlsxwriter"
    except Exception:
        try:
            import openpyxl  # type: ignore
            engine = "openpyxl"
        except Exception:
            engine = None

    if engine is None:
        st.sidebar.info("Instala 'openpyxl' o 'xlsxwriter' para habilitar descarga XLSX.")
    else:
        # Preparar DataFrame a exportar (ordenado por fechas si procede)
        try:
            df_export = sort_df_by_dates(df_ver) if (isinstance(df_ver, pd.DataFrame) and not df_ver.empty) else df_ver.copy()
        except Exception:
            df_export = df_ver.copy() if isinstance(df_ver, pd.DataFrame) else pd.DataFrame()

        with pd.ExcelWriter(bio, engine=engine) as writer:
            try:
                df_export.to_excel(writer, index=False, sheet_name="Filtrados")
            except Exception:
                # fallback: intentar convertir todo a strings y volver a escribir
                try:
                    df_export.astype(str).to_excel(writer, index=False, sheet_name="Filtrados")
                except Exception:
                    pass
        bio.seek(0)
        if st.sidebar.download_button(
            "⬇️ Descargar Excel (filtrados)",
            data=bio.getvalue(),
            file_name="clientes_filtrados.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key="dl_filtrados_sidebar"
        ):
            try:
                actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                append_historial(
                "", "", "", "", "", "",
                "Descarga de Excel filtrados",
                action="DESCARGA ZIP",  # o "DOCUMENTOS", según quieras categorizarlo
                actor=actor
                )
            except Exception:
                    pass
        
        # NUEVO: Descarga Excel por pestañas de asesores
        bio_asesores = io.BytesIO()
        try:
            # Preparar datos por asesor
            df_for_asesores = df_export.copy()
            if not df_for_asesores.empty and 'asesor' in df_for_asesores.columns:
                # Normalizar asesores vacíos
                df_for_asesores['asesor'] = df_for_asesores['asesor'].fillna("(Sin asesor)").replace({"": "(Sin asesor)"})
                
                # Obtener lista de asesores únicos
                asesores_unicos = sorted(df_for_asesores['asesor'].unique().tolist())
                
                if len(asesores_unicos) > 0:
                    with pd.ExcelWriter(bio_asesores, engine=engine) as writer:
                        # Crear una pestaña por cada asesor
                        for asesor in asesores_unicos:
                            try:
                                # Filtrar datos por asesor
                                df_asesor = df_for_asesores[df_for_asesores['asesor'] == asesor].copy()
                                if not df_asesor.empty:
                                    # Limpiar nombre del asesor para usar como nombre de pestaña (máximo 31 caracteres)
                                    sheet_name = str(asesor)[:31] if asesor else "Sin_asesor"
                                    # Reemplazar caracteres problemáticos para nombres de pestaña
                                    sheet_name = sheet_name.replace("/", "_").replace("\\", "_").replace("*", "_").replace("?", "_").replace(":", "_").replace("[", "_").replace("]", "_")
                                    
                                    df_asesor.to_excel(writer, index=False, sheet_name=sheet_name)
                            except Exception:
                                # Si falla crear la pestaña, continúa con el siguiente asesor
                                continue
                    
                    bio_asesores.seek(0)
                    if st.sidebar.download_button(
                        "📊 Descargar Excel (por asesores)",
                        data=bio_asesores.getvalue(),
                        file_name="clientes_por_asesores.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key="dl_asesores_sidebar",
                        help="Descarga Excel con una pestaña por cada asesor"
                    ):
                        try:
                            actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                            append_historial(
                            "", "", "", "", "", "",
                            "Descarga de Excel por asesores",
                            action="DESCARGA ZIP ASESOR",
                            actor=actor
                            )
                        except Exception:
                            pass
        except Exception:
            # Si algo falla en la creación del Excel por asesores, no mostrar el botón
            pass
except Exception:
    # no bloquear la UI si algo falla
    pass

# ---------- Main UI con pestañas ----------  # NEW
# Header profesional
render_professional_header()

tab_dash, tab_cli, tab_docs, tab_import, tab_hist = st.tabs(
    ["📊 Dashboard", "📋 Clientes", "📎 Documentos", "📥 Importar", "🗂️ Historial"]
)

# ===== Dashboard =====
with tab_dash:
    # Cargar datos frescos para el dashboard
    df_cli = cargar_y_corregir_clientes()
    
    if df_cli.empty:
        st.info("Sin clientes aún.")
    else:
        # Botón para descargar presentación
        col_titulo, col_pptx = st.columns([4, 1])
        with col_titulo:
            st.subheader("📊 KPIs Principales")
        with col_pptx:
            # Generar PowerPoint del dashboard
            pptx_data = generar_presentacion_dashboard(df_cli)
            st.download_button(
                label="📊 Descargar",
                data=pptx_data,
                file_name=f"dashboard_kapitaliza_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                help="Descargar presentación completa con gráficas"
            )
        
        # Preparar datos para KPIs
        total_clientes = len(df_cli)
        estatus_counts = df_cli["estatus"].fillna("").value_counts()

        # Calcular KPIs principales con lógica corregida
        dispersados = estatus_counts.get("DISPERSADO", 0)

        # Clasificar automáticamente todos los estatus que empiecen con "RECH" como rechazados
        rechazados = sum([
            count for estatus, count in estatus_counts.items()
            if estatus and (estatus.startswith("RECH") or estatus.startswith("REC"))
        ])

        # === NUEVO: calcular propuestas (clientes con monto de propuesta) ===
        analisis_tmp = calcular_analisis_financiero(df_cli)
        total_propuestas = int(analisis_tmp.get('clientes_con_monto', 0))
        dispersados_con_monto = int(analisis_tmp.get('dispersados_con_monto', 0))

        # Propuestas que aún no se han dispersado
        propuestas_no_dispersadas = max(0, total_propuestas - dispersados_con_monto)
        
        # Mostrar KPIs profesionales en columnas
        kpi1, kpi2, kpi3, kpi4 = st.columns(4)
        
        with kpi1:
            st.metric(
                label="👥 Total de Clientes",
                value=total_clientes,
                help="Número total de clientes en la base de datos"
            )
        
        with kpi2:
            # Mostrar dispersados pero en relación a las propuestas (monto)
            tasa_exito = (dispersados_con_monto / total_propuestas * 100) if total_propuestas > 0 else 0
            st.metric(
                label="✅ Dispersados (Éxito)",
                value=dispersados_con_monto,
                delta=f"{tasa_exito:.1f}% de propuestas",
                help="Clientes dispersados considerando solo propuestas (monto)"
            )

        with kpi3:
            # Mostrar propuestas que aún no se han dispersado (no-dispersado / total propuestas)
            tasa_proceso = (propuestas_no_dispersadas / total_propuestas * 100) if total_propuestas > 0 else 0
            st.metric(
                label="⏳ Propuestas no dispersadas",
                value=propuestas_no_dispersadas,
                delta=f"{tasa_proceso:.1f}% de propuestas",
                help="Propuestas con monto que aún no se han dispersado (no-dispersado/total propuestas)"
            )
        
        with kpi4:
            tasa_rechazo = (rechazados / total_clientes * 100) if total_clientes > 0 else 0
            st.metric(
                label="❌ Rechazados",
                value=rechazados,
                delta=f"{tasa_rechazo:.1f}% del total",
                help="Clientes rechazados por diversos motivos"
            )
        
        # 💹 Top Estatus por Monto
        st.markdown("##### 💹 Top Estatus por Monto")
        
        # Calcular métricas financieras para el ranking
        analisis_financiero = calcular_analisis_financiero(df_cli)
        
        # Mostrar total de presupuesto general
        total_presupuesto = analisis_financiero['total_propuesto']
        st.markdown(f"**Total Presupuesto General: {formatear_monto(total_presupuesto)}**")
        
        if not analisis_financiero['montos_por_estatus'].empty:
            # Filtrar solo estatus con monto > 0
            estatus_con_monto = analisis_financiero['montos_por_estatus'][
                analisis_financiero['montos_por_estatus'][('monto_propuesta_num', 'sum')] > 0
            ]
            
            # Obtener top estatus por monto propuesto
            top_estatus = estatus_con_monto.sort_values(
                ('monto_propuesta_num', 'sum'), ascending=False
            ).head(5)
            
            col1, col2 = st.columns(2)
            
            # Dividir en dos columnas para mostrar mejor con texto más pequeño
            for i, (estatus, data) in enumerate(top_estatus.iterrows()):
                monto_total = data[('monto_propuesta_num', 'sum')]
                cantidad = data[('monto_propuesta_num', 'count')]
                promedio = data[('monto_propuesta_num', 'mean')]
                
                with col1 if i % 2 == 0 else col2:
                    # Usar markdown para texto más pequeño
                    st.markdown(f"""
                    <div style="
                        border: 1px solid #e1e5e9;
                        border-radius: 8px;
                        padding: 12px;
                        margin: 4px 0;
                        background-color: #f8f9fa;
                    ">
                        <div style="font-size: 12px; color: #6c757d; font-weight: 500;">
                            {estatus}
                        </div>
                        <div style="font-size: 18px; font-weight: bold; color: #212529; margin: 4px 0;">
                            {formatear_monto(monto_total)}
                        </div>
                        <div style="font-size: 11px; color: #6c757d;">
                            {int(cantidad)} clientes • Promedio: {formatear_monto(promedio)}
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
        else:
            st.info("No hay datos con montos registrados para mostrar ranking.")
        
        st.markdown("---")
        
        #  ANÁLISIS FINANCIERO



        
        st.markdown("---")
        
        # �🔄 TABS SECUNDARIAS PARA ANÁLISIS DETALLADO
        dash_tab1, dash_tab2, dash_tab3 = st.tabs([
            "📊 Por Estatus", 
            "📅 Por Fecha", 
            "🏢 Por Sucursal/Asesor"
        ])
        
        # 📊 TAB 1: POR ESTATUS
        with dash_tab1:
            st.subheader("Análisis por Estatus")
            
            # Solo mostrar estatus con valores > 0 (más limpio)
            # Primero limpiar y normalizar estatus vacíos o nulos
            df_temp = df_cli.copy()
            df_temp["estatus_clean"] = df_temp["estatus"].fillna("(Sin estatus)").replace("", "(Sin estatus)")
            df_temp["estatus_clean"] = df_temp["estatus_clean"].apply(lambda x: x.strip() if str(x).strip() else "(Sin estatus)")
            
            estatus_data = df_temp["estatus_clean"].value_counts()
            estatus_data = estatus_data[estatus_data > 0]  # Filtrar solo valores > 0
            
            if estatus_data.empty:
                st.info("No hay datos de estatus para mostrar.")
            else:
                # Convertir a DataFrame y agregar porcentajes
                estatus_df = estatus_data.reset_index()
                estatus_df.columns = ["estatus", "cantidad"]
                estatus_df["porcentaje"] = (estatus_df["cantidad"] / total_clientes * 100).round(1)
                
                # Ordenar por cantidad (descendente)
                estatus_df = estatus_df.sort_values("cantidad", ascending=False)
                
                col1, col2 = st.columns([1, 2])
                
                with col1:
                    # Métricas con porcentajes del total
                    st.markdown("**Distribución detallada:**")
                    for _, row in estatus_df.iterrows():
                        st.metric(
                            label=row["estatus"],
                            value=f"{int(row['cantidad'])} ({row['porcentaje']}%)"
                        )
                
                with col2:
                    # Gráfico horizontal (más fácil de leer)
                    chart = alt.Chart(estatus_df).mark_bar(
                        cornerRadiusTopRight=4,
                        cornerRadiusBottomRight=4
                    ).encode(
                        x=alt.X("cantidad:Q", 
                               axis=alt.Axis(title="Cantidad de Clientes"),
                               scale=alt.Scale(domain=[0, estatus_df["cantidad"].max() * 1.1])),
                        y=alt.Y("estatus:N", 
                               sort="-x", 
                               axis=alt.Axis(title="Estatus")),
                        color=alt.Color("estatus:N", 
                                       scale=alt.Scale(scheme="category20"),
                                       legend=None),
                        tooltip=[
                            alt.Tooltip("estatus:N", title="Estatus"),
                            alt.Tooltip("cantidad:Q", title="Cantidad"),
                            alt.Tooltip("porcentaje:Q", title="Porcentaje (%)", format=".1f")
                        ]
                    ).properties(
                        height=max(300, len(estatus_df) * 40),
                        title="Distribución de Clientes por Estatus"
                    )
                    
                    # Agregar etiquetas con valores
                    text = alt.Chart(estatus_df).mark_text(
                        align="left",
                        baseline="middle",
                        dx=5,
                        fontSize=11,
                        fontWeight="bold"
                    ).encode(
                        x=alt.X("cantidad:Q"),
                        y=alt.Y("estatus:N", sort="-x"),
                        text=alt.Text("cantidad:Q", format="d"),
                        color=alt.value("black")
                    )
                    
                    st.altair_chart(chart + text, use_container_width=True)
        
        # 📅 TAB 2: POR FECHA
        with dash_tab2:
            st.subheader("Análisis Temporal")
            
            # Selector de período predefinido
            periodo_options = {
                "Último mes": 30,
                "Últimos 3 meses": 90,
                "Últimos 6 meses": 180,
                "Último año": 365,
                "Todo el histórico": None,
                "Personalizado": "custom"
            }
            
            col_periodo, col_custom = st.columns([1, 2])
            
            with col_periodo:
                periodo_seleccionado = st.selectbox(
                    "Seleccionar período:",
                    list(periodo_options.keys()),
                    index=0
                )
            
            # Determinar fechas según selección
            date_col = None
            for c in ("fecha_ingreso", "fecha_dispersion"):
                if c in df_cli.columns:
                    date_col = c
                    break
            
            if date_col is None:
                st.warning("No se encontraron columnas de fecha válidas.")
            else:
                df_temp = df_cli.copy()
                # Aplicar el manejo flexible de fechas
                try:
                    df_temp[date_col] = parse_dates_flexible(df_temp[date_col])
                except Exception:
                    df_temp[date_col] = pd.to_datetime(df_temp[date_col], errors="coerce")
                df_temp = df_temp.dropna(subset=[date_col])
                
                if df_temp.empty:
                    st.warning("No hay datos con fechas válidas.")
                else:
                    today = pd.Timestamp.now()
                    
                    if periodo_seleccionado == "Personalizado":
                        with col_custom:
                            fecha_min = df_temp[date_col].min().date()
                            fecha_max = df_temp[date_col].max().date()
                            
                            rango_fechas = st.date_input(
                                "Seleccionar rango:",
                                value=(fecha_min, fecha_max),
                                min_value=fecha_min,
                                max_value=fecha_max
                            )
                            
                            if isinstance(rango_fechas, tuple) and len(rango_fechas) == 2:
                                fecha_inicio, fecha_fin = rango_fechas
                            else:
                                fecha_inicio = fecha_fin = rango_fechas
                    else:
                        dias = periodo_options[periodo_seleccionado]
                        if dias is None:  # Todo el histórico
                            fecha_inicio = df_temp[date_col].min().date()
                            fecha_fin = df_temp[date_col].max().date()
                        else:
                            fecha_fin = today.date()
                            fecha_inicio = (today - pd.Timedelta(days=dias)).date()
                    
                    # Filtrar datos por período
                    mask = (df_temp[date_col].dt.date >= fecha_inicio) & (df_temp[date_col].dt.date <= fecha_fin)
                    df_periodo = df_temp[mask].copy()
                    
                    if df_periodo.empty:
                        st.info("No hay datos en el período seleccionado.")
                    else:
                        # KPIs del período
                        st.markdown("**KPIs del período seleccionado:**")
                        
                        total_periodo = len(df_periodo)
                        dispersados_periodo = len(df_periodo[df_periodo["estatus"] == "DISPERSADO"])
                        tasa_conversion = (dispersados_periodo / total_periodo * 100) if total_periodo > 0 else 0
                        
                        kpi_col1, kpi_col2, kpi_col3 = st.columns(3)
                        
                        with kpi_col1:
                            st.metric("Total en período", total_periodo)
                        with kpi_col2:
                            st.metric("Dispersados", dispersados_periodo)
                        with kpi_col3:
                            st.metric("Tasa de conversión", f"{tasa_conversion:.1f}%")
                        
                        st.markdown("---")
                        
                        # Gráfico de evolución temporal (por mes)
                        st.markdown("**Evolución mensual:**")
                        
                        df_mensual = df_periodo.copy()
                        df_mensual["año_mes"] = df_mensual[date_col].dt.to_period("M")
                        evolucion = df_mensual.groupby("año_mes").size().reset_index(name="cantidad")
                        evolucion["fecha"] = evolucion["año_mes"].astype(str)
                        
                        if len(evolucion) > 0:
                            line_chart = alt.Chart(evolucion).mark_line(
                                point=True,
                                strokeWidth=3
                            ).encode(
                                x=alt.X("fecha:O", axis=alt.Axis(title="Mes", labelAngle=-45)),
                                y=alt.Y("cantidad:Q", axis=alt.Axis(title="Número de Clientes")),
                                tooltip=["fecha:O", "cantidad:Q"]
                            ).properties(
                                height=300,
                                title="Evolución de Clientes por Mes"
                            )
                            
                            st.altair_chart(line_chart, use_container_width=True)
                        
                        # Gráfico de dona: Distribución por estatus en el período
                        st.markdown("**Distribución por estatus en el período:**")
                        
                        estatus_periodo = df_periodo["estatus"].fillna("(Sin estatus)").value_counts()
                        estatus_periodo_df = estatus_periodo.reset_index()
                        estatus_periodo_df.columns = ["estatus", "cantidad"]
                        # Agregar porcentaje como columna calculada
                        estatus_periodo_df["porcentaje"] = (estatus_periodo_df["cantidad"] / total_periodo * 100).round(1)
                        
                        dona_chart = alt.Chart(estatus_periodo_df).mark_arc(
                            innerRadius=50,
                            outerRadius=120
                        ).encode(
                            theta=alt.Theta("cantidad:Q"),
                            color=alt.Color("estatus:N", scale=alt.Scale(scheme="category20")),
                            tooltip=[
                                alt.Tooltip("estatus:N", title="Estatus"),
                                alt.Tooltip("cantidad:Q", title="Cantidad"),
                                alt.Tooltip("porcentaje:Q", title="Porcentaje (%)", format=".1f")
                            ]
                        ).properties(
                            height=300,
                            title="Distribución por Estatus"
                        )
                        
                        st.altair_chart(dona_chart, use_container_width=True)
        
        # 🏢 TAB 3: POR SUCURSAL/ASESOR
        with dash_tab3:
            st.subheader("Análisis por Organización")
            
            # Sub-tabs: Sucursales, Asesores, Fuentes
            sub_tab1, sub_tab2, sub_tab3 = st.tabs(["🏢 Sucursales", "👥 Asesores", "📢 Fuentes"])
            
            # SUB-TAB: SUCURSALES
            with sub_tab1:
                sucursal_counts = df_cli["sucursal"].fillna("(Sin sucursal)").value_counts()
                sucursal_counts = sucursal_counts[sucursal_counts > 0]
                
                if sucursal_counts.empty:
                    st.info("No hay datos de sucursales para mostrar.")
                else:
                    sucursal_df = sucursal_counts.reset_index()
                    sucursal_df.columns = ["sucursal", "cantidad"]
                    sucursal_df["porcentaje"] = (sucursal_df["cantidad"] / total_clientes * 100).round(1)
                    
                    col1, col2 = st.columns([1, 2])
                    
                    with col1:
                        st.markdown("**Distribución por sucursal:**")
                        st.dataframe(
                            sucursal_df.assign(
                                **{"Porcentaje": sucursal_df["porcentaje"].astype(str) + "%"}
                            )[["sucursal", "cantidad", "Porcentaje"]],
                            use_container_width=True,
                            hide_index=True
                        )
                    
                    with col2:
                        chart_suc = alt.Chart(sucursal_df).mark_arc(
                            innerRadius=30,
                            outerRadius=100
                        ).encode(
                            theta=alt.Theta("cantidad:Q"),
                            color=alt.Color("sucursal:N", scale=alt.Scale(scheme="set3")),
                            tooltip=[
                                alt.Tooltip("sucursal:N", title="Sucursal"),
                                alt.Tooltip("cantidad:Q", title="Cantidad"),
                                alt.Tooltip("porcentaje:Q", title="Porcentaje (%)", format=".1f")
                            ]
                        ).properties(
                            height=300,
                            title="Distribución por Sucursal"
                        )
                        
                        st.altair_chart(chart_suc, use_container_width=True)
            
            # SUB-TAB: ASESORES
            with sub_tab2:
                # Para la sección de asesores del dashboard, usar datos filtrados EXCEPTO por asesor
                # para mostrar todos los asesores (misma lógica que en tab_asesores)
                try:
                    # Aplicar filtros SIN incluir el filtro de asesor
                    sucursal_for_filter = df_cli["sucursal"].fillna("").replace({"": "(Sin sucursal)"})
                    fuente_for_filter = df_cli["fuente"].fillna("").replace({"": "(Sin fuente)"})
                    
                    # Aplicar solo filtros de sucursal, estatus y fuente (NO asesor)
                    if not f_suc or len(f_suc) == 0 or set(f_suc) == set(SUC_ALL):
                        suc_mask_dash = pd.Series(True, index=df_cli.index)
                    else:
                        suc_mask_dash = sucursal_for_filter.isin(f_suc)

                    if not f_est or len(f_est) == 0 or set(f_est) == set(EST_ALL):
                        est_mask_dash = pd.Series(True, index=df_cli.index)
                    else:
                        est_mask_dash = df_cli["estatus"].isin(f_est)

                    if not f_fuente or len(f_fuente) == 0 or set(f_fuente) == set(FUENTE_ALL):
                        fuente_mask_dash = pd.Series(True, index=df_cli.index)
                    else:
                        fuente_mask_dash = fuente_for_filter.isin(f_fuente)
                    
                    # Datos filtrados para asesores (sin filtro de asesor)
                    df_asesor_dash = df_cli[suc_mask_dash & est_mask_dash & fuente_mask_dash]
                except Exception:
                    # Fallback: usar todos los datos
                    df_asesor_dash = df_cli
                
                asesor_counts = df_asesor_dash["asesor"].fillna("(Sin asesor)").value_counts()
                asesor_counts = asesor_counts[asesor_counts > 0]
                
                if asesor_counts.empty:
                    st.info("No hay datos de asesores para mostrar.")
                else:
                    # Mostrar TODOS los asesores (sin límite de Top 10)
                    asesor_final = asesor_counts
                    
                    asesor_df = asesor_final.reset_index()
                    asesor_df.columns = ["asesor", "cantidad"]
                    # Usar el total de clientes filtrados para calcular porcentajes
                    total_filtrados = len(df_asesor_dash)
                    asesor_df["porcentaje"] = (asesor_df["cantidad"] / total_filtrados * 100).round(1) if total_filtrados > 0 else 0
                    
                    col1, col2 = st.columns([1, 2])
                    
                    with col1:
                        st.markdown("**Top asesores:**")
                        st.dataframe(
                            asesor_df.assign(
                                **{"Porcentaje": asesor_df["porcentaje"].astype(str) + "%"}
                            )[["asesor", "cantidad", "Porcentaje"]],
                            use_container_width=True,
                            hide_index=True
                        )
                    
                    with col2:
                        # Gráfico horizontal para asesores (mejor legibilidad)
                        chart_ases = alt.Chart(asesor_df).mark_bar(
                            cornerRadiusTopRight=4,
                            cornerRadiusBottomRight=4
                        ).encode(
                            x=alt.X("cantidad:Q", axis=alt.Axis(title="Cantidad")),
                            y=alt.Y("asesor:N", sort="-x", axis=alt.Axis(title="Asesor")),
                            color=alt.Color("asesor:N", scale=alt.Scale(scheme="category20"), legend=None),
                            tooltip=[
                                alt.Tooltip("asesor:N", title="Asesor"),
                                alt.Tooltip("cantidad:Q", title="Cantidad"),
                                alt.Tooltip("porcentaje:Q", title="Porcentaje (%)", format=".1f")
                            ]
                        ).properties(
                            height=max(300, len(asesor_df) * 30),
                            title="Distribución por Asesor"
                        )
                        
                        st.altair_chart(chart_ases, use_container_width=True)
            
            # SUB-TAB: FUENTES
            with sub_tab3:
                fuente_counts = df_cli["fuente"].fillna("(Sin fuente)").value_counts()
                fuente_counts = fuente_counts[fuente_counts > 0]
                
                if fuente_counts.empty:
                    st.info("No hay datos de fuentes para mostrar.")
                else:
                    fuente_df = fuente_counts.reset_index()
                    fuente_df.columns = ["fuente", "cantidad"]
                    fuente_df["porcentaje"] = (fuente_df["cantidad"] / total_clientes * 100).round(1)
                    
                    col1, col2 = st.columns([1, 2])
                    
                    with col1:
                        st.markdown("**Distribución por fuente:**")
                        st.dataframe(
                            fuente_df.assign(
                                **{"Porcentaje": fuente_df["porcentaje"].astype(str) + "%"}
                            )[["fuente", "cantidad", "Porcentaje"]],
                            use_container_width=True,
                            hide_index=True
                        )
                    
                    with col2:
                        # Gráfico de dona para fuentes
                        chart_fuente = alt.Chart(fuente_df).mark_arc(
                            innerRadius=40,
                            outerRadius=110
                        ).encode(
                            theta=alt.Theta("cantidad:Q"),
                            color=alt.Color("fuente:N", scale=alt.Scale(scheme="pastel1")),
                            tooltip=[
                                alt.Tooltip("fuente:N", title="Fuente"),
                                alt.Tooltip("cantidad:Q", title="Cantidad"),
                                alt.Tooltip("porcentaje:Q", title="Porcentaje (%)", format=".1f")
                            ]
                        ).properties(
                            height=300,
                            title="Distribución por Fuente de Captación"
                        )
                        
                        st.altair_chart(chart_fuente, use_container_width=True)

        # 💰 ANÁLISIS FINANCIERO AVANZADO - CARTERA KAPITALIZA
        st.markdown("---")
        st.subheader(" Análisis Financiero — Cartera Kapitaliza")
        
        # Preparar datos con limpieza
        df_temp = df_cli.copy()
        
        def limpiar_monto_simple(monto_str):
            if pd.isna(monto_str) or str(monto_str).strip() == "":
                return 0.0
            try:
                import re
                clean = re.sub(r'[,$\s]', '', str(monto_str))
                return float(clean)
            except:
                return 0.0
        
        # Usar monto_final para dispersados, monto_propuesta para el resto
        df_temp['monto_analisis'] = df_temp.apply(
            lambda row: limpiar_monto_simple(row['monto_final']) if row['estatus'] == 'DISPERSADO' 
            else limpiar_monto_simple(row['monto_propuesta']), axis=1
        )
        
        # Filtrar solo clientes con monto > 0
        df_analisis = df_temp[df_temp['monto_analisis'] > 0].copy()
        
        if df_analisis.empty:
            st.info("No hay datos con montos registrados para análisis.")
        else:
            # === MODELO FINANCIERO REFINADO ===
            # Tres componentes: Probabilidad de Conversión, Factor de Retorno, Riesgo
            prob_conversion = {
                "DISPERSADO": 1.00,
                "APROB. CON PROPUESTA": 0.75,
                "PROPUESTA": 0.75,  # Alias para propuesta
                "PEND. ACEPT. CLIENTE": 0.65,
                "PENDIENTE CLIENTE": 0.65,  # Alias
                "PEND. DOC. PARA EVALUACION": 0.45,
                "PENDIENTE DOC": 0.45,  # Alias
                "EN ONBOARDING": 0.55,
                "RECH. CLIENTE CANCELA": 0.10,
                "RECH. SOBREENDEUDAMIENTO": 0.05,
            }
            
            factor_retorno = {
                "DISPERSADO": 1.00,
                "APROB. CON PROPUESTA": 0.85,
                "PROPUESTA": 0.85,
                "PEND. ACEPT. CLIENTE": 0.80,
                "PENDIENTE CLIENTE": 0.80,
                "PEND. DOC. PARA EVALUACION": 0.70,
                "PENDIENTE DOC": 0.70,
                "EN ONBOARDING": 0.75,
                "RECH. CLIENTE CANCELA": 0.00,
                "RECH. SOBREENDEUDAMIENTO": 0.00,
            }
            
            riesgo_pct = {
                "DISPERSADO": 5,
                "APROB. CON PROPUESTA": 20,
                "PROPUESTA": 20,
                "PEND. ACEPT. CLIENTE": 30,
                "PENDIENTE CLIENTE": 30,
                "PEND. DOC. PARA EVALUACION": 45,
                "PENDIENTE DOC": 45,
                "EN ONBOARDING": 40,
                "RECH. CLIENTE CANCELA": 90,
                "RECH. SOBREENDEUDAMIENTO": 95,
            }
            
            # Para otros estatus de rechazo, aplicar valores altos
            for estatus in df_analisis['estatus'].unique():
                if estatus and (estatus.startswith("RECH") or estatus.startswith("REC")):
                    if estatus not in prob_conversion:
                        prob_conversion[estatus] = 0.05
                        factor_retorno[estatus] = 0.00
                        riesgo_pct[estatus] = 95
            
            # Aplicar modelo refinado
            df_analisis["Probabilidad de Conversión"] = df_analisis["estatus"].map(prob_conversion).fillna(0.5)
            df_analisis["Factor Retorno"] = df_analisis["estatus"].map(factor_retorno).fillna(0.5)
            df_analisis["Riesgo (%)"] = df_analisis["estatus"].map(riesgo_pct).fillna(50)
            
            # Cálculos financieros
            df_analisis["Monto Esperado"] = df_analisis["monto_analisis"] * df_analisis["Probabilidad de Conversión"]
            df_analisis["Retorno Esperado"] = df_analisis["monto_analisis"] * df_analisis["Factor Retorno"]
            df_analisis["Riesgo Estimado (%)"] = df_analisis["Riesgo (%)"]  # Para compatibilidad
            
            # === MÉTRICAS GLOBALES REFINADAS ===
            total_cartera = df_analisis["monto_analisis"].sum()
            total_monto_esperado = df_analisis["Monto Esperado"].sum()
            total_retorno = df_analisis["Retorno Esperado"].sum()
            prom_riesgo = df_analisis["Riesgo (%)"].mean()
            prom_conversion = df_analisis["Probabilidad de Conversión"].mean() * 100
            
            # === DIAGNÓSTICO EJECUTIVO AUTOMATIZADO ===
            st.markdown("##### 🧠 Diagnóstico Financiero")
            
            # Resumen ejecutivo
            st.markdown(f"""
            **📊 Resumen Ejecutivo:**
            - Cartera total: **{formatear_monto(total_cartera)}**
            - Conversión esperada: **{formatear_monto(total_monto_esperado)}** ({(total_monto_esperado/total_cartera*100):.1f}% de la cartera)
            - Retorno esperado: **{formatear_monto(total_retorno)}**
            - Riesgo promedio: **{prom_riesgo:.1f}%**
            - Conversión media: **{prom_conversion:.1f}%**
            """)
            
            st.markdown("---")
            
            # Análisis automatizado tipo analista
            if prom_riesgo < 30 and prom_conversion > 70:
                st.success("✅ **Cartera saludable.** Alta probabilidad de conversión y bajo nivel de riesgo. Se recomienda mantener políticas actuales de aprobación.")
            elif prom_riesgo < 60 and prom_conversion > 40:
                st.warning("⚠️ **Cartera moderada.** Algunos clientes requieren seguimiento. Enfocar esfuerzos en etapas de aceptación y documentación.")
            else:
                st.error("🚨 **Riesgo alto.** La mayoría de las solicitudes tienen baja conversión o alta cancelación. Revisar criterios de originación y documentación.")
            
            # Análisis específico de oportunidades (incluir todos los pendientes)
            # Definir estatus que se consideran "pendientes" o en proceso
            estatus_pendientes = df_analisis[~df_analisis['estatus'].str.contains('DISPERSADO|RECHAZADO', na=False, case=False)]
            
            if not estatus_pendientes.empty:
                total_pendientes = len(estatus_pendientes)
                monto_pendientes = estatus_pendientes['monto_analisis'].sum()
                pct_pendientes = (total_pendientes / len(df_analisis) * 100)
                
                # Obtener breakdown por estatus pendiente
                breakdown_pendientes = estatus_pendientes.groupby('estatus').agg({
                    'nombre': 'count',
                    'monto_analisis': 'sum'
                }).rename(columns={'nombre': 'cantidad'})
                
                # Crear texto detallado con los estatus pendientes
                detalle_estatus = []
                for estatus, row in breakdown_pendientes.iterrows():
                    detalle_estatus.append(f"{row['cantidad']} {estatus}")
                
                detalle_texto = ", ".join(detalle_estatus[:3])  # Mostrar los primeros 3
                if len(detalle_estatus) > 3:
                    detalle_texto += f" y {len(detalle_estatus) - 3} más"
                
                st.info(f"💡 **Oportunidad:** {total_pendientes} créditos pendientes ({pct_pendientes:.0f}% del total) representan {formatear_monto(monto_pendientes)} en oportunidades. Incluye: {detalle_texto}. Priorizar seguimiento puede mejorar conversión.")
            
            # === RESUMEN POR ESTATUS ===
            col_tabla, col_grafico = st.columns([1, 1])
            
            with col_tabla:
                st.markdown("##### 🔍 Distribución por Estatus")
                resumen = df_analisis.groupby("estatus").agg({
                    "monto_analisis": "sum",
                    "Retorno Esperado": "sum",
                    "Riesgo Estimado (%)": "mean",
                    "nombre": "count"  # Contar clientes
                }).rename(columns={"nombre": "Clientes"}).reset_index()
                
                # Formatear para mostrar
                resumen_display = resumen.copy()
                resumen_display["Monto Total"] = resumen_display["monto_analisis"].apply(lambda x: f"${x:,.0f}")
                resumen_display["Retorno Esperado"] = resumen_display["Retorno Esperado"].apply(lambda x: f"${x:,.0f}")
                resumen_display["Riesgo (%)"] = resumen_display["Riesgo Estimado (%)"].apply(lambda x: f"{x:.1f}%")
                
                st.dataframe(
                    resumen_display[["estatus", "Clientes", "Monto Total", "Retorno Esperado", "Riesgo (%)"]],
                    use_container_width=True,
                    hide_index=True
                )
            
            with col_grafico:
                pass  # Columna vacía
            
            st.markdown("---")
            st.caption("© CRM Kapitaliza ")

# ===== Clientes (alta + edición) =====
with tab_cli:
    # Cargar datos frescos para la pestaña de clientes
    df_cli = cargar_y_corregir_clientes()
    
    st.subheader("➕ Agregar cliente")
    with st.expander("Formulario de alta", expanded=False):  # UI más limpia

        # --- NEW: eliminar el selectbox "Asesor" (se pide quitar el "botoncito").
        # Mantener sólo el checkbox para crear un nuevo asesor y el input para su nombre.
        st.checkbox("Nuevo asesor (marca para escribir nombre y apellido)", key="form_new_asesor_toggle", help="Marca para ingresar manualmente el nombre y apellido del asesor")
        if st.session_state.get("form_new_asesor_toggle", False):
            st.text_input("Nombre y apellido del nuevo asesor", placeholder="Ej. Juan Pérez", key="form_nuevo_asesor")
        # --- END NEW ---

        with st.form("form_alta_cliente", clear_on_submit=True):
            c1, c2, c3 = st.columns(3)
            with c1:
                id_n = st.text_input("ID (opcional)", key="form_id")
                nombre_n = st.text_input("Nombre *")
                sucursal_n = st.selectbox("Sucursal *", SUCURSALES)

                # REPLACED: permitir elegir un asesor existente dentro del form,
                # o usar el "Nuevo asesor" si el checkbox (fuera del form) está marcado.
                raw_ases = [a for a in df_cli["asesor"].fillna("").unique() if str(a).strip()]
                asesores_exist = sorted(list(dict.fromkeys([_norm_sin_asesor_label(a) for a in raw_ases])))
                # Construir opciones sin duplicados (asegurando la etiqueta estándar "(Sin asesor)")
                asesores_choices = list(dict.fromkeys(["(Sin asesor)"] + asesores_exist))
                asesor_select = st.selectbox("Asesor", asesores_choices, key="form_ases_select")

                if st.session_state.get("form_new_asesor_toggle", False):
                    # si el usuario marcó "Nuevo asesor" usamos el texto ingresado (tiene prioridad)
                    asesor_n = st.session_state.get("form_nuevo_asesor", "").strip()
                else:
                    # usar la selección del selectbox (o '' si eligió "(Sin asesor)")
                    asesor_n = "" if asesor_select == "(Sin asesor)" else asesor_select

                analista_n = st.text_input("Analista")
            with c2:
                fecha_ingreso_n = st.date_input("Fecha ingreso", value=date.today())
                fecha_dispersion_n = st.date_input("Fecha dispersión", value=date.today())
                estatus_n = st.selectbox("Estatus", ESTATUS_OPCIONES, index=0)
                segundo_estatus_n = st.selectbox("Segundo estatus", SEGUNDO_ESTATUS_OPCIONES, index=0)
            with c3:
                monto_prop_n = st.text_input("Monto propuesta", value="")
                monto_final_n = st.text_input("Monto final", value="")
                score_n = st.text_input("Score", value="")
                telefono_n = st.text_input("Teléfono")
                correo_n = st.text_input("Correo")
                fuente_n = st.text_input("Fuente", value="")
            obs_n = st.text_area("Observaciones")

            st.markdown("**Documentos:**")
            
            # Opción para elegir dónde guardar
            col_storage1, col_storage2 = st.columns(2)
            with col_storage1:
                usar_google_drive = st.checkbox(
                    "📤 Guardar en Google Drive", 
                    value=st.session_state.get('drive_creds') is not None,
                    disabled=st.session_state.get('drive_creds') is None,
                    help="Guarda documentos en Google Drive (requiere conexión)"
                )
            with col_storage2:
                if st.session_state.get('drive_creds') is None:
                    st.info("💡 Conecta Google Drive en el sidebar para activar esta opción")
            
            up_estado = st.file_uploader("Estado de cuenta", type=DOC_CATEGORIAS["estado_cuenta"], accept_multiple_files=True, key="doc_estado")
            up_buro   = st.file_uploader("Buró de crédito", type=DOC_CATEGORIAS["buro_credito"], accept_multiple_files=True, key="doc_buro")
            up_solic  = st.file_uploader("Solicitud", type=DOC_CATEGORIAS["solicitud"], accept_multiple_files=True, key="doc_solic")
            up_otros = st.file_uploader("Otros", type=None, accept_multiple_files=True, key="doc_otros")
            st.markdown("Dar doble click para confirmar.")
            if st.form_submit_button("Guardar cliente"):
                
                if not nombre_n.strip():
                    st.warning("El nombre es obligatorio.")
                else:
                    # validar nuevo asesor si corresponde
                    if st.session_state.get("form_new_asesor_toggle", False) and not st.session_state.get("form_nuevo_asesor", "").strip():
                        st.warning("Cuando seleccionas 'Nuevo asesor' debes ingresar el nombre y apellido del asesor.")
                    else:
                        # usar ID proporcionado si existe y es único; si no, generar uno nuevo
                        provided = (id_n or "").strip()
                        if provided:
                            # sanitizar y validar unicidad
                            cid_candidate = safe_name(provided)
                            if cid_candidate in df_cli["id"].astype(str).tolist():
                                st.warning(f"El ID '{cid_candidate}' ya existe. Elige otro o deja vacío para generar uno.")
                                # no continuar con la creación
                                st.stop()
                            cid = cid_candidate
                        else:
                            cid = nuevo_id_cliente(df_cli)
                        # usar asesor_n calculado arriba (puede ser '')
                        asesor_final = find_matching_asesor(asesor_n.strip(), df_cli)
                        nuevo = {
                            "id": cid,
                            "nombre": nombre_n.strip(),
                            "sucursal": sucursal_n,
                            "asesor": asesor_final,
                            "fecha_ingreso": str(fecha_ingreso_n),
                            "fecha_dispersion": str(fecha_dispersion_n),
                            "estatus": estatus_n,
                            "monto_propuesta": str(monto_prop_n).strip(),
                            "monto_final": str(monto_final_n).strip(),
                            "segundo_estatus": segundo_estatus_n,
                            "observaciones": obs_n.strip(),
                            "score": str(score_n).strip(),
                            "telefono": telefono_n.strip(),
                            "correo": correo_n.strip(),
                            "analista": analista_n.strip(),
                            "fuente": fuente_n.strip(),
                        }
                        base = pd.concat([df_cli, pd.DataFrame([nuevo])], ignore_index=True)
                        guardar_clientes(base)
                        # registrar creación en historial
                        actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                        append_historial(cid, nuevo.get("nombre", ""), "", nuevo.get("estatus", ""), "", nuevo.get("segundo_estatus", ""), f"Creado por {actor}", action="CLIENTE AGREGADO", actor=actor)

                        # Guardar documentos (auto refresh al terminar) — acumular y registrar 1 sola entrada en historial
                        subidos_lote = []
                        if up_estado:   subidos_lote += subir_docs(cid, up_estado,   prefijo="estado_", usar_drive=usar_google_drive)
                        if up_buro:     subidos_lote += subir_docs(cid, up_buro,     prefijo="buro_", usar_drive=usar_google_drive)
                        if up_solic:    subidos_lote += subir_docs(cid, up_solic,    prefijo="solic_", usar_drive=usar_google_drive)
                        #if up_contrato: subidos_lote += subir_docs(cid, up_contrato, prefijo="contrato_", usar_drive=usar_google_drive)
                        if up_otros:    subidos_lote += subir_docs(cid, up_otros,    prefijo="otros_", usar_drive=usar_google_drive)

                        if subidos_lote:
                            actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                            append_historial(
                                cid, nuevo.get("nombre",""),
                                nuevo.get("estatus",""), nuevo.get("estatus",""),
                                nuevo.get("segundo_estatus",""), nuevo.get("segundo_estatus",""),
                                f"Subidos: {', '.join(subidos_lote)}",
                                action="DOCUMENTOS", actor=actor
                            )

                        st.success(f"Cliente {cid} creado ✅")
                        do_rerun()  # NEW: refresca todo

    st.subheader("📋 Lista de clientes")
    
    # Usar los datos ya filtrados del sidebar (df_ver)
    # que incluye todos los filtros aplicados correctamente
    df_clientes_mostrar = df_ver.copy()
    
    if df_clientes_mostrar.empty:
        st.info("No hay clientes con los filtros seleccionados.")
    else:
        colcfg = {
            "id": st.column_config.TextColumn("ID", disabled=True),
            "nombre": st.column_config.TextColumn("Nombre"),
            "sucursal": st.column_config.SelectboxColumn("Sucursal", options=[""]+SUCURSALES, required=False),
            "asesor": st.column_config.TextColumn("Asesor"),
            "fecha_ingreso": st.column_config.TextColumn("Fecha ingreso (YYYY-MM-DD)"),
            "fecha_dispersion": st.column_config.TextColumn("Fecha dispersión (YYYY-MM-DD)"),
            "estatus": st.column_config.SelectboxColumn("Estatus", options=ESTATUS_OPCIONES, required=True),
            "monto_propuesta": st.column_config.TextColumn("Monto propuesta"),
            "monto_final": st.column_config.TextColumn("Monto final"),
            "segundo_estatus": st.column_config.SelectboxColumn("Segundo estatus", options=SEGUNDO_ESTATUS_OPCIONES),
            "observaciones": st.column_config.TextColumn("Observaciones"),
            "score": st.column_config.TextColumn("Score"),
            "telefono": st.column_config.TextColumn("Teléfono"),
            "correo": st.column_config.TextColumn("Correo"),
            "analista": st.column_config.TextColumn("Analista"),
            "fuente": st.column_config.TextColumn("Fuente"),
        }

        df_clientes_mostrar["sucursal"] = df_clientes_mostrar["sucursal"].where(df_clientes_mostrar["sucursal"].isin(SUCURSALES), "")
        # antes de mostrar el editor, ordenar df_clientes_mostrar por fechas asc
        df_clientes_mostrar = sort_df_by_dates(df_clientes_mostrar)  # apply ordering
        # FIX: data_editor no acepta ColumnDataKind.DATETIME si la columna está configurada como TextColumn.
        # Convertir las columnas de fecha a strings 'YYYY-MM-DD' para mantener compatibilidad con column_config.
        for _dcol in ("fecha_ingreso", "fecha_dispersion"):
            if _dcol in df_clientes_mostrar.columns:
                try:
                    # Aplicar el manejo flexible de fechas
                    df_temp_col = parse_dates_flexible(df_clientes_mostrar[_dcol])
                    df_clientes_mostrar[_dcol] = df_temp_col.dt.date.astype(str).replace("NaT", "")
                except Exception:
                    df_clientes_mostrar[_dcol] = df_clientes_mostrar[_dcol].astype(str).fillna("")
        ed = st.data_editor(
            df_clientes_mostrar,
            use_container_width=True,
            hide_index=True,
            column_config=colcfg,
            key="editor_clientes"
        )

        st.markdown("### Cambio de estatus")
        # asegurar ids_quick esté siempre definido para evitar NameError si df_clientes_mostrar no tiene 'id' o no es DataFrame
        try:
            ids_quick = (df_clientes_mostrar["id"].tolist() if (isinstance(df_clientes_mostrar, pd.DataFrame) and "id" in df_clientes_mostrar.columns) else [])
            ids_quick = [x for x in ids_quick if str(x).strip()]
            try:
                ids_quick = sorted(ids_quick, key=lambda s: (len(str(s)), str(s)))
            except Exception:
                pass
        except Exception:
            ids_quick = []
        if ids_quick:
            col_q1, col_q2, col_q3, col_q4 = st.columns([2,2,2,3])
            with col_q1:
                cid_quick = st.selectbox("Cliente", ids_quick, format_func=lambda x: f"{x} - {get_nombre_by_id(x)}")
                nombre_q = get_nombre_by_id(cid_quick)
                estatus_actual = get_field_by_id(cid_quick, "estatus")
                seg_actual = get_field_by_id(cid_quick, "segundo_estatus")
            with col_q2:
                nuevo_estatus = st.selectbox("Nuevo estatus", ESTATUS_OPCIONES, index=ESTATUS_OPCIONES.index(estatus_actual) if estatus_actual in ESTATUS_OPCIONES else 0)
            with col_q3:
                nuevo_seg = st.selectbox("Segundo estatus", SEGUNDO_ESTATUS_OPCIONES, index=SEGUNDO_ESTATUS_OPCIONES.index(seg_actual) if seg_actual in SEGUNDO_ESTATUS_OPCIONES else 0)
            with col_q4:
                obs_q = st.text_input("Observaciones (opcional)")
                if st.button("Actualizar estatus"):
                    base = df_cli.set_index("id")
                    base.at[cid_quick, "estatus"] = nuevo_estatus
                    base.at[cid_quick, "segundo_estatus"] = nuevo_seg
                    df_cli = base.reset_index()
                    guardar_clientes(df_cli)
                    # registrar en historial quién hizo el cambio (modificar)
                    actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                    append_historial(cid_quick, nombre_q, estatus_actual, nuevo_estatus, seg_actual, nuevo_seg, obs_q, action="ESTATUS MODIFICADO", actor=actor)
                    st.success(f"Estatus actualizado para {cid_quick} ✅")
                    do_rerun()

        col_save, col_del = st.columns([1,1])
        with col_save:
            if st.button("💾 Guardar cambios"):
                # conservar copia original para detectar cambios y registrar historial
                original_df = df_cli.copy()
                base = df_cli.set_index("id")
                for _, row in ed.iterrows():
                    cid = row["id"]
                    for k in COLUMNS:
                        if k == "id":
                            continue
                        base.at[cid, k] = str(row.get(k, ""))
                # NORMALIZAR/UNIFICAR asesores en el dataframe antes de guardar
                for idx in base.index:
                    base.at[idx, "asesor"] = find_matching_asesor(base.at[idx, "asesor"], base.reset_index())
                df_cli = base.reset_index()
                # registrar en historial los cambios por fila (si hay diferencias relevantes)
                try:
                    actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                    for idx in df_cli.index:
                        cid = df_cli.at[idx, "id"]
                        old_row = original_df[original_df["id"] == cid]
                        if not old_row.empty:
                            old_row = old_row.iloc[0]
                            # comparar estatus y segundo estatus y también detectar cambios en otras columnas
                            diffs = []
                            for c in COLUMNS:
                                if c == "id":
                                    continue
                                oldv = str(old_row.get(c, ""))
                                newv = str(df_cli.at[idx, c])
                                if oldv != newv:
                                    diffs.append(c)
                            if diffs:
                                est_old = old_row.get("estatus", "")
                                est_new = df_cli.at[idx, "estatus"] if "estatus" in df_cli.columns else ""
                                seg_old = old_row.get("segundo_estatus", "")
                                seg_new = df_cli.at[idx, "segundo_estatus"] if "segundo_estatus" in df_cli.columns else ""
                                obs = "Campos cambiados: " + ",".join(diffs)
                                append_historial(cid, df_cli.at[idx, "nombre"], est_old, est_new, seg_old, seg_new, obs, action="ESTATUS MODIFICADO", actor=actor)
                except Exception:
                    pass

                guardar_clientes(df_cli)
                st.success("Cambios guardados ✅")
                # Forzar reconstrucción de filtros de asesores en el sidebar
                try:
                    for _k in ("f_ases", "f_ases_ms", "f_ases_all"):
                        st.session_state.pop(_k, None)
                    st.session_state["_filters_token"] = st.session_state.get("_filters_token", 0) + 1
                except Exception:
                    pass
                do_rerun()

        with col_del:
            st.caption("Eliminar cliente (Siempre dar doble click para confirmar)")
            
            if can("delete_client"):
                    if ids_quick:
                        # mostrar opciones con 'ID - Nombre' para permitir borrar por nombre visualmente
                        opts = [""] + [f"{cid} - {get_nombre_by_id(cid)}" if get_nombre_by_id(cid) else str(cid) for cid in ids_quick]
                        sel = st.selectbox("Cliente a eliminar (ID - Nombre)", opts)
                        # extraer id del texto seleccionado
                        cid_del = ""
                        if sel:
                            if " - " in sel:
                                cid_del = sel.split(" - ", 1)[0]
                            else:
                                cid_del = sel
                    if cid_del and st.button("🗑️ Eliminar seleccionado"):
                        # registrar antes de eliminar
                        try:
                            nombre_del = get_nombre_by_id(cid_del)
                        except Exception:
                            nombre_del = ""
                        actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                        append_historial(cid_del, nombre_del, "", "", "", "", f"Eliminado por {actor}", action="CLIENTE ELIMINADO", actor=actor)
                        df_cli = eliminar_cliente(cid_del, df_cli, borrar_historial=False)
                        st.success(f"Cliente {cid_del} eliminado ✅")
                        do_rerun()
            else:
                st.info("No tienes permiso para eliminar clientes.")

# ===== Documentos (por cliente) =====
# Safety: garantizar que las pestañas fueron creadas; si no, volver a crearlas para evitar NameError.
if 'tab_docs' not in globals():
    tab_dash, tab_cli, tab_docs, tab_import, tab_hist = st.tabs(
        ["📊 Dashboard", "📋 Clientes", "📎 Documentos", "📥 Importar", "🗂️ Historial"]
    )

with tab_docs:
    st.subheader("📎 Documentos por cliente")
    # Cargar datos frescos para la pestaña de documentos
    df_cli = cargar_y_corregir_clientes()
    if df_cli.empty:
        st.info("No hay clientes aún.")
    else:
        ids = (df_cli["id"].tolist() if (isinstance(df_cli, pd.DataFrame) and "id" in df_cli.columns) else [])
        ids = [x for x in ids if str(x).strip()]
        try:
            ids = sorted(ids, key=lambda s: (len(str(s)), str(s)))
        except Exception:
            pass
        cid_sel = st.selectbox(
            "Selecciona cliente",
            [""] + ids,
            format_func=lambda x: "—" if x == "" else f"{x} - {get_nombre_by_id(x)}",
            key="docs_cid_sel"
        )
        if cid_sel:
            estatus_cliente_sel = get_field_by_id(cid_sel, "estatus")
            st.markdown("#### Subir documentos")
            
            # Opción para elegir dónde guardar
            col_storage_e1, col_storage_e2 = st.columns(2)
            with col_storage_e1:
                usar_google_drive_e = st.checkbox(
                    "📤 Guardar en Google Drive", 
                    value=st.session_state.get('drive_creds') is not None,
                    disabled=st.session_state.get('drive_creds') is None,
                    help="Guarda documentos en Google Drive (requiere conexión)",
                    key=f"drive_option_{cid_sel}"
                )
            with col_storage_e2:
                if st.session_state.get('drive_creds') is None:
                    st.info("💡 Conecta Google Drive en el sidebar para activar esta opción")
            
            # Unificar los uploaders en un solo formulario con un botón "Subir archivos"
            form_key = f"form_subir_docs_{cid_sel}"
            with st.form(form_key, clear_on_submit=True):
                up_estado_e = st.file_uploader("Estado de cuenta", type=DOC_CATEGORIAS["estado_cuenta"], accept_multiple_files=True, key=f"estado_{cid_sel}")
                up_buro_e = st.file_uploader("Buró de crédito", type=DOC_CATEGORIAS["buro_credito"], accept_multiple_files=True, key=f"buro_{cid_sel}")
                up_solic_e = st.file_uploader("Solicitud", type=DOC_CATEGORIAS["solicitud"], accept_multiple_files=True, key=f"solic_{cid_sel}")
                up_otros_e = st.file_uploader("Otros", type=DOC_CATEGORIAS["otros"], accept_multiple_files=True, key=f"otros_{cid_sel}")
                if _is_dispersion(estatus_cliente_sel):
                    up_contrato_e = st.file_uploader("Contrato ", type=DOC_CATEGORIAS["contrato"], accept_multiple_files=True, key=f"contrato_{cid_sel}")
                else:
                    up_contrato_e = None
                submitted = st.form_submit_button("⬆️ Subir archivos")
                if submitted:
                    subidos_lote = []
                    if up_estado_e:   subidos_lote += subir_docs(cid_sel, up_estado_e,   prefijo="estado_", usar_drive=usar_google_drive_e)
                    if up_buro_e:     subidos_lote += subir_docs(cid_sel, up_buro_e,     prefijo="buro_", usar_drive=usar_google_drive_e)
                    if up_solic_e:    subidos_lote += subir_docs(cid_sel, up_solic_e,    prefijo="solic_", usar_drive=usar_google_drive_e)
                    if up_otros_e:    subidos_lote += subir_docs(cid_sel, up_otros_e,    prefijo="otros_", usar_drive=usar_google_drive_e)
                    if up_contrato_e: subidos_lote += subir_docs(cid_sel, up_contrato_e, prefijo="contrato_", usar_drive=usar_google_drive_e)

                    if subidos_lote:
                        # Limpia uploaders
                        for k in (f"estado_{cid_sel}", f"buro_{cid_sel}", f"solic_{cid_sel}", f"otros_{cid_sel}", f"contrato_{cid_sel}"):
                            st.session_state.pop(k, None)

                        # 1 sola línea en historial
                        actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                        try:
                            nombre_x = get_nombre_by_id(cid_sel)
                            est_x    = get_field_by_id(cid_sel, "estatus")
                            seg_x    = get_field_by_id(cid_sel, "segundo_estatus")
                        except Exception:
                            nombre_x = est_x = seg_x = ""

                        append_historial(
                            str(cid_sel), nombre_x,
                            est_x, est_x, seg_x, seg_x,
                            f"Subidos: {', '.join(subidos_lote)}",
                            action="DOCUMENTOS", actor=actor
                        )

                        # refresco inmediato (token)
                        tok_key = f"docs_token_{cid_sel}"
                        st.session_state[tok_key] = st.session_state.get(tok_key, 0) + 1
                        st.success(f"Archivo(s) subido(s): {len(subidos_lote)} ✅")
                    else:
                        st.info("No seleccionaste archivos para subir.")

            st.markdown("—")
            if _is_dispersion(estatus_cliente_sel):
                st.success("Estatus actual: DISPERSADO — sube el Contrato en el formulario de arriba.")
            else:
                st.info("Para subir el Contrato, cambia el estatus del cliente a **en dispersión**.")

            # token para forzar re-render de botones de descarga cuando haya uploads
            tok_key = f"docs_token_{cid_sel}"
            tok = st.session_state.get(tok_key, 0)

            files = listar_docs_cliente(cid_sel)
            if files:
                st.markdown("#### Archivos del cliente")
                # mapping explícito de prefijos usados al guardar
                prefix_map = {
                    "estado_cuenta": "estado_",
                    "buro_credito": "buro_",
                    "solicitud": "solic_",
                    "contrato": "contrato_",
                    "otros": "otros_",
                }
                for cat in DOC_CATEGORIAS.keys():
                    pref = prefix_map.get(cat, cat.split('_')[0] + "_")
                    cat_files = [f for f in files if f.name.startswith(pref)]
                    if cat_files:
                        st.write(f"• {cat.replace('_',' ').title()}:")
                        for f in cat_files:
                                    # Flow: user clicks a request button which prepares the bytes and
                                    # registers the download in historial; then a download_button
                                    # appears where the user can complete the download.
                                    req_key = f"dl_req_{cid_sel}_{tok}_{f.name}"
                                    blob_key = f"dl_blob_{cid_sel}_{f.name}"
                                    btn_label = f"{f.name}"
                                    #if st.button(btn_label, key=req_key):
                                    try:
                                        data_bytes = f.read_bytes()
                                    except Exception:
                                        data_bytes = b""

                                    if st.download_button(
                                        f"⬇️Descargar {f.name}",
                                        data=data_bytes,
                                        file_name=f.name,
                                        key=f"dl_btn_{cid_sel}_{tok}_{f.name}"
                                    ):
                                            # Registrar en historial la descarga
                                        try:
                                            nombre_x = get_nombre_by_id(cid_sel)
                                            est_x    = get_field_by_id(cid_sel, "estatus")
                                            seg_x    = get_field_by_id(cid_sel, "segundo_estatus")
                                        except Exception:
                                            nombre_x = est_x = seg_x = ""
                                        actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                                        append_historial(
                                            str(cid_sel), nombre_x,
                                            est_x, est_x, seg_x, seg_x,
                                            f"Descargado: {f.name}",
                                            action="DESCARGA DOCUMENTO",   # o el que uses en ACTION_LABELS
                                            actor=actor
                                        )


                                    # --- Acciones adicionales: Eliminar / Reemplazar ---
                                    a1, a2 = st.columns([1, 2])
                                    with a1:
                                        del_key = f"del_file_{cid_sel}_{f.name}"
                                        if st.button("Eliminar", key=del_key):
                                            try:
                                                # borrar archivo físico
                                                f.unlink()
                                                # limpiar blobs relacionados
                                                st.session_state.pop(blob_key, None)
                                                # forzar refresh de botones
                                                st.session_state[tok_key] = st.session_state.get(tok_key, 0) + 1
                                                # historial
                                                try:
                                                    actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                                                    append_historial(str(cid_sel), get_nombre_by_id(cid_sel), "", "", "", "", f"Eliminado: {f.name}", action="DOCUMENTOS", actor=actor)
                                                except Exception:
                                                    pass
                                                st.success(f"Archivo eliminado: {f.name}")
                                            except Exception as e:
                                                st.error(f"No se pudo eliminar {f.name}: {e}")

                            
                # Después de listar todas las categorías, ofrecer ZIP del cliente (una sola vez)
                if st.button("📦 Descargar carpeta (ZIP)", key=f"zip_cliente_{cid_sel}"):
                    zip_buffer = io.BytesIO()
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                        for f in files:
                            if f.is_file():
                                zf.write(f, arcname=f.name)
                    zip_buffer.seek(0)
                    # registrar en historial la descarga del ZIP del cliente
                    try:
                        actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                        append_historial(str(cid_sel), get_nombre_by_id(cid_sel), "", "", "", "", f"ZIP cliente preparado ({len(files)} archivos)", action="DESCARGA ZIP CLIENTE", actor=actor)
                    except Exception:
                        pass
                    st.session_state[f"_last_zip_{cid_sel}"] = zip_buffer.getvalue()
                if st.session_state.get(f"_last_zip_{cid_sel}"):
                    st.download_button(
                        "⬇️ Descargar ZIP del cliente",
                        data=st.session_state.get(f"_last_zip_{cid_sel}"),
                        file_name=f"{safe_name(cid_sel)}_{safe_name(get_nombre_by_id(cid_sel))}.zip",
                        mime="application/zip",
                        key=f"dl_zip_cliente_{cid_sel}"
                    )

            if can("delete_client"):
                st.markdown("---")
                st.error("⚠️ Eliminar cliente (borra su carpeta y su historial).")
                st.caption("Siempre dar doble click para confirmar.")
                if st.button(f"🗑️ Eliminar {cid_sel}", key=f"del_{cid_sel}"):
                    # registrar en historial antes de eliminar
                    nombre_del = get_nombre_by_id(cid_sel)
                    actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                    append_historial(cid_sel, nombre_del, "", "", "", "", f"Eliminado por {actor}", action="CLIENTE ELIMINADO", actor=actor)
                    df_cli = eliminar_cliente(cid_sel, df_cli, borrar_historial=False)
                    st.success(f"Cliente {cid_sel} eliminado ✅")
                    do_rerun()
            else:
                st.info("Solo el administrador puede eliminar clientes.")

# ===== Importar (SOLO Excel) =====  # NEW: ZIP eliminado
with tab_import:
    st.subheader("📥 Importar clientes desde Excel (.xlsx)")
    st.caption("Descarga la plantilla, mapea columnas y ejecuta la importación.")

    import_cols_required = [
        "nombre","sucursal","asesor","fecha_ingreso","fecha_dispersion",
        "estatus","monto_propuesta","monto_final","segundo_estatus",
        "observaciones","score","telefono","correo","analista"
    ]
    import_cols_optional = ["id", "fuente"]  # si viene, permite actualizar por ID

    cta1, cta2 = st.columns([1,3])
    with cta1:
        st.caption("Plantilla de importación deshabilitada.")

    def _read_excel_uploaded(file) -> pd.DataFrame:
        try:
            df = pd.read_excel(file, dtype=str).fillna("")
            if sum(str(c).startswith("Unnamed") for c in df.columns) > len(df.columns) * 0.6:
                headers = [str(x).strip() for x in df.iloc[0].tolist()]
                df = df.iloc[1:].copy()
                df.columns = headers
            df.columns = [str(c).strip() for c in df.columns]
            return df
        except Exception as e:
            st.error(f"Error leyendo Excel: {e}")
            return pd.DataFrame()

    up_excel = st.file_uploader("Sube tu Excel (.xlsx)", type=["xlsx"], accept_multiple_files=False, key="up_excel_main")

    if up_excel:
        df_imp_raw = _read_excel_uploaded(up_excel)
        if df_imp_raw.empty:
            st.warning("El Excel está vacío o no se pudo leer.")
        else:
            with st.expander("Vista previa", expanded=True):
                st.dataframe(sort_df_by_dates(df_imp_raw).head(10), use_container_width=True)

            st.markdown("#### Mapeo de columnas")
            df_cols = df_imp_raw.columns.tolist()
            mapping = {}
            map_cols = import_cols_optional + import_cols_required
            M1, M2, M3 = st.columns(3)
            for i, col_needed in enumerate(map_cols):
                col = [M1, M2, M3][i % 3]
                mapping[col_needed] = col.selectbox(
                    f"Excel → {col_needed}",
                    ["(no asignar)"] + df_cols,
                    index=(df_cols.index(col_needed) + 1) if col_needed in df_cols else 0,
                    key=f"map_{col_needed}"
                )

            def _build_norm_df(df_src, mp):
                out = pd.DataFrame()
                for k in map_cols:
                    src = mp.get(k)
                    if src and src != "(no asignar)" and src in df_src.columns:
                        out[k] = df_src[src].astype(str).fillna("")
                    else:
                        out[k] = ""
                return out

            df_norm = _build_norm_df(df_imp_raw, mapping)

            # --- Canonizar valores frente a catálogos existentes para evitar duplicados parecidos ---
            ESTATUS_SYNONYMS = {
                "en revision": "EN REVISIÓN",
                "en revisión": "EN REVISIÓN",
                "revision": "EN REVISIÓN",
                "revisión": "EN REVISIÓN",
            }
            SEGUNDO_ESTATUS_SYNONYMS = {
                # agrega sinónimos si los conoces
            }

            def _canon_est(x: str) -> str:
                try:
                    return canonicalize_from_catalog(x, ESTATUS_OPCIONES, extra_synonyms=ESTATUS_SYNONYMS, min_ratio=0.90)
                except Exception:
                    return x

            def _canon_seg(x: str) -> str:
                try:
                    return canonicalize_from_catalog(x, SEGUNDO_ESTATUS_OPCIONES, extra_synonyms=SEGUNDO_ESTATUS_SYNONYMS, min_ratio=0.90)
                except Exception:
                    return x

            def _canon_suc(x: str) -> str:
                try:
                    return canonicalize_from_catalog(x, SUCURSALES, extra_synonyms=None, min_ratio=0.92)
                except Exception:
                    return x

            for col, fn in [
                ("estatus", _canon_est),
                ("segundo_estatus", _canon_seg),
                ("sucursal", _canon_suc),
            ]:
                if col in df_norm.columns:
                    try:
                        df_norm[col] = df_norm[col].astype(str).map(fn)
                    except Exception:
                        df_norm[col] = df_norm[col]

            # Detectar nuevos valores que no estén en los catálogos actuales
            nuevas_suc = sorted(set(df_norm.loc[df_norm["sucursal"].ne(""), "sucursal"]) - set(SUCURSALES))
            nuevos_est = sorted(set(df_norm.loc[df_norm["estatus"].ne(""), "estatus"]) - set(ESTATUS_OPCIONES))
            nuevos_seg = sorted(set(df_norm.loc[df_norm["segundo_estatus"].ne(""), "segundo_estatus"]) - set(SEGUNDO_ESTATUS_OPCIONES))

            # Agregar automáticamente y persistir
            if nuevas_suc:
                SUCURSALES.extend([s for s in nuevas_suc if s.strip()])
                save_sucursales(SUCURSALES)
                st.info(f"Se agregaron {len(nuevas_suc)} sucursal(es): {', '.join(nuevas_suc)}")

            if nuevos_est:
                ESTATUS_OPCIONES.extend([e for e in nuevos_est if e.strip()])
                save_estatus(ESTATUS_OPCIONES)
                st.info(f"Se agregaron {len(nuevos_est)} estatus: {', '.join(nuevos_est)}")

            if nuevos_seg:
                SEGUNDO_ESTATUS_OPCIONES.extend([e for e in nuevos_seg if e.strip() or e == ""])
                save_segundo_estatus(SEGUNDO_ESTATUS_OPCIONES)
                st.info(f"Se agregaron {len(nuevos_seg)} segundo estatus: {', '.join([x if x else '(vacío)' for x in nuevos_seg])}")

            with st.expander("Previsualización mapeada", expanded=False):
                st.dataframe(sort_df_by_dates(df_norm).head(10), use_container_width=True)

            st.markdown("#### Modo de importación")
            modo = st.radio(
                "¿Cómo quieres importar?",
                ["Agregar (solo nuevos)", "Actualizar por ID (si coincide)", "Upsert por Nombre+Teléfono"],
                horizontal=True,
                key="modo_import"
            )

            # Normalizar fechas a str si vienen tipo fecha
            for fcol in ["fecha_ingreso","fecha_dispersion"]:
                try:
                    df_norm[fcol] = pd.to_datetime(df_norm[fcol], errors="ignore").astype(str).replace("NaT","")
                except Exception:
                    pass

        if st.button("🚀 Importar ahora", type="primary", key="btn_importar_2"):
            base = df_cli.copy()

            def _nuevo_id_local(df):
                base_id = 1000
                try:
                    if not df.empty and "id" in df.columns:
                        nums = []
                        for x in df["id"].astype(str):
                            if str(x).startswith("C"):
                                try:
                                    nums.append(int(str(x).lstrip("C")))
                                except Exception:
                                    continue
                        if nums:
                            base_id = max(nums) + 1
                        else:
                            base_id = base_id + len(df) + 1
                except Exception:
                    base_id = base_id + 1
                return f"C{base_id}"

            actualizados = 0
            agregados = 0

            df_norm_obj = locals().get('df_norm', None)
            if df_norm_obj is not None and (not getattr(df_norm_obj, 'empty', True)):
                for _, r in df_norm_obj.iterrows():
                    r = r.fillna("")
                    rid = str(r.get("id", "")).strip()
                    rnombre = str(r.get("nombre", "")).strip()
                    rtel = str(r.get("telefono", "")).strip()

                    idx = None
                    if modo == "Actualizar por ID (si coincide)" and rid:
                        hit = base.index[base["id"] == rid].tolist()
                        idx = hit[0] if hit else None
                    elif modo == "Upsert por Nombre+Teléfono" and rnombre and rtel:
                        hits = base.index[(base["nombre"] == rnombre) & (base["telefono"] == rtel)].tolist()
                        idx = hits[0] if hits else None

                    # construir registro y mapear asesor contra el base actual
                    registro = {k: str(r.get(k, "")) for k in COLUMNS if k != "id"}
                    registro["asesor"] = find_matching_asesor(registro.get("asesor", ""), base)

                    if idx is not None:
                        for k, v in registro.items():
                            base.at[idx, k] = v
                        actualizados += 1
                        try:
                            actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                            cid_up = base.at[idx, "id"] if "id" in base.columns else (registro.get("id","") or "")
                            append_historial(cid_up, registro.get("nombre",""), "", registro.get("estatus",""), "", registro.get("segundo_estatus",""), f"Importación - actualizado", action="ESTATUS MODIFICADO", actor=actor)
                        except Exception:
                            pass
                    else:
                        if modo == "Agregar (solo nuevos)":
                            if rnombre and rtel and not base[(base["nombre"] == rnombre) & (base["telefono"] == rtel)].empty:
                                continue
                        new_id = rid if rid and (base["id"] != rid).all() else _nuevo_id_local(base)
                        nuevo = {"id": new_id, **registro}
                        base = pd.concat([base, pd.DataFrame([nuevo])], ignore_index=True)
                        agregados += 1
                        try:
                            actor = (current_user() or {}).get("user") or (current_user() or {}).get("email")
                            append_historial(new_id, nuevo.get("nombre",""), "", nuevo.get("estatus",""), "", nuevo.get("segundo_estatus",""), f"Importación - creado", action="CLIENTE AGREGADO", actor=actor)
                        except Exception:
                            pass

            try:
                base = _fix_missing_or_duplicate_ids(base)
            except Exception:
                pass
            guardar_clientes(base)
            st.success(f"Importación completada ✅  |  Agregados: {agregados}  ·  Actualizados: {actualizados}")

            # Limpieza del estado del mapeo para que no “se quede” la UI
            for k in list(st.session_state.keys()):
                if str(k).startswith("map_") or k in ("up_excel_main", "modo_import"):
                    st.session_state.pop(k, None)

            do_rerun()
        # (Se eliminó una copia duplicada del bloque "Historial de movimientos" aquí)
               
# ===== Historial =====
with tab_hist:
    # Mostrar el historial solo a administradores
    if not is_admin():
        st.warning("Solo los administradores pueden ver el historial.")
    else:
        st.subheader("🗂️ Historial de movimientos")
        
        # Lazy loading: solo cargar cuando el usuario lo solicite
        col1, col2 = st.columns([2, 1])
        with col1:
            if st.button("📂 Cargar Historial", key="load_hist", help="Cargar historial desde Google Sheets"):
                st.session_state["hist_loaded"] = True
        with col2:
            if st.session_state.get("hist_loaded", False):
                if st.button("🔄 Refrescar", key="refresh_hist_main", help="Actualizar historial"):
                    st.session_state["force_historial_reload"] = True
                    st.rerun()
        
        # Solo cargar y mostrar historial si está activado
        if st.session_state.get("hist_loaded", False):
            try:
                # Verificar si se solicitó recarga forzada
                force_reload = st.session_state.get("force_historial_reload", False)
                if force_reload:
                    st.session_state["force_historial_reload"] = False
                
                with st.spinner("Cargando historial desde Google Sheets..."):
                    dfh = cargar_historial(force_reload=force_reload)
            except Exception:
                dfh = pd.DataFrame()

            if dfh is None or dfh.empty:
                st.info("No hay registros en el historial.")
            else:
                # asegurarnos de tener columna de timestamp como datetime para ordenar
                try:
                    dfh["_ts_dt"] = pd.to_datetime(dfh["ts"], errors="coerce")
                    dfh = dfh.sort_values("_ts_dt", ascending=False)
                    dfh = dfh.drop(columns=["_ts_dt"])
                except Exception:
                    pass

                # Mostrar filtros simples
                cols_top = st.columns([3,2,2,2])
                with cols_top[0]:
                    qid = st.text_input("Filtrar por ID de cliente (parcial)")
                with cols_top[1]:
                    qactor = st.selectbox("Actor", ["TODOS"] + sorted([str(x) for x in sorted(set(dfh.get("actor",[])))]) , index=0)
                with cols_top[2]:
                    # Etiquetas amigables para las acciones
                    ACTION_LABELS = ["TODOS", "CLIENTE AGREGADO", "DESCARGA ZIP", "DESCARGA ZIP CLIENTE","DESCARGA ZIP ASESOR","DESCARGA DOCUMENTO", "DOCUMENTOS", "CLIENTE ELIMINADO", "ESTATUS MODIFICADO"]
                    qaction = st.selectbox("Acción", ACTION_LABELS, index=0)
                with cols_top[3]:
                    # Botón más pequeño y compacto para refrescar historial
                    st.markdown('<div class="small-refresh-button">', unsafe_allow_html=True)
                    if st.button("🔄 Refrescar", key="refresh_historial", use_container_width=False, help="Actualizar historial desde Google Sheets"):
                        # Forzar recarga del historial usando force_reload
                        st.session_state["force_historial_reload"] = True
                        st.rerun()
                    st.markdown('</div>', unsafe_allow_html=True)

                df_show = dfh.copy()

                # --- Filtro por rango de fechas (columna 'ts') ---
                try:
                    ts_all = pd.to_datetime(dfh.get('ts', pd.Series(dtype='object')), errors='coerce')
                    if not ts_all.dropna().empty:
                        min_ts = ts_all.min()
                        max_ts = ts_all.max()
                        # rango por defecto: últimas 30 días o todo el rango si es menor
                        default_end = max_ts.date()
                        default_start = (max_ts - pd.Timedelta(days=30)).date() if (max_ts - pd.Timedelta(days=30)) > min_ts else min_ts.date()
                        dr = st.date_input("Filtrar historial por fecha (desde → hasta)", value=(default_start, default_end), key="hist_date_range")
                        start_date, end_date = (dr if isinstance(dr, tuple) else (dr, dr))
                        if start_date and end_date:
                            mask_dates = pd.to_datetime(df_show.get('ts', pd.Series(dtype='object')), errors='coerce').dt.date
                            df_show = df_show[mask_dates.between(start_date, end_date)]
                except Exception:
                    pass
                if qid:
                    df_show = df_show[df_show["id"].astype(str).str.contains(qid, case=False, na=False)]
                if qactor and qactor != "TODOS":
                    df_show = df_show[df_show["actor"].astype(str) == qactor]
                if qaction and qaction != "TODOS":
                    qa = qaction
                    if qa == "CLIENTE AGREGADO":
                        df_show = df_show[df_show["action"].astype(str) == "CLIENTE AGREGADO"]
                    elif qa == "CLIENTE ELIMINADO":
                        df_show = df_show[df_show["action"].astype(str) == "CLIENTE ELIMINADO"]
                    elif qa == "ESTATUS MODIFICADO":
                        df_show = df_show[df_show["action"].astype(str) == "ESTATUS MODIFICADO"]
                    elif qa == "DESCARGA ZIP":
                        df_show = df_show[df_show["action"].astype(str) == "DESCARGA ZIP"]
                    elif qa == "DESCARGA ZIP CLIENTE":
                        df_show = df_show[df_show["action"].astype(str) == "DESCARGA ZIP CLIENTE"]
                    elif qa == "DOCUMENTOS":
                        df_show = df_show[df_show["action"].astype(str) == "DOCUMENTOS"]
                    elif qa == "DESCARGA ZIP ASESOR":
                        df_show = df_show[df_show["action"].astype(str) == "DESCARGA ZIP ASESOR"]
                    elif qa == "DESCARGA DOCUMENTO":
                        df_show = df_show[df_show["action"].astype(str) == "DESCARGA DOCUMENTO"]
                    else:
                        # OTROS: no filtrar por action; mostrar todo lo que no cae en las categorías anteriores
                        pass

                st.dataframe(df_show.reset_index(drop=True), use_container_width=True, hide_index=True)

                try:
                    csv_bytes = df_show.to_csv(index=False, encoding="utf-8")
                    st.download_button("⬇️ Descargar historial filtrado (CSV)", data=csv_bytes, file_name="historial_filtrado.csv", mime="text/csv")
                except Exception:
                    pass
                if st.button("🗑️ Borrar historial"):
                    try:
                        # Crear un CSV vacío con las columnas correctas
                        cols = ["id","nombre","estatus_old","estatus_new","segundo_old","segundo_new","observaciones","action","actor","ts"]
                        pd.DataFrame(columns=cols).to_csv(HISTORIAL_CSV, index=False, encoding="utf-8")
                        st.success("Historial eliminado correctamente.")
                        do_rerun()
                    except Exception as e:
                        st.error(f"Error al borrar historial: {e}")
        else:
            st.info("👆 Haz clic en 'Cargar Historial' para ver los registros desde Google Sheets.")